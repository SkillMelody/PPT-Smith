from __future__ import annotations

import json
import importlib.util
import subprocess
import sys
from pathlib import Path
from typing import Any

import pytest
from jsonschema import Draft202012Validator
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
FIXTURES = ROOT / "tests" / "fixtures" / "qa"
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from ppt_qa.package_inspector import inspect_package
from ppt_qa.report import inspection_to_dict
from ppt_qa.slide_inspector import inspect_slides
GENERATOR = FIXTURES / "generate_fixtures.py"
spec = importlib.util.spec_from_file_location("qa_fixture_generator", GENERATOR)
assert spec is not None and spec.loader is not None
generator = importlib.util.module_from_spec(spec)
spec.loader.exec_module(generator)
save_whole_slide_image = generator.save_whole_slide_image


def load_json(path: Path) -> dict[str, Any]:
    with path.open("r", encoding="utf-8") as handle:
        data = json.load(handle)
    assert isinstance(data, dict)
    return data


def inspect_fixture(name: str) -> dict[str, Any]:
    fixture = FIXTURES / name
    package = inspect_package(fixture / "deck.pptx", ppt_ir=load_json(fixture / "ppt-ir.json"))
    if not any(issue.severity == "fatal" for issue in package.issues):
        package = inspect_slides(
            fixture / "deck.pptx",
            package,
            ppt_ir=load_json(fixture / "ppt-ir.json"),
            style=load_json(fixture / "style-contract.json"),
            delivery=load_json(fixture / "delivery-plan.json"),
        )
    return inspection_to_dict(package)


def issue_codes(report: dict[str, Any]) -> set[str]:
    return {issue["issue_code"] for issue in report["issues"]}


def inspect_generated(
    tmp_path: Path,
    build_slide: Any,
    *,
    style: dict[str, Any] | None = None,
    ppt_ir: dict[str, Any] | None = None,
) -> dict[str, Any]:
    deck_path = tmp_path / "generated.pptx"
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    build_slide(slide)
    deck.save(deck_path)
    package = inspect_package(deck_path)
    return inspection_to_dict(
        inspect_slides(
            deck_path,
            package,
            ppt_ir=ppt_ir,
            style=style
            or {
                "colors": {
                    "background": "#FFFFFF",
                    "primary": "#0A2233",
                    "text_primary": "#0A2233",
                    "border": "#E4E0D7",
                },
                "typography": {
                    "minimum_body_size_pt": 10,
                    "minimum_footnote_size_pt": 8,
                    "font_primary": ["Aptos"],
                },
            },
        )
    )


def test_duplicate_title_and_message_do_not_reduce_native_text_ratio(
    tmp_path: Path,
) -> None:
    title = "同一句判断不应在可编辑比例分母里重复计算"

    def build(slide: Any) -> None:
        box = slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(10),
            Inches(0.6),
        )
        run = box.text_frame.paragraphs[0].add_run()
        run.text = title
        run.font.name = "Aptos"
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(10, 34, 51)

    report = inspect_generated(
        tmp_path,
        build,
        ppt_ir={
            "slides": [
                {
                    "id": "S01",
                    "title": title,
                    "message": title,
                    "objects": [],
                }
            ]
        },
    )
    assert report["metrics"]["native_text_ratio"] == 1.0


def test_valid_pptx_package_passes_and_matches_schema() -> None:
    report = inspect_fixture("valid-native-deck")
    assert report["status"] == "passed"
    assert report["issues"] == []
    schema = load_json(ROOT / "schemas" / "pptx-inspection.schema.json")
    Draft202012Validator(schema).validate(report)


def test_missing_media_fails() -> None:
    report = inspect_fixture("missing-media")
    assert "PPTX_MISSING_MEDIA" in issue_codes(report)
    assert report["status"] == "failed"


def test_broken_relationship_fails() -> None:
    report = inspect_fixture("broken-relationship")
    assert "PPTX_MISSING_MEDIA" in issue_codes(report)
    assert report["status"] == "failed"


def test_whole_slide_image_is_detected() -> None:
    report = inspect_fixture("whole-slide-image")
    assert "PPTX_WHOLE_SLIDE_RASTER" in issue_codes(report)


def test_background_image_with_native_overlay_is_allowed(tmp_path: Path) -> None:
    save_whole_slide_image(tmp_path, name="background", delivery_route="background_image", overlay=True)
    package = inspect_package(tmp_path / "deck.pptx", ppt_ir=load_json(tmp_path / "ppt-ir.json"))
    report = inspection_to_dict(
        inspect_slides(
            tmp_path / "deck.pptx",
            package,
            ppt_ir=load_json(tmp_path / "ppt-ir.json"),
            style=load_json(tmp_path / "style-contract.json"),
            delivery=load_json(tmp_path / "delivery-plan.json"),
        )
    )
    assert "PPTX_WHOLE_SLIDE_RASTER" not in issue_codes(report)


def test_font_below_minimum_fails() -> None:
    report = inspect_fixture("font-too-small")
    assert "PPTX_FONT_SIZE_BELOW_MINIMUM" in issue_codes(report)


def test_named_visual_boundary_is_not_an_empty_text_box(tmp_path: Path) -> None:
    def build(slide: Any) -> None:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(3),
        )
        shape.name = "Boundary:test:primary"
        shape.fill.background()
        shape.line.color.rgb = RGBColor(10, 34, 51)

    report = inspect_generated(tmp_path, build)
    assert "PPTX_EMPTY_TEXT_BOX" not in issue_codes(report)


def test_empty_auto_shape_is_not_an_empty_text_box(tmp_path: Path) -> None:
    def build(slide: Any) -> None:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(3),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(228, 224, 215)
        shape.line.color.rgb = RGBColor(228, 224, 215)

    report = inspect_generated(tmp_path, build)
    assert "PPTX_EMPTY_TEXT_BOX" not in issue_codes(report)


def test_semantic_footer_uses_footnote_minimum(tmp_path: Path) -> None:
    def build(slide: Any) -> None:
        box = slide.shapes.add_textbox(
            Inches(0.6),
            Inches(7.0),
            Inches(8),
            Inches(0.2),
        )
        box.name = "source_note"
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "来源：测试材料"
        run.font.name = "Aptos"
        run.font.size = Pt(8.5)
        run.font.color.rgb = RGBColor(10, 34, 51)

    report = inspect_generated(tmp_path, build)
    assert "PPTX_FONT_SIZE_BELOW_MINIMUM" not in issue_codes(report)


def test_severe_text_overflow_is_a_hard_error(tmp_path: Path) -> None:
    def build(slide: Any) -> None:
        box = slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4.51),
            Inches(0.24),
        )
        box.name = "body_copy"
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "这是一个用于复现严重文本溢出的长句，文字在很矮的容器中必然无法完整显示，需要阻断交付。"
        run.font.name = "Aptos"
        run.font.size = Pt(17)
        run.font.color.rgb = RGBColor(10, 34, 51)

    report = inspect_generated(tmp_path, build)
    issue = next(
        item for item in report["issues"] if item["issue_code"] == "TEXT_OVERFLOW_RISK"
    )
    assert issue["severity"] == "error"
    assert report["status"] == "failed"


def test_short_single_line_label_is_not_estimated_as_overflow(tmp_path: Path) -> None:
    def build(slide: Any) -> None:
        box = slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(0.7),
            Inches(0.125),
        )
        box.name = "metric_label"
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "39%"
        run.font.name = "Aptos"
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(10, 34, 51)

    report = inspect_generated(tmp_path, build)
    assert "TEXT_OVERFLOW_RISK" not in issue_codes(report)


def test_low_contrast_text_over_filled_shape_is_a_hard_error(tmp_path: Path) -> None:
    def build(slide: Any) -> None:
        background = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(5),
            Inches(1.2),
        )
        background.name = "Material:dark_callout"
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(10, 34, 51)
        background.line.color.rgb = RGBColor(10, 34, 51)
        box = slide.shapes.add_textbox(
            Inches(1.3),
            Inches(1.3),
            Inches(4.4),
            Inches(0.6),
        )
        box.name = "closing_principle"
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "文字颜色不能与深色底色相同"
        run.font.name = "Aptos"
        run.font.size = Pt(15)
        run.font.color.rgb = RGBColor(10, 34, 51)

    report = inspect_generated(tmp_path, build)
    assert "TEXT_LOW_CONTRAST" in issue_codes(report)
    assert report["status"] == "failed"


def test_contrast_uses_topmost_fill_under_visible_left_aligned_text(
    tmp_path: Path,
) -> None:
    def build(slide: Any) -> None:
        background = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(5),
            Inches(0.5),
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(241, 239, 234)
        background.line.color.rgb = RGBColor(241, 239, 234)
        foreground = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(1.5),
            Inches(0.5),
        )
        foreground.fill.solid()
        foreground.fill.fore_color.rgb = RGBColor(44, 74, 110)
        foreground.line.color.rgb = RGBColor(44, 74, 110)
        box = slide.shapes.add_textbox(
            Inches(1.15),
            Inches(1.12),
            Inches(3),
            Inches(0.2),
        )
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "Scaling somewhere"
        run.font.name = "Aptos"
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(255, 255, 255)

    report = inspect_generated(tmp_path, build)
    assert "TEXT_LOW_CONTRAST" not in issue_codes(report)


def test_arrow_text_requires_safe_horizontal_inset(tmp_path: Path) -> None:
    def build(slide: Any) -> None:
        arrow = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON,
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(1),
        )
        arrow.name = "Component:roadmap:chevron"
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(44, 74, 110)
        arrow.line.color.rgb = RGBColor(44, 74, 110)
        box = slide.shapes.add_textbox(
            Inches(1.08),
            Inches(1.2),
            Inches(1.5),
            Inches(0.5),
        )
        box.name = "Component:roadmap:label"
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "文字落入箭头凹口"
        run.font.name = "Aptos"
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(255, 255, 255)

    report = inspect_generated(tmp_path, build)
    assert "TEXT_UNSAFE_SHAPE_INSET" in issue_codes(report)


def test_semantic_metric_wall_is_not_reported_as_text_fragmentation(
    tmp_path: Path,
) -> None:
    def build(slide: Any) -> None:
        for index in range(6):
            for role, text in (("value", f"{index + 1}%"), ("label", f"指标 {index + 1}")):
                box = slide.shapes.add_textbox(
                    Inches(0.7 + (index % 3) * 3.8),
                    Inches(1.0 + (index // 3) * 2.3 + (0 if role == "value" else 0.7)),
                    Inches(3.2),
                    Inches(0.5),
                )
                box.name = f"Component:metric_wall:{role}:{index + 1:02d}"
                run = box.text_frame.paragraphs[0].add_run()
                run.text = text
                run.font.name = "Aptos"
                run.font.size = Pt(24 if role == "value" else 11)
                run.font.color.rgb = RGBColor(10, 34, 51)

    report = inspect_generated(tmp_path, build)
    assert "PPTX_TEXT_FRAGMENTATION" not in issue_codes(report)


def test_color_drift_detected() -> None:
    report = inspect_fixture("color-drift")
    assert "STYLE_COLOR_DRIFT" in issue_codes(report)


def test_fragmentation_thresholds_fire() -> None:
    report = inspect_fixture("fragmented-svg-conversion")
    codes = issue_codes(report)
    assert "PPTX_TINY_OBJECT_OVERLOAD" in codes
    assert "PPTX_TEXT_FRAGMENTATION" in codes


def test_out_of_bounds_detected() -> None:
    report = inspect_fixture("object-out-of-bounds")
    assert {"PPTX_TEXT_OUT_OF_BOUNDS", "GEOMETRY_OBJECT_OUT_OF_BOUNDS"} & issue_codes(report)


@pytest.mark.parametrize(
    ("fixture", "expected"),
    [
        ("rasterized-table", "PPTX_TABLE_NOT_NATIVE"),
        ("rasterized-chart", "PPTX_CHART_NOT_NATIVE"),
        ("route-deviation", "PPTX_RASTER_ROUTE_UNDECLARED"),
    ],
)
def test_rasterized_route_fixtures_produce_expected_issues(fixture: str, expected: str) -> None:
    report = inspect_fixture(fixture)
    assert expected in issue_codes(report)
    assert "PPTX_RASTER_ROUTE_UNDECLARED" in issue_codes(report)


def run_cli(*args: str) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, str(ROOT / "scripts" / "inspect_pptx.py"), *args],
        cwd=str(ROOT),
        text=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=False,
    )


def test_cli_exit_codes_0_1_2(tmp_path: Path) -> None:
    valid = FIXTURES / "valid-native-deck"
    whole = FIXTURES / "whole-slide-image"
    result_ok = run_cli(
        str(valid / "deck.pptx"),
        "--ppt-ir",
        str(valid / "ppt-ir.json"),
        "--style",
        str(valid / "style-contract.json"),
        "--delivery",
        str(valid / "delivery-plan.json"),
        "--output",
        str(tmp_path / "valid.json"),
    )
    assert result_ok.returncode == 0, result_ok.stderr
    result_error = run_cli(
        str(whole / "deck.pptx"),
        "--ppt-ir",
        str(whole / "ppt-ir.json"),
        "--style",
        str(whole / "style-contract.json"),
        "--delivery",
        str(whole / "delivery-plan.json"),
        "--output",
        str(tmp_path / "whole.json"),
    )
    assert result_error.returncode == 1, result_error.stderr
    result_unreadable = run_cli(str(tmp_path / "missing.pptx"))
    assert result_unreadable.returncode == 2
