from __future__ import annotations

import json
import sys
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "scripts"))

from builders.base import BuildPlan
from builders.python_pptx_adapter import PythonPptxAdapter
from ppt_qa.package_inspector import inspect_package


def build_drill_down_fixture(tmp_path: Path) -> Path:
    style = json.loads(
        (ROOT / "references" / "template-packs" / "azure-insight-architecture.json").read_text(
            encoding="utf-8"
        )
    )
    plan = BuildPlan(
        builder="python_pptx",
        slides=[
            {
                "id": "S01",
                "title": "指标逐级穿透",
                "style_contract": style,
                "objects": [
                    {
                        "id": "drill-1",
                        "component_type": "drill_down_stair",
                        "editability": "native_required",
                        "delivery_plan": {"selected_route": "native_diagram"},
                        "content": {
                            "levels": [
                                {"label": "集团指标", "detail": "经营结果"},
                                {"label": "业务板块", "detail": "行业与区域"},
                                {"label": "业务流程", "detail": "过程指标"},
                                {"label": "业务单据", "detail": "原始记录"},
                            ]
                        },
                    }
                ],
            }
        ],
    )
    result = PythonPptxAdapter().build(plan, tmp_path)
    assert result.status == "created", result.errors
    assert result.fallbacks == []
    return Path(result.pptx)


def test_drill_down_stair_keeps_focus_path_continuous(tmp_path: Path) -> None:
    deck_path = build_drill_down_fixture(tmp_path)
    inspection = inspect_package(deck_path)
    deck = Presentation(deck_path)
    shapes = list(deck.slides[0].shapes)

    assert inspection.status == "passed"
    with zipfile.ZipFile(deck_path) as package:
        assert not any(name.startswith("ppt/media/") for name in package.namelist())
    assert sum(shape.name.startswith("Component:drill_down_stair:level:") for shape in shapes) == 4
    connectors = [
        shape
        for shape in shapes
        if shape.name.startswith("Connector:drill_down_stair:focus:")
    ]
    assert len(connectors) == 3
    for connector in connectors:
        xml = etree.tostring(connector.element).decode("utf-8")
        assert "stCxn" in xml and "endCxn" in xml
        assert connector.height < Inches(0.4)
