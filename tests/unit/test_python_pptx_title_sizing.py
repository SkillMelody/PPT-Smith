from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT))

from builders.base import BuildPlan
from builders.python_pptx_adapter import PythonPptxAdapter


def test_long_slide_title_receives_a_two_line_safe_textbox(tmp_path: Path) -> None:
    title = "高绩效组织在路线图、领导协同和敏捷交付上形成组合能力"
    result = PythonPptxAdapter().build(BuildPlan(builder="python_pptx", slides=[{
        "id": "S01", "title": title, "slide_role": "data", "style_contract": {"typography": {"title_sizes_pt": {"slide": 23}}}, "objects": [],
    }]), tmp_path)
    assert result.status == "created", result.errors
    deck = Presentation(result.pptx)
    title_shape = next(shape for shape in deck.slides[0].shapes if shape.text == title)
    assert title_shape.height >= Inches(0.78)
    run = title_shape.text_frame.paragraphs[0].runs[0]
    assert run.font.size.pt == 20.0
    assert run.font.name == "PingFang SC"
