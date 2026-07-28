from __future__ import annotations

import json
import sys
import zipfile
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "scripts"))

from builders.base import BuildPlan
from builders.python_pptx_adapter import PythonPptxAdapter
from ppt_qa.package_inspector import inspect_package


def build_phase_roadmap_fixture(tmp_path: Path) -> Path:
    style = json.loads(
        (ROOT / "references" / "template-packs" / "azure-insight-architecture.json").read_text(
            encoding="utf-8"
        )
    )
    plan = BuildPlan(
        builder="python_pptx",
        slides=[
            {
                "id": "S01",
                "title": "三阶段能力演进",
                "style_contract": style,
                "objects": [
                    {
                        "id": "roadmap-1",
                        "component_type": "phase_roadmap",
                        "editability": "native_required",
                        "delivery_plan": {"selected_route": "native_diagram"},
                        "content": {
                            "phases": [
                                {"label": "一期", "period": "Q1", "objective": "统一底座"},
                                {"label": "二期", "period": "Q2–Q3", "objective": "规模扩展"},
                                {"label": "三期", "period": "Q4", "objective": "价值闭环"},
                            ],
                            "milestones": [
                                {"label": "数据就绪", "phase": 0},
                                {"label": "试点上线", "phase": 0},
                                {"label": "规模推广", "phase": 1},
                                {"label": "经营验收", "phase": 2},
                            ],
                        },
                    }
                ],
            }
        ],
    )
    result = PythonPptxAdapter().build(plan, tmp_path)
    assert result.status == "created", result.errors
    assert result.fallbacks == []
    return Path(result.pptx)


def test_phase_roadmap_keeps_chronology_and_milestones_editable(tmp_path: Path) -> None:
    deck_path = build_phase_roadmap_fixture(tmp_path)
    inspection = inspect_package(deck_path)
    deck = Presentation(deck_path)
    shapes = list(deck.slides[0].shapes)

    assert inspection.status == "passed"
    with zipfile.ZipFile(deck_path) as package:
        assert not any(name.startswith("ppt/media/") for name in package.namelist())
    phases = sorted(
        (
            shape
            for shape in shapes
            if shape.name.startswith("Component:phase_roadmap:phase:")
        ),
        key=lambda shape: shape.left,
    )
    assert len(phases) == 3
    assert [shape.text.splitlines()[0] for shape in phases] == ["一期", "二期", "三期"]
    assert sum(shape.name.startswith("Component:phase_roadmap:milestone:") for shape in shapes) == 4
