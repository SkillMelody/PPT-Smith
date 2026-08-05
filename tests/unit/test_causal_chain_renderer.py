from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "scripts"))

from builders.base import BuildPlan
from builders.python_pptx_adapter import PythonPptxAdapter
from ppt_qa.package_inspector import inspect_package
from ppt_qa.slide_inspector import inspect_slides


def test_causal_chain_renders_a_native_editable_main_path(tmp_path: Path) -> None:
    style = json.loads((ROOT / "references" / "template-packs" / "azure-insight-architecture.json").read_text(encoding="utf-8"))
    plan = BuildPlan(builder="python_pptx", slides=[{
        "id": "S01", "title": "验证前置降低返工", "style_contract": style,
        "objects": [{
            "id": "causal-1", "component_type": "causal_chain", "editability": "native_required",
            "delivery_plan": {"selected_route": "native_diagram"},
            "content": {"nodes": [
                {"label": "前置验证", "role": "原因", "priority": "primary"},
                {"label": "减少返工", "role": "机制", "priority": "primary"},
                {"label": "可靠交付", "role": "结果", "priority": "primary"},
            ]},
        }],
    }])
    result = PythonPptxAdapter().build(plan, tmp_path)
    assert result.status == "created", result.errors
    deck_path = Path(result.pptx)
    package = inspect_package(deck_path)
    inspection = inspect_slides(deck_path, package, include_raw_xml=True)
    assert package.status == "passed"
    codes = {issue.issue_code for slide in inspection.slides for issue in slide.issues}
    assert "PPTX_CONNECTOR_THROUGH_NODE" not in codes
    deck = Presentation(deck_path)
    names = [shape.name for shape in deck.slides[0].shapes]
    assert sum(name.startswith("Component:causal_chain:node:") for name in names) == 3
    assert sum(name.startswith("Connector:causal_chain:main:") for name in names) == 2
