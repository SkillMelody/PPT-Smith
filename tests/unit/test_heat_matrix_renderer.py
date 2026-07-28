from __future__ import annotations

import json
import sys
import zipfile
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "scripts"))

from builders.base import BuildPlan
from builders.python_pptx_adapter import PythonPptxAdapter
from ppt_qa.package_inspector import inspect_package


def build_heat_matrix_fixture(tmp_path: Path) -> Path:
    style = json.loads(
        (ROOT / "references" / "template-packs" / "azure-insight-architecture.json").read_text(
            encoding="utf-8"
        )
    )
    plan = BuildPlan(
        builder="python_pptx",
        slides=[
            {
                "id": "S01",
                "title": "职能采用热力矩阵",
                "style_contract": style,
                "objects": [
                    {
                        "id": "heat-1",
                        "component_type": "heat_matrix",
                        "editability": "native_required",
                        "delivery_plan": {"selected_route": "native_diagram"},
                        "content": {
                            "rows": ["技术", "运营", "营销"],
                            "columns": ["制造", "零售", "金融", "医疗"],
                            "values": [
                                [24, 22, 18, 15],
                                [20, 17, 16, 12],
                                [14, 13, 11, 9],
                            ],
                            "value_suffix": "%",
                            "highlighted_cells": [[0, 0], [0, 1]],
                        },
                    }
                ],
            }
        ],
    )
    result = PythonPptxAdapter().build(plan, tmp_path)
    assert result.status == "created", result.errors
    assert result.fallbacks == []
    return Path(result.pptx)


def test_heat_matrix_creates_native_editable_cells(tmp_path: Path) -> None:
    deck_path = build_heat_matrix_fixture(tmp_path)
    inspection = inspect_package(deck_path)
    deck = Presentation(deck_path)
    shapes = list(deck.slides[0].shapes)

    assert inspection.status == "passed"
    with zipfile.ZipFile(deck_path) as package:
        assert not any(name.startswith("ppt/media/") for name in package.namelist())
    assert sum(shape.name.startswith("Component:heat_matrix:cell:") for shape in shapes) == 12
    texts = {shape.text for shape in shapes if hasattr(shape, "text_frame") and shape.text.strip()}
    assert {"技术", "运营", "24%", "22%"} <= texts
