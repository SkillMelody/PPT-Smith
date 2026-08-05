from __future__ import annotations

import json
import sys
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "scripts"))

from builders.base import BuildPlan
from builders.python_pptx_adapter import PythonPptxAdapter
from ppt_qa.package_inspector import inspect_package
from ppt_qa.slide_inspector import inspect_slides


def build_layered_architecture_fixture(
    tmp_path: Path,
    *,
    boundary_label: str | None = "统一技术平台",
) -> Path:
    style = json.loads(
        (ROOT / "references" / "template-packs" / "azure-insight-architecture.json").read_text(
            encoding="utf-8"
        )
    )
    plan = BuildPlan(
        builder="python_pptx",
        slides=[
            {
                "id": "S01",
                "title": "五层技术架构",
                "style_contract": style,
                "objects": [
                    {
                        "id": "arch-1",
                        "component_type": "layered_architecture",
                        "editability": "native_required",
                        "delivery_plan": {"selected_route": "native_diagram"},
                        "content": {
                            **(
                                {"boundary_label": boundary_label}
                                if boundary_label is not None
                                else {}
                            ),
                            "layers": [
                                {"label": "体验层", "nodes": ["门户", "移动端"]},
                                {"label": "业务层", "nodes": ["营销", "运营"]},
                                {"label": "能力层", "nodes": ["智能体", "知识库"]},
                                {"label": "数据层", "nodes": ["指标", "主数据"]},
                                {"label": "基础设施层", "nodes": ["计算", "安全"]},
                            ],
                        },
                    }
                ],
            }
        ],
    )
    result = PythonPptxAdapter().build(plan, tmp_path)
    assert result.status == "created", result.errors
    assert result.fallbacks == []
    return Path(result.pptx)


def test_layered_architecture_preserves_layers_boundaries_and_main_flow(tmp_path: Path) -> None:
    deck_path = build_layered_architecture_fixture(tmp_path)
    inspection = inspect_package(deck_path)
    slide_inspection = inspect_slides(deck_path, inspection, include_raw_xml=True)
    deck = Presentation(deck_path)
    shapes = list(deck.slides[0].shapes)

    assert inspection.status == "passed"
    codes = {
        issue.issue_code
        for slide_result in slide_inspection.slides
        for issue in slide_result.issues
    }
    assert "PPTX_FONT_SIZE_BELOW_MINIMUM" not in codes
    assert "PPTX_CONNECTOR_THROUGH_NODE" not in codes
    with zipfile.ZipFile(deck_path) as package:
        assert not any(name.startswith("ppt/media/") for name in package.namelist())
    assert sum(shape.name.startswith("Component:layered_architecture:layer:") for shape in shapes) == 5
    assert sum(shape.name.startswith("Boundary:") for shape in shapes) >= 1
    boundary_label = next(
        shape
        for shape in shapes
        if shape.name == "Component:layered_architecture:boundary_label"
    )
    first_layer = next(
        shape
        for shape in shapes
        if shape.name == "Component:layered_architecture:layer:0"
    )
    assert boundary_label.text == "统一技术平台"
    assert boundary_label.top + boundary_label.height <= first_layer.top
    connectors = [
        shape
        for shape in shapes
        if shape.name.startswith("Connector:layered_architecture:main:")
    ]
    assert len(connectors) == 4
    for connector in connectors:
        xml = etree.tostring(connector.element).decode("utf-8")
        assert "stCxn" in xml and "endCxn" in xml


def test_layered_architecture_omits_empty_boundary_label(tmp_path: Path) -> None:
    deck_path = build_layered_architecture_fixture(tmp_path, boundary_label=None)
    deck = Presentation(deck_path)
    names = {shape.name for shape in deck.slides[0].shapes}

    assert "Boundary:layered_architecture:primary" in names
    assert "Component:layered_architecture:boundary_label" not in names
