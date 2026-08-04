from __future__ import annotations

import json
import sys
import zipfile
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "scripts"))

from builders.base import BuildPlan
from builders.python_pptx_adapter import PythonPptxAdapter
from ppt_qa.package_inspector import inspect_package


def _style() -> dict[str, object]:
    return json.loads(
        (ROOT / "references" / "template-packs" / "azure-insight-architecture.json").read_text(
            encoding="utf-8"
        )
    )


def test_product_ui_overview_is_source_labeled_native_workspace_sketch(tmp_path: Path) -> None:
    result = PythonPptxAdapter().build(
        BuildPlan(
            builder="python_pptx",
            slides=[
                {
                    "id": "S01",
                    "title": "工作台 UI 总览",
                    "style_contract": _style(),
                    "objects": [
                        {
                            "id": "prd-ui-overview-1",
                            "component_type": "product_ui_overview",
                            "editability": "native_required",
                            "delivery_plan": {"selected_route": "native_diagram"},
                            "content": {
                                "source_refs": ["prd#workspace-layout"],
                                "spaces": [
                                    {"id": "topbar", "label": "TopBar"},
                                    {"id": "canvas", "label": "Canvas"},
                                    {"id": "inspector", "label": "Inspector"},
                                ],
                                "relationships": [
                                    {"from": "topbar", "to": "canvas", "label": "进入工作区"},
                                    {"from": "canvas", "to": "inspector", "label": "选区联动"},
                                ],
                            },
                        }
                    ],
                }
            ],
        ),
        tmp_path,
    )

    assert result.status == "created", result.errors
    assert result.fallbacks == []
    deck_path = Path(result.pptx)
    assert inspect_package(deck_path).status == "passed"
    with zipfile.ZipFile(deck_path) as package:
        assert not any(name.startswith("ppt/media/") for name in package.namelist())

    shapes = list(Presentation(deck_path).slides[0].shapes)
    names = {shape.name for shape in shapes}
    assert "Label:product_ui_overview:derived" in names
    assert "Source:product_ui_overview:0" in names
    assert sum(name.startswith("Component:product_ui_overview:space:") for name in names) == 3
    assert sum(name.startswith("Connector:product_ui_overview:relationship:") for name in names) == 2
    assert all(shape.shape_type != 13 for shape in shapes)


def test_interaction_storyboard_is_source_labeled_native_state_flow(tmp_path: Path) -> None:
    result = PythonPptxAdapter().build(
        BuildPlan(
            builder="python_pptx",
            slides=[
                {
                    "id": "S01",
                    "title": "协作交互路径",
                    "style_contract": _style(),
                    "objects": [
                        {
                            "id": "prd-storyboard-1",
                            "component_type": "interaction_storyboard",
                            "editability": "native_required",
                            "delivery_plan": {"selected_route": "native_diagram"},
                            "content": {
                                "source_refs": ["prd#streaming-diff"],
                                "steps": [
                                    {"id": "select", "kind": "action", "label": "选区并触发"},
                                    {"id": "stream", "kind": "async_state", "label": "流式生成中"},
                                    {"id": "diff", "kind": "result", "label": "查看 Diff"},
                                    {"id": "adopt", "kind": "decision", "label": "采纳或撤销"},
                                ],
                                "transitions": [
                                    {"from": "select", "to": "stream"},
                                    {"from": "stream", "to": "diff"},
                                    {"from": "diff", "to": "adopt"},
                                ],
                            },
                        }
                    ],
                }
            ],
        ),
        tmp_path,
    )

    assert result.status == "created", result.errors
    assert result.fallbacks == []
    deck_path = Path(result.pptx)
    assert inspect_package(deck_path).status == "passed"
    with zipfile.ZipFile(deck_path) as package:
        assert not any(name.startswith("ppt/media/") for name in package.namelist())

    shapes = list(Presentation(deck_path).slides[0].shapes)
    names = {shape.name for shape in shapes}
    assert "Label:interaction_storyboard:derived" in names
    assert "Source:interaction_storyboard:0" in names
    assert sum(name.startswith("Component:interaction_storyboard:step:") for name in names) == 4
    assert sum(name.startswith("Connector:interaction_storyboard:transition:") for name in names) == 3
    assert all(shape.shape_type != 13 for shape in shapes)
