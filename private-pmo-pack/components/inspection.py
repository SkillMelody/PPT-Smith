"""Static inspection of editable PMO component PPTX packages."""

from __future__ import annotations

from pathlib import Path
from typing import Any
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def inspect_component_pptx(
    path: Path,
    *,
    expected_slide_count: int,
    expected_prefixes: list[str],
    minimum_native_shapes: int,
    require_chart: bool = False,
    require_table: bool = False,
    require_object_bounds: bool = False,
) -> dict[str, Any]:
    if not path.exists() or path.stat().st_size == 0:
        raise FileNotFoundError(path)
    prs = Presentation(path)
    counts = {
        "slide_count": len(prs.slides),
        "picture_count": 0,
        "chart_count": 0,
        "table_count": 0,
        "connector_count": 0,
        "text_shape_count": 0,
        "native_shape_count": 0,
        "total_shape_count": 0,
    }
    shape_names: list[str] = []
    outside_shape_names: list[str] = []
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    for slide in prs.slides:
        for shape in slide.shapes:
            counts["total_shape_count"] += 1
            shape_names.append(shape.name)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                counts["picture_count"] += 1
            else:
                counts["native_shape_count"] += 1
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                counts["connector_count"] += 1
            if getattr(shape, "has_chart", False):
                counts["chart_count"] += 1
            if getattr(shape, "has_table", False):
                counts["table_count"] += 1
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text:
                counts["text_shape_count"] += 1
            x1 = min(int(shape.left), int(shape.left + shape.width))
            x2 = max(int(shape.left), int(shape.left + shape.width))
            y1 = min(int(shape.top), int(shape.top + shape.height))
            y2 = max(int(shape.top), int(shape.top + shape.height))
            if x1 < -1000 or y1 < -1000 or x2 > slide_width + 1000 or y2 > slide_height + 1000:
                outside_shape_names.append(shape.name)

    with ZipFile(path) as archive:
        members = archive.namelist()
    media_files = [name for name in members if name.startswith("ppt/media/") and not name.endswith("/")]
    chart_parts = [name for name in members if name.startswith("ppt/charts/chart") and name.endswith(".xml")]
    workbook_parts = [name for name in members if name.startswith("ppt/embeddings/") and not name.endswith("/")]
    missing_prefixes = [prefix for prefix in expected_prefixes if not any(name.startswith(prefix) for name in shape_names)]
    gates = {
        "slide_count_matches": counts["slide_count"] == expected_slide_count,
        "no_picture_shapes": counts["picture_count"] == 0,
        "no_media_files": len(media_files) == 0,
        "native_shapes_present": counts["native_shape_count"] >= minimum_native_shapes,
        "semantic_shape_names_present": not missing_prefixes,
    }
    if require_chart:
        gates["native_chart_present"] = counts["chart_count"] >= 1 and len(chart_parts) >= 1
    if require_table:
        gates["native_table_present"] = counts["table_count"] >= 1
    if require_object_bounds:
        gates["all_objects_within_slide"] = not outside_shape_names
    return {
        "path": str(path),
        **counts,
        "media_file_count": len(media_files),
        "chart_part_count": len(chart_parts),
        "embedded_workbook_count": len(workbook_parts),
        "missing_component_prefixes": missing_prefixes,
        "outside_shape_names": outside_shape_names,
        "gates": gates,
        "passed": all(gates.values()),
        "limitations": [
            "python-pptx emits native connectors but does not expose a reliable API to prove endpoint binding.",
            "Static package inspection cannot replace rendered visual review.",
        ],
    }


def inspect_phase1_pptx(path: Path) -> dict[str, Any]:
    return inspect_component_pptx(
        path,
        expected_slide_count=3,
        expected_prefixes=["doughnut.", "roadmap.", "gantt."],
        minimum_native_shapes=41,
        require_chart=True,
    )


def inspect_portfolio_doughnut_pptx(path: Path) -> dict[str, Any]:
    return inspect_component_pptx(
        path,
        expected_slide_count=1,
        expected_prefixes=[
            "doughnut.native_chart",
            "doughnut.leader.",
            "doughnut.legend_fallback.",
            "doughnut.readout.",
        ],
        minimum_native_shapes=35,
        require_chart=True,
        require_object_bounds=True,
    )


def inspect_phase2_p0_pptx(path: Path) -> dict[str, Any]:
    return inspect_component_pptx(
        path,
        expected_slide_count=4,
        expected_prefixes=["kpi.", "milestone.", "heatmap.", "raci."],
        minimum_native_shapes=90,
        require_table=True,
    )
