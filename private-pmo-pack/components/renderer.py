"""Native python-pptx renderers for PMO component-system phases."""

from __future__ import annotations

from datetime import date, timedelta
from pathlib import Path
from typing import Any, Iterable, Mapping

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from components.ir import validate_component_ir
from components.layout import (
    layout_doughnut,
    layout_gantt,
    layout_kpi_status,
    layout_milestone,
    layout_raci,
    layout_risk_heat_map,
    layout_roadmap,
)
from components.qa import check_text_bounds, check_text_overflow, legend_completeness, qa_component
from components.tokens import get_token


DEFAULT_PHASE1_IRS: list[dict[str, Any]] = [
    {
        "type": "annotated_doughnut",
        "title": "Investment mix is concentrated on platform and data foundations",
        "center_metric": "100%",
        "center_caption": "portfolio",
        "categories": [
            {"label": "Platform", "value": 38, "color_role": "sequence.phase_1"},
            {"label": "Data", "value": 27, "color_role": "sequence.phase_2"},
            {"label": "Adoption", "value": 21, "color_role": "sequence.phase_3"},
            {"label": "Controls", "value": 14, "color_role": "sequence.phase_4"},
        ],
    },
    {
        "type": "roadmap",
        "title": "Four gated phases convert strategy into an auditable scale decision",
        "phases": [
            {
                "id": "P1",
                "name": "Discover",
                "period": "Jul 2026",
                "objective": "Align value, scope, and governance.",
                "outputs": ["Project charter", "Value baseline"],
                "exit_condition": "Scope and value baseline signed",
                "gate": "G1",
            },
            {
                "id": "P2",
                "name": "Design",
                "period": "Aug 2026",
                "objective": "Define target process and solution.",
                "outputs": ["Solution blueprint", "Control design"],
                "exit_condition": "Design authority approval",
                "gate": "G2",
            },
            {
                "id": "P3",
                "name": "Pilot",
                "period": "Sep-Oct 2026",
                "objective": "Prove operational and commercial evidence.",
                "outputs": ["Pilot release", "Adoption evidence"],
                "exit_condition": "Pilot thresholds achieved",
                "gate": "G3",
            },
            {
                "id": "P4",
                "name": "Scale",
                "period": "Nov-Dec 2026",
                "objective": "Industrialize delivery and benefits tracking.",
                "outputs": ["Scale runbook", "Benefits dashboard"],
                "exit_condition": "SteerCo scale decision",
                "gate": "G4",
            },
        ],
    },
    {
        "type": "gantt",
        "title": "Pilot schedule remains on plan, with data migration on watch",
        "start_date": "2026-07-01",
        "end_date": "2026-09-30",
        "reference_date": "2026-08-17",
        "tasks": [
            {
                "id": "T1",
                "label": "Governance setup",
                "owner": "PMO",
                "planned_intervals": [{"start": "2026-07-01", "end": "2026-07-24"}],
                "actual_intervals": [{"start": "2026-07-01", "end": "2026-07-22"}],
                "health": "healthy",
                "milestones": [{"date": "2026-07-24", "label": "Charter"}],
            },
            {
                "id": "T2",
                "label": "Solution design",
                "owner": "Architecture",
                "planned_intervals": [{"start": "2026-07-13", "end": "2026-08-21"}],
                "actual_intervals": [{"start": "2026-07-13", "end": "2026-08-14"}],
                "health": "healthy",
                "milestones": [{"date": "2026-08-21", "label": "Design"}],
            },
            {
                "id": "T3",
                "label": "Data migration",
                "owner": "Data",
                "planned_intervals": [{"start": "2026-07-27", "end": "2026-09-11"}],
                "actual_intervals": [{"start": "2026-07-27", "end": "2026-08-12"}],
                "health": "watch",
                "milestones": [],
            },
            {
                "id": "T4",
                "label": "Pilot build",
                "owner": "Product",
                "planned_intervals": [{"start": "2026-08-10", "end": "2026-09-18"}],
                "actual_intervals": [{"start": "2026-08-10", "end": "2026-08-15"}],
                "health": "healthy",
                "milestones": [{"date": "2026-09-18", "label": "Release"}],
            },
            {
                "id": "T5",
                "label": "Adoption readiness",
                "owner": "Change",
                "planned_intervals": [{"start": "2026-08-24", "end": "2026-09-30"}],
                "actual_intervals": [],
                "health": "neutral",
                "milestones": [{"date": "2026-09-30", "label": "Go-live"}],
            },
        ],
    },
]


DEFAULT_PHASE2_P0_IRS: list[dict[str, Any]] = [
    {
        "type": "kpi_status",
        "title": "Value delivery is on plan while adoption remains on watch",
        "metrics": [
            {"id": "K1", "label": "Run-rate value", "value": 8.4, "unit": "$m", "target": 8.0, "direction": "higher_is_better", "status": "healthy", "period": "Q3 2026", "note": "Validated by Finance", "trend_values": [6.1, 6.8, 7.4, 8.4]},
            {"id": "K2", "label": "Adoption", "value": 74, "unit": "%", "target": 80, "direction": "higher_is_better", "status": "watch", "period": "Aug 2026", "note": "Six points below target", "trend_values": [51, 60, 68, 74]},
            {"id": "K3", "label": "Cycle time", "value": 12.3, "unit": "days", "target": 10, "direction": "lower_is_better", "status": "risk", "period": "Aug 2026", "note": "Data handoffs constrain flow", "trend_values": [15.2, 14.1, 13.0, 12.3]},
            {"id": "K4", "label": "Control pass rate", "value": 96, "unit": "%", "target": 95, "direction": "higher_is_better", "status": "healthy", "period": "Aug 2026", "note": "Sample size 412"},
        ],
    },
    {
        "type": "milestone",
        "title": "Six evidence-backed milestones define the path to scale approval",
        "start_date": "2026-07-01",
        "end_date": "2026-12-15",
        "milestones": [
            {"id": "M1", "label": "Charter approved", "date": "2026-07-08", "owner": "PMO", "evidence": "Signed charter", "status": "healthy", "phase": "Discover"},
            {"id": "M2", "label": "Blueprint signed", "date": "2026-08-05", "owner": "Architecture", "evidence": "Design authority minute", "status": "healthy", "phase": "Design"},
            {"id": "M3", "label": "Data rehearsal", "date": "2026-09-02", "owner": "Data", "evidence": "Reconciliation pack", "status": "watch", "phase": "Pilot"},
            {"id": "M4", "label": "Pilot release", "date": "2026-10-01", "owner": "Product", "evidence": "Release record", "status": "healthy", "phase": "Pilot"},
            {"id": "M5", "label": "Adoption threshold", "date": "2026-11-05", "owner": "Change", "evidence": "Usage dashboard", "status": "watch", "phase": "Scale"},
            {"id": "M6", "label": "Scale decision", "date": "2026-12-10", "owner": "SteerCo", "evidence": "Decision log", "status": "neutral", "phase": "Scale"},
        ],
    },
    {
        "type": "risk_heat_map",
        "title": "Data readiness concentrates the critical exposure in two cells",
        "risks": [
            {"id": "R1", "label": "Source data quality", "likelihood": 5, "impact": 5, "owner": "Data", "status": "risk", "mitigation": "Daily cleanse burn-down"},
            {"id": "R2", "label": "Interface latency", "likelihood": 4, "impact": 5, "owner": "Engineering", "status": "risk", "mitigation": "Performance war room"},
            {"id": "R3", "label": "SME availability", "likelihood": 4, "impact": 4, "owner": "Operations", "status": "watch", "mitigation": "Named alternates"},
            {"id": "R4", "label": "Training completion", "likelihood": 4, "impact": 4, "owner": "Change", "status": "watch", "mitigation": "Manager escalation"},
            {"id": "R5", "label": "Vendor access", "likelihood": 4, "impact": 4, "owner": "Security", "status": "watch", "mitigation": "Pre-approved access window"},
            {"id": "R6", "label": "Benefits attribution", "likelihood": 2, "impact": 3, "owner": "Finance", "status": "neutral", "mitigation": "Reconcile baseline monthly"},
        ],
    },
    {
        "type": "raci",
        "title": "Single-point accountability is explicit across the pilot governance chain",
        "roles": ["Sponsor", "PMO", "Product", "Data", "Change"],
        "activities": ["Approve charter", "Sign blueprint", "Release pilot", "Validate benefits", "Approve scale"],
        "assignments": [
            {"activity": "Approve charter", "role": "Sponsor", "code": "A"}, {"activity": "Approve charter", "role": "PMO", "code": "R"}, {"activity": "Approve charter", "role": "Product", "code": "C"},
            {"activity": "Sign blueprint", "role": "Sponsor", "code": "A"}, {"activity": "Sign blueprint", "role": "Product", "code": "R"}, {"activity": "Sign blueprint", "role": "Data", "code": "C"},
            {"activity": "Release pilot", "role": "Sponsor", "code": "A"}, {"activity": "Release pilot", "role": "Product", "code": "R"}, {"activity": "Release pilot", "role": "Change", "code": "C"},
            {"activity": "Validate benefits", "role": "Sponsor", "code": "A"}, {"activity": "Validate benefits", "role": "PMO", "code": "R"}, {"activity": "Validate benefits", "role": "Data", "code": "C"},
            {"activity": "Approve scale", "role": "Sponsor", "code": "A"}, {"activity": "Approve scale", "role": "PMO", "code": "R"}, {"activity": "Approve scale", "role": "Product", "code": "C"}, {"activity": "Approve scale", "role": "Change", "code": "I"},
        ],
    },
]


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#"))


def _role_color(tokens: Mapping[str, Any], role: str) -> str:
    return str(get_token(tokens, f"semantic.{role}"))


class NativeComponentDeck:
    def __init__(self, tokens: Mapping[str, Any]) -> None:
        self.tokens = tokens
        self.prs = Presentation()
        self.prs.slide_width = Inches(float(tokens["slide"]["width_in"]))
        self.prs.slide_height = Inches(float(tokens["slide"]["height_in"]))
        self.blank_layout = self.prs.slide_layouts[6]
        self.component_records: list[dict[str, Any]] = []
        self.text_boxes: list[dict[str, Any]] = []

    def add_slide(self, title: str, kicker: str, component: str) -> Any:
        slide = self.prs.slides.add_slide(self.blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(_role_color(self.tokens, "surface.canvas"))
        self.text(slide, kicker.upper(), 0.56, 0.24, 4.8, 0.22, role="label", color_role="sequence.phase_3", bold=True, name=f"{component}.kicker")
        self.text(slide, title, 0.56, 0.52, 11.95, 0.78, role="slide_title", color_role="text.title", bold=True, name=f"{component}.title")
        self.line(slide, 0.56, 1.34, 12.78, 1.34, color_role="line.divider", width_pt=0.65, name=f"{component}.title_rule")
        self.text(
            slide,
            "BOARDROOM INK / NATIVE EDITABLE",
            9.75,
            7.10,
            3.0,
            0.18,
            role="footnote",
            color_role="text.muted",
            bold=True,
            align=PP_ALIGN.RIGHT,
            name=f"{component}.footer",
        )
        return slide

    def text(
        self,
        slide: Any,
        value: str,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        role: str = "body",
        color_role: str = "text.body",
        bold: bool | None = None,
        font_size_pt: float | None = None,
        align: PP_ALIGN = PP_ALIGN.LEFT,
        valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
        name: str,
    ) -> Any:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        shape.name = name
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.03)
        frame.margin_right = Inches(0.03)
        frame.margin_top = Inches(0.01)
        frame.margin_bottom = Inches(0.01)
        frame.vertical_anchor = valign
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = str(value)
        run.font.name = str(self.tokens["typography"]["font"])
        resolved_font_pt = (
            float(self.tokens["typography"]["roles"][role]["size_pt"])
            if font_size_pt is None
            else float(font_size_pt)
        )
        run.font.size = Pt(resolved_font_pt)
        run.font.bold = self.tokens["typography"]["roles"][role]["bold"] if bold is None else bold
        run.font.color.rgb = _rgb(_role_color(self.tokens, color_role))
        self.text_boxes.append(
            {
                "name": name,
                "text": str(value),
                "x": x,
                "y": y,
                "w": w,
                "h": h,
                "font_pt": resolved_font_pt,
            }
        )
        return shape

    def shape(
        self,
        slide: Any,
        kind: MSO_AUTO_SHAPE_TYPE,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        fill_role: str,
        line_role: str = "line.divider",
        line_width_pt: float = 0.65,
        name: str,
    ) -> Any:
        shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.name = name
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(_role_color(self.tokens, fill_role))
        shape.line.color.rgb = _rgb(_role_color(self.tokens, line_role))
        shape.line.width = Pt(line_width_pt)
        return shape

    def line(
        self,
        slide: Any,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        *,
        color_role: str,
        width_pt: float,
        name: str,
        elbow: bool = False,
    ) -> Any:
        connector_type = MSO_CONNECTOR.ELBOW if elbow else MSO_CONNECTOR.STRAIGHT
        connector = slide.shapes.add_connector(connector_type, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        connector.name = name
        connector.line.color.rgb = _rgb(_role_color(self.tokens, color_role))
        connector.line.width = Pt(width_pt)
        return connector

    def save(self, output: Path) -> None:
        output.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output)


def _render_doughnut(deck: NativeComponentDeck, ir: Mapping[str, Any]) -> None:
    slide = deck.add_slide(ir["title"], "Annotated doughnut", "doughnut")
    frame = {"x": 0.66, "y": 1.38, "w": 12.0, "h": 5.35}
    layout = layout_doughnut(ir, frame, deck.tokens)
    center = layout["center"]
    radius = layout["radius"]
    chart_data = ChartData()
    chart_data.categories = [category["label"] for category in ir["categories"]]
    chart_data.add_series("Share", [category["value"] for category in ir["categories"]])
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(center["x"] - radius),
        Inches(center["y"] - radius),
        Inches(radius * 2),
        Inches(radius * 2),
        chart_data,
    )
    chart_frame.name = "doughnut.native_chart"
    chart = chart_frame.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.doughnut_hole_size = 62
    plot.has_data_labels = False
    series = chart.series[0]
    for index, category in enumerate(ir["categories"]):
        point = series.points[index]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = _rgb(_role_color(deck.tokens, category["color_role"]))
        point.format.line.color.rgb = _rgb(_role_color(deck.tokens, "surface.canvas"))
        point.format.line.width = Pt(1.0)

    deck.text(
        slide,
        ir.get("center_metric", ""),
        center["x"] - 0.80,
        center["y"] - 0.40,
        1.60,
        0.50,
        role="metric",
        color_role="text.title",
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        name="doughnut.center_metric",
    )
    deck.text(
        slide,
        ir.get("center_caption", ""),
        center["x"] - 0.58,
        center["y"] + 0.22,
        1.16,
        0.22,
        role="footnote",
        color_role="text.muted",
        align=PP_ALIGN.CENTER,
        name="doughnut.center_caption",
    )

    if layout["mode"] == "annotated":
        for index, (label, leader) in enumerate(zip(layout["labels"], layout["leaders"])):
            last = leader["points"][-1]
            anchor = leader["points"][0]
            deck.line(
                slide,
                anchor["x"],
                anchor["y"],
                last["x"],
                last["y"],
                color_role="line.annotation",
                width_pt=1.5,
                name=f"doughnut.leader.{index + 1}",
                elbow=True,
            )
            deck.shape(
                slide,
                MSO_AUTO_SHAPE_TYPE.OVAL,
                anchor["x"] - 0.035,
                anchor["y"] - 0.035,
                0.07,
                0.07,
                fill_role=label["color_role"],
                line_role=label["color_role"],
                name=f"doughnut.anchor.{index + 1}",
            )
            value = f"{label['label']}  {label['value']:.0f}%"
            deck.text(
                slide,
                value,
                label["x"],
                label["y"],
                label["w"],
                label["h"],
                role="label",
                color_role="text.body",
                bold=True,
                align=PP_ALIGN.RIGHT if label["side"] == "left" else PP_ALIGN.LEFT,
                valign=MSO_ANCHOR.MIDDLE,
                name=f"doughnut.label.{index + 1}",
            )
    else:
        legend_x = frame["x"] + frame["w"] * 0.67
        for index, item in enumerate(layout["legend"]):
            y = frame["y"] + 0.65 + index * 0.52
            deck.shape(
                slide,
                MSO_AUTO_SHAPE_TYPE.OVAL,
                legend_x,
                y + 0.05,
                0.12,
                0.12,
                fill_role=item["color_role"],
                line_role=item["color_role"],
                name=f"doughnut.legend_dot.{index + 1}",
            )
            deck.text(
                slide,
                f"{item['label']}  {item['value']:.0f}%",
                legend_x + 0.22,
                y,
                3.25,
                0.24,
                role="label",
                color_role="text.body",
                name=f"doughnut.legend_label.{index + 1}",
            )

    deck.text(
        slide,
        "Native chart data remains editable; labels and leaders are native PowerPoint objects.",
        0.72,
        6.54,
        8.2,
        0.2,
        role="footnote",
        color_role="text.muted",
        name="doughnut.note",
    )
    deck.component_records.append(
        {
            "component": "annotated_doughnut",
            "layout_mode": layout["mode"],
            "reason_codes": layout["reason_codes"],
            "leader_count": len(layout["leaders"]),
            "chart_first_slice_angle_deg": layout["chart_first_slice_angle_deg"],
            "endpoint_segment_alignment_passed": all(
                segment["start_angle"] <= segment["endpoint_angle"] <= segment["end_angle"]
                for segment in layout["segments"]
            ),
        }
    )


def _render_portfolio_doughnut(deck: NativeComponentDeck, ir: Mapping[str, Any]) -> None:
    """Render the SteerCo poster variant using the shared doughnut layout."""

    slide = deck.add_slide(ir["title"], "Portfolio investment / SteerCo view", "doughnut")
    frame = {"x": 0.68, "y": 1.46, "w": 8.15, "h": 4.88}
    layout = layout_doughnut(ir, frame, deck.tokens)
    if layout["mode"] != "annotated":
        raise ValueError(f"portfolio doughnut requires annotated mode: {layout['reason_codes']}")
    center = layout["center"]
    radius = layout["radius"]

    chart_data = ChartData()
    chart_data.categories = [category["label"] for category in ir["categories"]]
    chart_data.add_series("Investment share", [category["value"] for category in ir["categories"]])
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(center["x"] - radius),
        Inches(center["y"] - radius),
        Inches(radius * 2),
        Inches(radius * 2),
        chart_data,
    )
    chart_frame.name = "doughnut.native_chart"
    chart = chart_frame.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.doughnut_hole_size = 64
    plot.has_data_labels = False
    series = chart.series[0]
    for index, category in enumerate(ir["categories"]):
        point = series.points[index]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = _rgb(_role_color(deck.tokens, category["color_role"]))
        point.format.line.color.rgb = _rgb(_role_color(deck.tokens, "surface.canvas"))
        point.format.line.width = Pt(1.5)

    deck.text(
        slide,
        ir["center_metric"],
        center["x"] - 0.88,
        center["y"] - 0.42,
        1.76,
        0.54,
        role="metric",
        font_size_pt=25,
        color_role="text.title",
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        name="doughnut.center_metric",
    )
    deck.text(
        slide,
        ir["center_caption"],
        center["x"] - 0.76,
        center["y"] + 0.22,
        1.52,
        0.24,
        role="footnote",
        font_size_pt=8.4,
        color_role="text.muted",
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        name="doughnut.center_caption",
    )

    # Leaders are explicit connected polyline segments so the rendered route
    # exactly matches deterministic layout geometry.
    for leader_index, (label, leader) in enumerate(zip(layout["labels"], layout["leaders"]), start=1):
        category = ir["categories"][label["category_index"]]
        points = leader["points"]
        for segment_index, (start, end) in enumerate(zip(points, points[1:]), start=1):
            deck.line(
                slide,
                start["x"],
                start["y"],
                end["x"],
                end["y"],
                color_role="line.annotation",
                width_pt=1.15,
                name=f"doughnut.leader.{leader_index}.segment.{segment_index}",
            )
        anchor = points[0]
        deck.shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.OVAL,
            anchor["x"] - 0.042,
            anchor["y"] - 0.042,
            0.084,
            0.084,
            fill_role=label["color_role"],
            line_role="surface.canvas",
            line_width_pt=0.7,
            name=f"doughnut.anchor.{leader_index}",
        )
        align = PP_ALIGN.RIGHT if label["side"] == "left" else PP_ALIGN.LEFT
        metric_text = f"{label['value']:.0f}%  ·  ¥{float(category['amount']):.2f}M"
        deck.text(
            slide,
            category["label"],
            label["x"],
            label["y"],
            label["w"],
            0.20,
            role="label",
            font_size_pt=9.8,
            color_role="text.title",
            bold=True,
            align=align,
            name=f"doughnut.label.{leader_index}.category",
        )
        deck.text(
            slide,
            metric_text,
            label["x"],
            label["y"] + 0.22,
            label["w"],
            0.28,
            role="section_title",
            font_size_pt=14.2,
            color_role="text.title",
            bold=True,
            align=align,
            name=f"doughnut.label.{leader_index}.metric",
        )
        status_dot_x = (
            label["x"] + label["w"] - 0.09
            if label["side"] == "left"
            else label["x"]
        )
        status_text_x = (
            label["x"]
            if label["side"] == "left"
            else label["x"] + 0.15
        )
        status_text_w = label["w"] - 0.15
        deck.shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.OVAL,
            status_dot_x,
            label["y"] + 0.56,
            0.075,
            0.075,
            fill_role=f"state.{category['status_role']}",
            line_role=f"state.{category['status_role']}",
            name=f"doughnut.label.{leader_index}.status_dot",
        )
        deck.text(
            slide,
            category["status_text"],
            status_text_x,
            label["y"] + 0.50,
            status_text_w,
            0.18,
            role="footnote",
            font_size_pt=8.2,
            color_role="text.muted",
            align=align,
            name=f"doughnut.label.{leader_index}.status",
        )

    # The fifth/small category intentionally falls back to one compact band.
    fallback = layout["legend_fallback"]
    if fallback:
        item = fallback[0]
        category = ir["categories"][item["index"]]
        deck.shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            0.70,
            6.40,
            8.13,
            0.53,
            fill_role="surface.paper",
            line_role="line.soft",
            line_width_pt=0.65,
            name="doughnut.legend_fallback.band",
        )
        deck.shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.OVAL,
            0.92,
            6.59,
            0.11,
            0.11,
            fill_role=item["color_role"],
            line_role=item["color_role"],
            name="doughnut.legend_fallback.category_dot",
        )
        deck.text(
            slide,
            category["label"],
            1.13,
            6.50,
            1.45,
            0.20,
            role="label",
            font_size_pt=9.2,
            color_role="text.title",
            bold=True,
            name="doughnut.legend_fallback.label",
        )
        deck.text(
            slide,
            f"{item['value']:.0f}%  ·  ¥{float(category['amount']):.2f}M",
            2.55,
            6.49,
            1.65,
            0.21,
            role="label",
            font_size_pt=10.2,
            color_role="text.title",
            bold=True,
            name="doughnut.legend_fallback.metric",
        )
        deck.shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.OVAL,
            4.43,
            6.59,
            0.075,
            0.075,
            fill_role=f"state.{category['status_role']}",
            line_role=f"state.{category['status_role']}",
            name="doughnut.legend_fallback.status_dot",
        )
        deck.text(
            slide,
            category["status_text"],
            4.57,
            6.51,
            1.35,
            0.18,
            role="footnote",
            font_size_pt=8.2,
            color_role="text.muted",
            name="doughnut.legend_fallback.status",
        )
        deck.text(
            slide,
            "Displayed amounts total ¥8.60M; each category is rounded to ¥0.01M.",
            5.66,
            6.47,
            2.90,
            0.28,
            role="footnote",
            font_size_pt=7.2,
            color_role="text.muted",
            align=PP_ALIGN.RIGHT,
            name="doughnut.legend_fallback.rounding_note",
        )

    # One restrained premium panel, not nested cards.
    panel_x, panel_y, panel_w, panel_h = 9.15, 1.55, 3.56, 5.37
    deck.shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        panel_x,
        panel_y,
        panel_w,
        panel_h,
        fill_role="surface.paper",
        line_role="line.soft",
        line_width_pt=0.75,
        name="doughnut.readout.panel",
    )
    deck.shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        panel_x,
        panel_y,
        0.07,
        0.72,
        fill_role="sequence.phase_3",
        line_role="sequence.phase_3",
        name="doughnut.readout.accent",
    )
    deck.text(
        slide,
        "PMO READOUT",
        panel_x + 0.28,
        panel_y + 0.25,
        2.3,
        0.22,
        role="label",
        font_size_pt=9.4,
        color_role="sequence.phase_3",
        bold=True,
        name="doughnut.readout.title",
    )
    deck.text(
        slide,
        "Where governance attention should move next",
        panel_x + 0.28,
        panel_y + 0.57,
        2.88,
        0.52,
        role="section_title",
        font_size_pt=15.0,
        color_role="text.title",
        bold=True,
        name="doughnut.readout.subtitle",
    )
    readout_top = panel_y + 1.35
    readout_gap = 1.22
    for index, readout in enumerate(ir.get("readouts", []), start=1):
        y = readout_top + (index - 1) * readout_gap
        deck.text(
            slide,
            f"0{index}",
            panel_x + 0.28,
            y,
            0.38,
            0.22,
            role="label",
            font_size_pt=8.4,
            color_role="text.muted",
            bold=True,
            name=f"doughnut.readout.{index}.number",
        )
        deck.shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.OVAL,
            panel_x + 0.73,
            y + 0.06,
            0.09,
            0.09,
            fill_role=f"state.{readout['status_role']}",
            line_role=f"state.{readout['status_role']}",
            name=f"doughnut.readout.{index}.status_dot",
        )
        deck.text(
            slide,
            readout["headline"],
            panel_x + 0.92,
            y - 0.03,
            2.33,
            0.44,
            role="body",
            font_size_pt=10.6,
            color_role="text.title",
            bold=True,
            name=f"doughnut.readout.{index}.headline",
        )
        deck.text(
            slide,
            readout["detail"],
            panel_x + 0.92,
            y + 0.48,
            2.28,
            0.34,
            role="footnote",
            font_size_pt=8.4,
            color_role="text.muted",
            name=f"doughnut.readout.{index}.detail",
        )
        if index < len(ir.get("readouts", [])):
            deck.line(
                slide,
                panel_x + 0.28,
                y + 0.98,
                panel_x + panel_w - 0.28,
                y + 0.98,
                color_role="line.soft",
                width_pt=0.55,
                name=f"doughnut.readout.{index}.divider",
            )

    deck.text(
        slide,
        "PPTSmith × MeowClaw Lab  ·  Native editable PMO sample",
        0.70,
        7.10,
        5.6,
        0.18,
        role="footnote",
        font_size_pt=7.4,
        color_role="text.muted",
        name="doughnut.brand_text",
    )
    deck.component_records.append(
        {
            "component": "annotated_doughnut",
            "variant": "portfolio_poster",
            "layout_mode": layout["mode"],
            "reason_codes": layout["reason_codes"],
            "leader_count": len(layout["leaders"]),
            "leader_segment_count": sum(len(item["points"]) - 1 for item in layout["leaders"]),
            "legend_fallback_count": len(layout["legend_fallback"]),
            "chart_first_slice_angle_deg": layout["chart_first_slice_angle_deg"],
            "endpoint_segment_alignment_passed": all(
                segment["start_angle"] <= segment["endpoint_angle"] <= segment["end_angle"]
                for segment in layout["segments"]
                if segment["index"] < int(ir.get("max_leaders", len(ir["categories"])))
            ),
        }
    )


def _render_roadmap(deck: NativeComponentDeck, ir: Mapping[str, Any]) -> None:
    slide = deck.add_slide(ir["title"], "Phased roadmap", "roadmap")
    frame = {"x": 0.66, "y": 1.44, "w": 12.0, "h": 5.32}
    layout = layout_roadmap(ir, frame, deck.tokens)
    phase_colors = ["sequence.phase_1", "sequence.phase_2", "sequence.phase_3", "sequence.phase_4", "sequence.phase_5"]
    for index, (phase, box) in enumerate(zip(ir["phases"], layout["phase_boxes"])):
        x = box["x"] - (0.50 if index else 0.0)
        w = box["w"] + (0.50 if index else 0.0)
        deck.shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.PENTAGON if index == 0 else MSO_AUTO_SHAPE_TYPE.CHEVRON,
            x,
            box["chevron_y"],
            w,
            box["chevron_h"],
            fill_role=phase_colors[index % len(phase_colors)],
            line_role=phase_colors[index % len(phase_colors)],
            line_width_pt=0.0,
            name=f"roadmap.phase.{index + 1}.chevron",
        )
        deck.text(
            slide,
            f"{index + 1:02d}  {phase['name']}\n{phase['period']}",
            box["x"] + 0.22,
            box["chevron_y"] + 0.19,
            box["w"] - 0.44,
            0.60,
            role="section_title",
            color_role="text.inverse",
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            name=f"roadmap.phase.{index + 1}.header",
        )
        deck.text(
            slide,
            phase["objective"],
            box["x"] + 0.16,
            box["output_y"],
            box["w"] - 0.32,
            0.52,
            role="body",
            color_role="text.title",
            bold=True,
            align=PP_ALIGN.CENTER,
            name=f"roadmap.phase.{index + 1}.objective",
        )
        output_text = "\n".join(f"• {item}" for item in phase["outputs"])
        deck.text(
            slide,
            output_text,
            box["x"] + 0.20,
            box["output_y"] + 0.70,
            box["w"] - 0.40,
            0.72,
            role="label",
            color_role="text.body",
            name=f"roadmap.phase.{index + 1}.outputs",
        )
        deck.shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            box["x"] + 0.10,
            box["summary_y"] - 0.82,
            box["w"] - 0.20,
            0.62,
            fill_role="surface.paper",
            line_role="line.divider",
            name=f"roadmap.phase.{index + 1}.exit_panel",
        )
        deck.text(
            slide,
            f"EXIT\n{phase['exit_condition']}",
            box["x"] + 0.18,
            box["summary_y"] - 0.74,
            box["w"] - 0.36,
            0.46,
            role="footnote",
            color_role="text.body",
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            name=f"roadmap.phase.{index + 1}.exit",
        )

    gate_y = frame["y"] + 1.34
    deck.line(
        slide,
        frame["x"] + 0.05,
        gate_y,
        frame["x"] + frame["w"] - 0.05,
        gate_y,
        color_role="line.divider",
        width_pt=1.0,
        name="roadmap.gate_channel",
    )
    for index, gate in enumerate(layout["gates"]):
        gate_x = min(gate["x"], frame["x"] + frame["w"] - 0.12)
        deck.shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.DIAMOND,
            gate_x - 0.10,
            gate_y - 0.10,
            0.20,
            0.20,
            fill_role="schedule.milestone",
            line_role="schedule.milestone",
            name=f"roadmap.gate.{index + 1}.diamond",
        )
        deck.text(
            slide,
            gate["label"],
            gate_x - 0.28,
            gate_y + 0.14,
            0.56,
            0.19,
            role="footnote",
            color_role="text.muted",
            bold=True,
            align=PP_ALIGN.CENTER,
            name=f"roadmap.gate.{index + 1}.label",
        )
    deck.shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        frame["x"],
        frame["y"] + frame["h"] - 0.44,
        frame["w"],
        0.36,
        fill_role="sequence.phase_1",
        line_role="sequence.phase_1",
        name="roadmap.duration_bar",
    )
    deck.text(
        slide,
        "TOTAL DURATION  |  JUL-DEC 2026  |  SCALE DECISION AT G4",
        frame["x"] + 0.20,
        frame["y"] + frame["h"] - 0.37,
        frame["w"] - 0.40,
        0.20,
        role="label",
        color_role="text.inverse",
        bold=True,
        align=PP_ALIGN.CENTER,
        name="roadmap.duration_label",
    )
    deck.component_records.append(
        {
            "component": "roadmap",
            "layout_mode": layout["mode"],
            "reason_codes": layout["reason_codes"],
            "phase_count": len(ir["phases"]),
            "gate_alignment_passed": layout["boundaries"] == [gate["x"] for gate in layout["gates"]],
        }
    )


def _month_boundaries(start: date, end: date) -> list[date]:
    current = date(start.year, start.month, 1)
    boundaries = [current]
    while current <= end:
        if current.month == 12:
            current = date(current.year + 1, 1, 1)
        else:
            current = date(current.year, current.month + 1, 1)
        boundaries.append(current)
    return boundaries


def _render_gantt(deck: NativeComponentDeck, ir: Mapping[str, Any]) -> None:
    slide = deck.add_slide(ir["title"], "Gantt timeline", "gantt")
    task_x, task_w = 0.60, 2.22
    owner_x, owner_w = 2.82, 0.78
    timeline = {"x": 3.60, "y": 1.46, "w": 9.08, "h": 4.95}
    layout = layout_gantt(ir, timeline, deck.tokens)
    header_y = timeline["y"]
    body_top = timeline["y"] + 0.72
    body_bottom = timeline["y"] + timeline["h"]
    deck.shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, task_x, header_y, task_w, 0.72, fill_role="sequence.phase_1", line_role="sequence.phase_1", name="gantt.task_header")
    deck.shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, owner_x, header_y, owner_w, 0.72, fill_role="sequence.phase_1", line_role="sequence.phase_1", name="gantt.owner_header")
    deck.text(slide, "TASK / DELIVERABLE", task_x + 0.12, header_y + 0.24, task_w - 0.24, 0.22, role="label", color_role="text.inverse", bold=True, name="gantt.task_header_label")
    deck.text(slide, "OWNER", owner_x + 0.05, header_y + 0.24, owner_w - 0.10, 0.22, role="label", color_role="text.inverse", bold=True, align=PP_ALIGN.CENTER, name="gantt.owner_header_label")
    deck.shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, timeline["x"], header_y, timeline["w"], 0.72, fill_role="surface.paper", line_role="line.divider", name="gantt.timeline_header")

    for index, month in enumerate(layout["month_spans"]):
        deck.text(
            slide,
            month["label"],
            month["x1"],
            header_y + 0.05,
            month["x2"] - month["x1"],
            0.23,
            role="label",
            color_role="text.title",
            bold=True,
            align=PP_ALIGN.CENTER,
            name=f"gantt.month.{index + 1}",
        )

    for column in layout["time_columns"]:
        column_start = date.fromisoformat(column["start"])
        column_end = date.fromisoformat(column["end"])
        deck.line(
            slide,
            column["x1"],
            header_y + 0.31,
            column["x1"],
            body_bottom,
            color_role="line.soft",
            width_pt=0.5,
            name=f"gantt.week_grid.{column['index']}",
        )
        deck.text(
            slide,
            f"W{column['index']}\n{column_start.day:02d}-{column_end.day:02d}",
            column["x1"],
            header_y + 0.37,
            column["x2"] - column["x1"],
            0.30,
            role="footnote",
            color_role="text.muted",
            align=PP_ALIGN.CENTER,
            name=f"gantt.week_label.{column['index']}",
        )
    deck.line(
        slide,
        timeline["x"] + timeline["w"],
        header_y + 0.31,
        timeline["x"] + timeline["w"],
        body_bottom,
        color_role="line.soft",
        width_pt=0.5,
        name="gantt.week_grid.end",
    )

    for index, row in enumerate(layout["rows"]):
        y = row["y"]
        row_fill = "surface.paper" if index % 2 else "surface.canvas"
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, task_x, y - 0.10, task_w, row["h"], fill_role=row_fill, line_role="line.soft", line_width_pt=0.5, name=f"gantt.row.{index + 1}.task_cell")
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, owner_x, y - 0.10, owner_w, row["h"], fill_role=row_fill, line_role="line.soft", line_width_pt=0.5, name=f"gantt.row.{index + 1}.owner_cell")
        deck.text(slide, row["label"], task_x + 0.12, y - 0.02, task_w - 0.24, 0.25, role="label", color_role="text.body", bold=True, name=f"gantt.row.{index + 1}.task")
        deck.text(slide, row["owner"], owner_x + 0.04, y - 0.02, owner_w - 0.08, 0.25, role="footnote", color_role="text.muted", align=PP_ALIGN.CENTER, name=f"gantt.row.{index + 1}.owner")
        deck.line(slide, timeline["x"], y - 0.10 + row["h"], timeline["x"] + timeline["w"], y - 0.10 + row["h"], color_role="line.soft", width_pt=0.5, name=f"gantt.row.{index + 1}.rule")
        for interval_index, bar in enumerate(row["planned_bars"]):
            deck.shape(
                slide,
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                bar["x1"],
                y + 0.04,
                max(0.04, bar["x2"] - bar["x1"]),
                0.18,
                fill_role="schedule.planned",
                line_role="schedule.planned",
                line_width_pt=0.0,
                name=f"gantt.row.{index + 1}.planned.{interval_index + 1}",
            )
        for interval_index, bar in enumerate(row["actual_bars"]):
            deck.shape(
                slide,
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                bar["x1"],
                y + 0.085,
                max(0.04, bar["x2"] - bar["x1"]),
                0.09,
                fill_role="schedule.actual",
                line_role="schedule.actual",
                line_width_pt=0.0,
                name=f"gantt.row.{index + 1}.actual.{interval_index + 1}",
            )
        actual_endpoint = row["actual_bars"][-1]["x2"] if row["actual_bars"] else row["planned_bars"][0]["x1"]
        deck.shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.OVAL,
            actual_endpoint - 0.055,
            y + 0.075,
            0.11,
            0.11,
            fill_role=f"state.{row['health']}",
            line_role=f"state.{row['health']}",
            name=f"gantt.row.{index + 1}.health",
        )
        for milestone_index, milestone in enumerate(row["milestones"]):
            deck.shape(
                slide,
                MSO_AUTO_SHAPE_TYPE.DIAMOND,
                milestone["x"] - 0.09,
                y + 0.01,
                0.18,
                0.18,
                fill_role="schedule.milestone",
                line_role="schedule.milestone",
                name=f"gantt.row.{index + 1}.milestone.{milestone_index + 1}",
            )

    if layout["reference_x"] is not None:
        deck.line(
            slide,
            layout["reference_x"],
            body_top,
            layout["reference_x"],
            body_bottom,
            color_role="schedule.reference",
            width_pt=1.5,
            name="gantt.reference_line",
        )
        deck.text(
            slide,
            "TODAY",
            layout["reference_x"] - 0.30,
            body_top - 0.29,
            0.60,
            0.18,
            role="footnote",
            color_role="schedule.reference",
            bold=True,
            align=PP_ALIGN.CENTER,
            name="gantt.reference_label",
        )

    legend_items = [
        ("planned", "Planned", "schedule.planned", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE),
        ("actual", "Actual", "schedule.actual", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE),
        ("health_healthy", "On track", "state.healthy", MSO_AUTO_SHAPE_TYPE.OVAL),
        ("health_watch", "Watch", "state.watch", MSO_AUTO_SHAPE_TYPE.OVAL),
        ("health_neutral", "Not assessed", "state.neutral", MSO_AUTO_SHAPE_TYPE.OVAL),
        ("milestone", "Milestone", "schedule.milestone", MSO_AUTO_SHAPE_TYPE.DIAMOND),
        ("reference", "Today", "schedule.reference", MSO_AUTO_SHAPE_TYPE.RECTANGLE),
    ]
    legend_y = 6.60
    for index, (_, label, color_role, kind) in enumerate(legend_items):
        x = 3.65 + index * 1.24
        deck.shape(slide, kind, x, legend_y, 0.16, 0.12, fill_role=color_role, line_role=color_role, name=f"gantt.legend.{index + 1}.symbol")
        deck.text(slide, label, x + 0.22, legend_y - 0.02, 0.96, 0.18, role="footnote", color_role="text.muted", name=f"gantt.legend.{index + 1}.label")
    used_legend_roles = {"planned", "actual", "milestone", "reference"}
    used_legend_roles.update(f"health_{row['health']}" for row in layout["rows"])
    legend_qa = legend_completeness(
        used_legend_roles,
        {item[0] for item in legend_items},
    )
    deck.component_records.append(
        {
            "component": "gantt",
            "layout_mode": layout["mode"],
            "reason_codes": layout["reason_codes"],
            "row_count": len(ir["tasks"]),
            "timeline_column_count": len(layout["time_columns"]),
            "week_label_count": len(layout["time_columns"]),
            "month_span_count": len(layout["month_spans"]),
            "month_spans_cover_timeline": (
                layout["month_spans"][0]["x1"] == timeline["x"]
                and layout["month_spans"][-1]["x2"] == timeline["x"] + timeline["w"]
            ),
            "legend_completeness": legend_qa,
        }
    )


def _format_number(value: Any) -> str:
    numeric = float(value)
    if numeric.is_integer():
        return f"{numeric:,.0f}"
    return f"{numeric:,.1f}"


def _render_kpi_status(deck: NativeComponentDeck, ir: Mapping[str, Any]) -> None:
    slide = deck.add_slide(ir["title"], "KPI / status", "kpi")
    frame = {"x": 0.62, "y": 1.42, "w": 12.08, "h": 4.98}
    layout = layout_kpi_status(ir, frame, deck.tokens)
    quality = qa_component(ir, layout, deck.tokens)
    if not quality["passed"]:
        raise ValueError(f"KPI QA failed: {quality['reason_codes']}")
    page = layout["pages"][0]
    for index, card in enumerate(page["cards"]):
        metric = card["metric"]
        prefix = f"kpi.card.{index + 1}"
        deck.shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            card["x"],
            card["y"],
            card["w"],
            card["h"],
            fill_role="surface.panel",
            line_role="line.divider",
            line_width_pt=0.8,
            name=f"{prefix}.panel",
        )
        deck.shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            card["x"],
            card["y"],
            card["w"],
            0.09,
            fill_role=f"state.{metric['status']}",
            line_role=f"state.{metric['status']}",
            line_width_pt=0.0,
            name=f"{prefix}.status_bar",
        )
        deck.shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.OVAL,
            card["x"] + 0.20,
            card["y"] + 0.25,
            0.13,
            0.13,
            fill_role=f"state.{metric['status']}",
            line_role=f"state.{metric['status']}",
            name=f"{prefix}.status_dot",
        )
        deck.text(
            slide,
            f"{metric['id']}  {metric['label']}",
            card["x"] + 0.42,
            card["y"] + 0.19,
            card["w"] - 0.62,
            0.30,
            role="label",
            color_role="text.title",
            bold=True,
            name=f"{prefix}.label",
        )
        deck.text(
            slide,
            _format_number(metric["value"]),
            card["x"] + 0.20,
            card["y"] + 0.58,
            1.42,
            0.52,
            role="metric",
            color_role="text.title",
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
            name=f"{prefix}.value",
        )
        deck.text(
            slide,
            metric["unit"],
            card["x"] + 1.55,
            card["y"] + 0.76,
            0.72,
            0.24,
            role="label",
            color_role="text.muted",
            bold=True,
            name=f"{prefix}.unit",
        )
        variance_sign = "+" if card["variance"] >= 0 else ""
        deck.text(
            slide,
            f"TARGET  {_format_number(metric['target'])} {metric['unit']}\n"
            f"VARIANCE  {variance_sign}{_format_number(card['variance'])}",
            card["x"] + 0.20,
            card["y"] + 1.18,
            card["w"] * 0.48,
            0.48,
            role="footnote",
            color_role="text.body",
            bold=True,
            name=f"{prefix}.target",
        )
        trend_points = card["trend_points"]
        for point_index, (first, second) in enumerate(zip(trend_points, trend_points[1:])):
            deck.line(
                slide,
                first["x"],
                first["y"],
                second["x"],
                second["y"],
                color_role=f"state.{metric['status']}",
                width_pt=1.6,
                name=f"{prefix}.trend.{point_index + 1}",
            )
        if trend_points:
            for point_index, point in enumerate(trend_points):
                deck.shape(
                    slide,
                    MSO_AUTO_SHAPE_TYPE.OVAL,
                    point["x"] - 0.035,
                    point["y"] - 0.035,
                    0.07,
                    0.07,
                    fill_role=f"state.{metric['status']}",
                    line_role=f"state.{metric['status']}",
                    name=f"{prefix}.trend_dot.{point_index + 1}",
                )
        deck.text(
            slide,
            metric["note"],
            card["x"] + 0.20,
            card["y"] + card["h"] - 0.42,
            card["w"] - 0.40,
            0.24,
            role="footnote",
            color_role="text.muted",
            name=f"{prefix}.note",
        )
        deck.text(
            slide,
            metric["period"],
            card["x"] + card["w"] - 1.10,
            card["y"] + card["h"] - 0.42,
            0.90,
            0.24,
            role="footnote",
            color_role="text.muted",
            bold=True,
            align=PP_ALIGN.RIGHT,
            name=f"{prefix}.period",
        )
    statuses = ["healthy", "watch", "risk", "neutral"]
    for index, status in enumerate(statuses):
        x = 7.15 + index * 1.33
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x, 6.62, 0.11, 0.11, fill_role=f"state.{status}", line_role=f"state.{status}", name=f"kpi.legend.{status}.dot")
        deck.text(slide, deck.tokens["status_semantics"][status]["label"], x + 0.18, 6.58, 1.05, 0.20, role="footnote", color_role="text.muted", name=f"kpi.legend.{status}.label")
    deck.component_records.append(
        {
            "component": "kpi_status",
            "layout_mode": layout["mode"],
            "reason_codes": layout["reason_codes"],
            "metric_count": len(ir["metrics"]),
            "qa": quality,
        }
    )


def _render_milestone(deck: NativeComponentDeck, ir: Mapping[str, Any]) -> None:
    slide = deck.add_slide(ir["title"], "Milestone timeline", "milestone")
    frame = {"x": 0.72, "y": 1.50, "w": 11.88, "h": 4.88}
    layout = layout_milestone(ir, frame, deck.tokens)
    quality = qa_component(ir, layout, deck.tokens)
    if not quality["passed"]:
        raise ValueError(f"Milestone QA failed: {quality['reason_codes']}")
    page = layout["pages"][0]
    axis_y = page["axis_y"]
    deck.line(slide, frame["x"], axis_y, frame["x"] + frame["w"], axis_y, color_role="sequence.phase_1", width_pt=2.5, name="milestone.axis")
    deck.text(slide, ir["start_date"], frame["x"], axis_y + 0.13, 1.15, 0.20, role="footnote", color_role="text.muted", name="milestone.axis.start")
    deck.text(slide, ir["end_date"], frame["x"] + frame["w"] - 1.15, axis_y + 0.13, 1.15, 0.20, role="footnote", color_role="text.muted", align=PP_ALIGN.RIGHT, name="milestone.axis.end")
    for index, node in enumerate(page["nodes"]):
        milestone = node["milestone"]
        prefix = f"milestone.node.{index + 1}"
        box = node["label_box"]
        deck.line(slide, node["x"], node["y"], node["x"], node["leader_end_y"], color_role="line.annotation", width_pt=1.0, name=f"{prefix}.leader")
        deck.shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.DIAMOND,
            node["x"] - 0.11,
            node["y"] - 0.11,
            0.22,
            0.22,
            fill_role=f"state.{milestone['status']}",
            line_role="sequence.phase_1",
            line_width_pt=1.0,
            name=f"{prefix}.diamond",
        )
        deck.shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            box["x"],
            box["y"],
            box["w"],
            box["h"],
            fill_role="surface.paper",
            line_role="line.divider",
            line_width_pt=0.65,
            name=f"{prefix}.panel",
        )
        deck.text(
            slide,
            f"{milestone['id']}  {milestone['label']}\n{milestone['date']}  |  {milestone['owner']}\nEVIDENCE  {milestone['evidence']}",
            box["x"] + 0.10,
            box["y"] + 0.09,
            box["w"] - 0.20,
            box["h"] - 0.16,
            role="footnote",
            color_role="text.body",
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            name=f"{prefix}.label",
        )
    statuses = ["healthy", "watch", "risk", "neutral"]
    for index, status in enumerate(statuses):
        x = 7.15 + index * 1.33
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.DIAMOND, x, 6.62, 0.12, 0.12, fill_role=f"state.{status}", line_role=f"state.{status}", name=f"milestone.legend.{status}.diamond")
        deck.text(slide, deck.tokens["status_semantics"][status]["label"], x + 0.18, 6.58, 1.05, 0.20, role="footnote", color_role="text.muted", name=f"milestone.legend.{status}.label")
    deck.component_records.append(
        {
            "component": "milestone",
            "layout_mode": layout["mode"],
            "reason_codes": layout["reason_codes"],
            "milestone_count": len(ir["milestones"]),
            "qa": quality,
        }
    )


def _heat_cell_role(likelihood: int, impact: int) -> str:
    score = likelihood * impact
    if score >= 16:
        return "state.risk"
    if score >= 9:
        return "state.watch"
    if score >= 5:
        return "surface.paper"
    return "surface.canvas"


def _render_risk_heat_map(deck: NativeComponentDeck, ir: Mapping[str, Any]) -> None:
    slide = deck.add_slide(ir["title"], "Risk heat map", "heatmap")
    frame = {"x": 0.62, "y": 1.43, "w": 12.08, "h": 5.10}
    layout = layout_risk_heat_map(ir, frame, deck.tokens)
    quality = qa_component(ir, layout, deck.tokens)
    if not quality["passed"]:
        raise ValueError(f"Risk heat map QA failed: {quality['reason_codes']}")
    page = layout["pages"][0]
    grid = page["grid"]
    for impact in range(1, 6):
        for likelihood in range(1, 6):
            x = grid["x"] + (likelihood - 1) * page["cell_w"]
            y = grid["y"] + (5 - impact) * page["cell_h"]
            deck.shape(
                slide,
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                x,
                y,
                page["cell_w"],
                page["cell_h"],
                fill_role=_heat_cell_role(likelihood, impact),
                line_role="surface.canvas",
                line_width_pt=1.2,
                name=f"heatmap.cell.L{likelihood}.I{impact}",
            )
    for value in range(1, 6):
        deck.text(slide, str(value), grid["x"] + (value - 1) * page["cell_w"], grid["y"] + grid["h"] + 0.05, page["cell_w"], 0.20, role="label", color_role="text.muted", align=PP_ALIGN.CENTER, name=f"heatmap.axis.likelihood.{value}")
        deck.text(slide, str(value), grid["x"] - 0.30, grid["y"] + (5 - value) * page["cell_h"] + page["cell_h"] * 0.35, 0.22, 0.20, role="label", color_role="text.muted", align=PP_ALIGN.RIGHT, name=f"heatmap.axis.impact.{value}")
    deck.text(slide, "LIKELIHOOD  →", grid["x"] + grid["w"] * 0.31, grid["y"] + grid["h"] + 0.29, grid["w"] * 0.38, 0.22, role="label", color_role="text.title", bold=True, align=PP_ALIGN.CENTER, name="heatmap.axis.likelihood.title")
    deck.text(slide, "IMPACT  ↑", frame["x"], grid["y"] - 0.03, 0.80, 0.22, role="label", color_role="text.title", bold=True, name="heatmap.axis.impact.title")
    for index, point in enumerate(page["points"]):
        status = point["risk"]["status"]
        prefix = f"heatmap.point.{index + 1}"
        point_diameter = float(layout["point_diameter"])
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, point["x"] - point_diameter / 2, point["y"] - point_diameter / 2, point_diameter, point_diameter, fill_role=f"state.{status}", line_role="sequence.phase_1", line_width_pt=1.0, name=f"{prefix}.dot")
        deck.text(
            slide,
            point["id"],
            point["x"] - 0.15,
            point["y"] - 0.10,
            0.30,
            0.20,
            role="footnote",
            color_role="text.title" if status == "watch" else "text.inverse",
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            name=f"{prefix}.id",
        )
    panel = page["list_panel"]
    deck.shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, panel["x"], panel["y"], panel["w"], panel["h"], fill_role="surface.panel", line_role="line.divider", name="heatmap.register.panel")
    deck.text(slide, "RISK REGISTER  |  ID / EXPOSURE / OWNER / MITIGATION", panel["x"] + 0.18, panel["y"] + 0.16, panel["w"] - 0.36, 0.24, role="label", color_role="text.title", bold=True, name="heatmap.register.header")
    row_h = (panel["h"] - 0.58) / len(page["risks"])
    for index, risk in enumerate(page["risks"]):
        y = panel["y"] + 0.50 + index * row_h
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, panel["x"] + 0.17, y + 0.08, 0.10, 0.10, fill_role=f"state.{risk['status']}", line_role=f"state.{risk['status']}", name=f"heatmap.register.{index + 1}.status")
        deck.text(slide, risk["id"], panel["x"] + 0.34, y + 0.01, 0.42, 0.22, role="label", color_role="text.title", bold=True, name=f"heatmap.register.{index + 1}.id")
        deck.text(slide, risk["label"], panel["x"] + 0.78, y + 0.01, panel["w"] * 0.32, 0.22, role="label", color_role="text.body", bold=True, name=f"heatmap.register.{index + 1}.label")
        deck.text(slide, f"L{risk['likelihood']} × I{risk['impact']}\n{risk['owner']}", panel["x"] + panel["w"] * 0.48, y - 0.01, panel["w"] * 0.18, 0.40, role="footnote", color_role="text.muted", bold=True, name=f"heatmap.register.{index + 1}.owner")
        deck.text(slide, risk["mitigation"], panel["x"] + panel["w"] * 0.67, y - 0.01, panel["w"] * 0.30, 0.40, role="footnote", color_role="text.body", name=f"heatmap.register.{index + 1}.mitigation")
        if index < len(page["risks"]) - 1:
            deck.line(slide, panel["x"] + 0.16, y + row_h - 0.04, panel["x"] + panel["w"] - 0.16, y + row_h - 0.04, color_role="line.soft", width_pt=0.5, name=f"heatmap.register.{index + 1}.rule")
    deck.component_records.append(
        {
            "component": "risk_heat_map",
            "layout_mode": layout["mode"],
            "reason_codes": layout["reason_codes"],
            "risk_count": len(ir["risks"]),
            "qa": quality,
        }
    )


def _set_table_text(cell: Any, value: str, tokens: Mapping[str, Any], *, bold: bool, color_role: str, size_pt: float) -> None:
    frame = cell.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = value
    run.font.name = str(tokens["typography"]["font"])
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(_role_color(tokens, color_role))


def _render_raci(deck: NativeComponentDeck, ir: Mapping[str, Any]) -> None:
    slide = deck.add_slide(ir["title"], "RACI matrix", "raci")
    frame = {"x": 0.64, "y": 1.48, "w": 12.02, "h": 4.94}
    layout = layout_raci(ir, frame, deck.tokens)
    quality = qa_component(ir, layout, deck.tokens)
    if not quality["passed"]:
        raise ValueError(f"RACI QA failed: {quality['reason_codes']}")
    page = layout["pages"][0]
    table_shape = slide.shapes.add_table(
        len(page["rows"]) + 1,
        len(page["roles"]) + 1,
        Inches(frame["x"]),
        Inches(frame["y"]),
        Inches(frame["w"]),
        Inches(frame["h"]),
    )
    table_shape.name = "raci.native_table"
    table = table_shape.table
    table.columns[0].width = Inches(page["activity_w"])
    for column_index in range(1, len(page["roles"]) + 1):
        table.columns[column_index].width = Inches(page["role_w"])
    table.rows[0].height = Inches(page["row_h"])
    for row_index in range(1, len(page["rows"]) + 1):
        table.rows[row_index].height = Inches(page["row_h"])
    header = table.cell(0, 0)
    header.fill.solid()
    header.fill.fore_color.rgb = _rgb(_role_color(deck.tokens, "sequence.phase_1"))
    _set_table_text(header, "ACTIVITY / DECISION", deck.tokens, bold=True, color_role="text.inverse", size_pt=9)
    for column_index, role in enumerate(page["roles"], start=1):
        cell = table.cell(0, column_index)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(_role_color(deck.tokens, "sequence.phase_1"))
        _set_table_text(cell, role, deck.tokens, bold=True, color_role="text.inverse", size_pt=8.5)
    code_styles = {
        "A": ("sequence.phase_1", "text.inverse"),
        "R": ("sequence.phase_3", "text.inverse"),
        "C": ("surface.paper", "text.title"),
        "I": ("line.soft", "text.body"),
        "": ("surface.canvas", "text.muted"),
    }
    for row_index, row in enumerate(page["rows"], start=1):
        activity_cell = table.cell(row_index, 0)
        activity_cell.fill.solid()
        activity_cell.fill.fore_color.rgb = _rgb(
            _role_color(deck.tokens, "surface.paper" if row_index % 2 else "surface.canvas")
        )
        _set_table_text(activity_cell, row["activity"], deck.tokens, bold=True, color_role="text.body", size_pt=9)
        activity_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        for column_index, code in enumerate(row["codes"], start=1):
            cell = table.cell(row_index, column_index)
            fill_role, text_role = code_styles[code]
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(_role_color(deck.tokens, fill_role))
            _set_table_text(cell, code, deck.tokens, bold=True, color_role=text_role, size_pt=12)
    legend = [
        ("A", "Accountable", "sequence.phase_1", "text.inverse"),
        ("R", "Responsible", "sequence.phase_3", "text.inverse"),
        ("C", "Consulted", "surface.paper", "text.title"),
        ("I", "Informed", "line.soft", "text.body"),
    ]
    for index, (code, label, fill_role, text_role) in enumerate(legend):
        x = 7.25 + index * 1.33
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, 6.61, 0.24, 0.18, fill_role=fill_role, line_role="line.divider", name=f"raci.legend.{code}.box")
        deck.text(slide, code, x, 6.59, 0.24, 0.18, role="footnote", color_role=text_role, bold=True, align=PP_ALIGN.CENTER, name=f"raci.legend.{code}.code")
        deck.text(slide, label, x + 0.31, 6.58, 0.92, 0.20, role="footnote", color_role="text.muted", name=f"raci.legend.{code}.label")
    deck.component_records.append(
        {
            "component": "raci",
            "layout_mode": layout["mode"],
            "reason_codes": layout["reason_codes"],
            "role_count": len(ir["roles"]),
            "activity_count": len(ir["activities"]),
            "qa": quality,
        }
    )


def _count_shapes(prs: Presentation) -> dict[str, int]:
    counts = {
        "slide_count": len(prs.slides),
        "picture_count": 0,
        "chart_count": 0,
        "table_count": 0,
        "text_shape_count": 0,
        "native_shape_count": 0,
        "connector_count": 0,
        "total_shape_count": 0,
    }
    for slide in prs.slides:
        for shape in slide.shapes:
            counts["total_shape_count"] += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                counts["picture_count"] += 1
            else:
                counts["native_shape_count"] += 1
            if getattr(shape, "has_chart", False):
                counts["chart_count"] += 1
            if getattr(shape, "has_table", False):
                counts["table_count"] += 1
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text:
                counts["text_shape_count"] += 1
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                counts["connector_count"] += 1
    return counts


def _check_object_bounds(prs: Presentation, *, tolerance_emu: int = 1000) -> dict[str, Any]:
    outside: list[str] = []
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    for slide in prs.slides:
        for shape in slide.shapes:
            x1 = min(int(shape.left), int(shape.left + shape.width))
            x2 = max(int(shape.left), int(shape.left + shape.width))
            y1 = min(int(shape.top), int(shape.top + shape.height))
            y2 = max(int(shape.top), int(shape.top + shape.height))
            if x1 < -tolerance_emu or y1 < -tolerance_emu or x2 > slide_w + tolerance_emu or y2 > slide_h + tolerance_emu:
                outside.append(shape.name)
    return {"passed": not outside, "outside": outside}


def build_phase1_sample_deck(
    output: Path,
    *,
    tokens: Mapping[str, Any],
    components: Iterable[Mapping[str, Any]] | None = None,
) -> dict[str, Any]:
    """Build one editable sample slide for each Phase 1 component."""

    selected = [validate_component_ir(item) for item in (components or DEFAULT_PHASE1_IRS)]
    required_types = ["annotated_doughnut", "roadmap", "gantt"]
    actual_types = [item["type"] for item in selected]
    if actual_types != required_types:
        raise ValueError(f"phase 1 sample requires component order {required_types}")
    deck = NativeComponentDeck(tokens)
    _render_doughnut(deck, selected[0])
    _render_roadmap(deck, selected[1])
    _render_gantt(deck, selected[2])
    deck.save(output)
    counts = _count_shapes(deck.prs)
    slide_region = {
        "x": 0.0,
        "y": 0.0,
        "w": float(tokens["slide"]["width_in"]),
        "h": float(tokens["slide"]["height_in"]),
    }
    bounds_qa = check_text_bounds(deck.text_boxes, slide_region)
    overflow_qa = check_text_overflow(deck.text_boxes, min_font_pt=7.5)
    return {
        **counts,
        "output": str(output),
        "components": deck.component_records,
        "text_bounds": bounds_qa,
        "text_overflow": overflow_qa,
        "native_editability": {
            "message_objects_rasterized": False,
            "whole_slide_pictures": 0,
            "connector_endpoint_binding_verified": False,
        },
    }


def build_portfolio_doughnut_deck(
    output: Path,
    *,
    tokens: Mapping[str, Any],
    component: Mapping[str, Any],
) -> dict[str, Any]:
    """Build one native editable portfolio-investment poster slide."""

    selected = validate_component_ir(component)
    if selected["type"] != "annotated_doughnut":
        raise ValueError("portfolio poster requires an annotated_doughnut component")
    deck = NativeComponentDeck(tokens)
    _render_portfolio_doughnut(deck, selected)
    deck.save(output)
    counts = _count_shapes(deck.prs)
    slide_region = {
        "x": 0.0,
        "y": 0.0,
        "w": float(tokens["slide"]["width_in"]),
        "h": float(tokens["slide"]["height_in"]),
    }
    return {
        **counts,
        "output": str(output),
        "components": deck.component_records,
        "object_bounds": _check_object_bounds(deck.prs),
        "text_bounds": check_text_bounds(deck.text_boxes, slide_region),
        "text_overflow": check_text_overflow(deck.text_boxes, min_font_pt=7.0),
        "native_editability": {
            "message_objects_rasterized": False,
            "whole_slide_pictures": 0,
            "package_media_files": 0,
            "connector_endpoint_binding_verified": False,
        },
    }


def build_phase2_p0_sample_deck(
    output: Path,
    *,
    tokens: Mapping[str, Any],
    components: Iterable[Mapping[str, Any]] | None = None,
) -> dict[str, Any]:
    """Build one native editable sample slide for each Phase 2 P0 component."""

    selected = [validate_component_ir(item) for item in (components or DEFAULT_PHASE2_P0_IRS)]
    required_types = ["kpi_status", "milestone", "risk_heat_map", "raci"]
    actual_types = [item["type"] for item in selected]
    if actual_types != required_types:
        raise ValueError(f"phase 2 P0 sample requires component order {required_types}")
    deck = NativeComponentDeck(tokens)
    _render_kpi_status(deck, selected[0])
    _render_milestone(deck, selected[1])
    _render_risk_heat_map(deck, selected[2])
    _render_raci(deck, selected[3])
    deck.save(output)
    counts = _count_shapes(deck.prs)
    slide_region = {
        "x": 0.0,
        "y": 0.0,
        "w": float(tokens["slide"]["width_in"]),
        "h": float(tokens["slide"]["height_in"]),
    }
    bounds_qa = check_text_bounds(deck.text_boxes, slide_region)
    overflow_qa = check_text_overflow(deck.text_boxes, min_font_pt=7.5)
    component_failures = [
        {
            "component": record["component"],
            "reason_codes": record["qa"]["reason_codes"],
        }
        for record in deck.component_records
        if not record["qa"]["passed"]
    ]
    return {
        **counts,
        "output": str(output),
        "components": deck.component_records,
        "text_bounds": bounds_qa,
        "text_overflow": overflow_qa,
        "component_qa": {
            "passed": not component_failures,
            "failures": component_failures,
        },
        "native_editability": {
            "message_objects_rasterized": False,
            "whole_slide_pictures": 0,
            "package_media_files": 0,
        },
    }
