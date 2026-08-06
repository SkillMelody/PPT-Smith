from __future__ import annotations

import json
import math
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from components.inspection import inspect_phase2_p0_pptx
from components.ir import ComponentValidationError, validate_component_ir
from components.layout import (
    layout_kpi_status,
    layout_milestone,
    layout_raci,
    layout_risk_heat_map,
)
from components.qa import qa_component
from components.renderer import build_phase2_p0_sample_deck
from components.tokens import resolve_tokens


def load_fixture(name: str) -> dict:
    return json.loads((ROOT / "tests" / "fixtures" / name).read_text(encoding="utf-8"))


def normal_components() -> list[dict]:
    return [
        {
            "type": "kpi_status",
            "title": "Transformation value is on plan while cycle time remains at risk",
            "metrics": [
                {
                    "id": "K1",
                    "label": "Run-rate value",
                    "value": 8.4,
                    "unit": "$m",
                    "target": 8.0,
                    "direction": "higher_is_better",
                    "status": "healthy",
                    "period": "Q3 2026",
                    "note": "Validated by Finance",
                    "trend_values": [6.1, 6.8, 7.4, 8.4],
                },
                {
                    "id": "K2",
                    "label": "Adoption",
                    "value": 74,
                    "unit": "%",
                    "target": 80,
                    "direction": "higher_is_better",
                    "status": "watch",
                    "period": "Aug 2026",
                    "note": "Six points below target",
                    "trend_values": [51, 60, 68, 74],
                },
                {
                    "id": "K3",
                    "label": "Cycle time",
                    "value": 12.3,
                    "unit": "days",
                    "target": 10,
                    "direction": "lower_is_better",
                    "status": "risk",
                    "period": "Aug 2026",
                    "note": "Data handoffs are the constraint",
                    "trend_values": [15.2, 14.1, 13.0, 12.3],
                },
                {
                    "id": "K4",
                    "label": "Control pass rate",
                    "value": 96,
                    "unit": "%",
                    "target": 95,
                    "direction": "higher_is_better",
                    "status": "healthy",
                    "period": "Aug 2026",
                    "note": "Sample size 412",
                },
            ],
        },
        {
            "type": "milestone",
            "title": "Six evidence-backed milestones define the path to scale approval",
            "start_date": "2026-07-01",
            "end_date": "2026-12-15",
            "milestones": [
                {"id": "M1", "label": "Charter approved", "date": "2026-07-08", "owner": "PMO", "evidence": "Signed charter", "status": "healthy", "phase": "Discover"},
                {"id": "M2", "label": "Blueprint signed", "date": "2026-08-05", "owner": "Architecture", "evidence": "Design authority minute", "status": "healthy", "phase": "Design"},
                {"id": "M3", "label": "Data rehearsal", "date": "2026-09-02", "owner": "Data", "evidence": "Reconciliation pack", "status": "watch", "phase": "Pilot"},
                {"id": "M4", "label": "Pilot release", "date": "2026-10-01", "owner": "Product", "evidence": "Release record", "status": "healthy", "phase": "Pilot"},
                {"id": "M5", "label": "Adoption threshold", "date": "2026-11-05", "owner": "Change", "evidence": "Usage dashboard", "status": "watch", "phase": "Scale"},
                {"id": "M6", "label": "Scale decision", "date": "2026-12-10", "owner": "SteerCo", "evidence": "Decision log", "status": "neutral", "phase": "Scale"},
            ],
        },
        {
            "type": "risk_heat_map",
            "title": "Data readiness concentrates the critical exposure in two cells",
            "risks": [
                {"id": "R1", "label": "Source data quality", "likelihood": 5, "impact": 5, "owner": "Data", "status": "risk", "mitigation": "Daily cleanse burn-down"},
                {"id": "R2", "label": "Interface latency", "likelihood": 4, "impact": 5, "owner": "Engineering", "status": "risk", "mitigation": "Performance war room"},
                {"id": "R3", "label": "SME availability", "likelihood": 4, "impact": 4, "owner": "Operations", "status": "watch", "mitigation": "Named alternates"},
                {"id": "R4", "label": "Training completion", "likelihood": 4, "impact": 4, "owner": "Change", "status": "watch", "mitigation": "Manager escalation"},
                {"id": "R5", "label": "Vendor access", "likelihood": 4, "impact": 4, "owner": "Security", "status": "watch", "mitigation": "Pre-approved access window"},
                {"id": "R6", "label": "Benefits attribution", "likelihood": 2, "impact": 3, "owner": "Finance", "status": "neutral", "mitigation": "Reconcile baseline monthly"},
            ],
        },
        {
            "type": "raci",
            "title": "Single-point accountability is explicit across the pilot governance chain",
            "roles": ["Sponsor", "PMO", "Product", "Data", "Change"],
            "activities": ["Approve charter", "Sign blueprint", "Release pilot", "Validate benefits", "Approve scale"],
            "assignments": [
                {"activity": "Approve charter", "role": "Sponsor", "code": "A"},
                {"activity": "Approve charter", "role": "PMO", "code": "R"},
                {"activity": "Approve charter", "role": "Product", "code": "C"},
                {"activity": "Sign blueprint", "role": "Sponsor", "code": "A"},
                {"activity": "Sign blueprint", "role": "Product", "code": "R"},
                {"activity": "Sign blueprint", "role": "Data", "code": "C"},
                {"activity": "Release pilot", "role": "Sponsor", "code": "A"},
                {"activity": "Release pilot", "role": "Product", "code": "R"},
                {"activity": "Release pilot", "role": "Change", "code": "C"},
                {"activity": "Validate benefits", "role": "Sponsor", "code": "A"},
                {"activity": "Validate benefits", "role": "PMO", "code": "R"},
                {"activity": "Validate benefits", "role": "Data", "code": "C"},
                {"activity": "Approve scale", "role": "Sponsor", "code": "A"},
                {"activity": "Approve scale", "role": "PMO", "code": "R"},
                {"activity": "Approve scale", "role": "Product", "code": "C"},
                {"activity": "Approve scale", "role": "Change", "code": "I"},
            ],
        },
    ]


class Phase2P0ComponentTests(unittest.TestCase):
    def setUp(self) -> None:
        legacy = json.loads((ROOT / "contracts" / "mckinsey-pmo.style.json").read_text(encoding="utf-8"))
        self.tokens = resolve_tokens(legacy)
        self.components = normal_components()

    def test_all_four_semantic_irs_validate_without_geometry(self) -> None:
        for component in self.components:
            validated = validate_component_ir(component)
            self.assertEqual(component["type"], validated["type"])
            serialized = json.dumps(validated)
            self.assertNotIn('"x"', serialized)
            self.assertNotIn('"shape"', serialized)

    def test_semantic_blocking_gates_reject_invalid_contracts(self) -> None:
        invalid_kpi = dict(self.components[0])
        invalid_kpi["metrics"] = [dict(self.components[0]["metrics"][0], direction="up", value="eight")]
        with self.assertRaises(ComponentValidationError) as kpi_error:
            validate_component_ir(invalid_kpi)
        self.assertTrue({"metric_direction_invalid", "metric_value_invalid"}.issubset(kpi_error.exception.reason_codes))

        invalid_milestone = dict(self.components[1])
        invalid_milestone["milestones"] = [
            dict(self.components[1]["milestones"][0], id="M1", evidence=""),
            dict(self.components[1]["milestones"][1], id="M1", date="2027-01-01"),
        ]
        with self.assertRaises(ComponentValidationError) as milestone_error:
            validate_component_ir(invalid_milestone)
        self.assertTrue(
            {"milestone_id_duplicate", "milestone_evidence_missing", "milestone_outside_timeline"}.issubset(
                milestone_error.exception.reason_codes
            )
        )

        invalid_risk = dict(self.components[2])
        invalid_risk["risks"] = [dict(self.components[2]["risks"][0], likelihood=6, mitigation="")]
        with self.assertRaises(ComponentValidationError) as risk_error:
            validate_component_ir(invalid_risk)
        self.assertTrue({"risk_scale_invalid", "risk_mitigation_missing"}.issubset(risk_error.exception.reason_codes))

        invalid_raci = dict(self.components[3])
        invalid_raci["assignments"] = [
            {"activity": "Approve charter", "role": "Sponsor", "code": "A"},
            {"activity": "Approve charter", "role": "PMO", "code": "A"},
            {"activity": "Sign blueprint", "role": "Product", "code": "X"},
        ]
        with self.assertRaises(ComponentValidationError) as raci_error:
            validate_component_ir(invalid_raci)
        self.assertTrue(
            {"raci_accountable_count_invalid", "raci_responsible_missing", "raci_code_invalid"}.issubset(
                raci_error.exception.reason_codes
            )
        )

    def test_kpi_layout_splits_above_six_and_trends_remain_in_card(self) -> None:
        layout = layout_kpi_status(
            self.components[0],
            {"x": 0.6, "y": 1.4, "w": 12.1, "h": 5.2},
            self.tokens,
        )
        self.assertEqual("single", layout["mode"])
        self.assertEqual(4, len(layout["pages"][0]["cards"]))
        self.assertEqual(layout["pages"][0]["cards"][0]["y"], layout["pages"][0]["cards"][1]["y"])
        self.assertGreater(layout["pages"][0]["cards"][2]["y"], layout["pages"][0]["cards"][0]["y"])
        for card in layout["pages"][0]["cards"]:
            for point in card["trend_points"]:
                self.assertGreaterEqual(point["x"], card["trend_box"]["x"])
                self.assertLessEqual(point["x"], card["trend_box"]["x"] + card["trend_box"]["w"])
                self.assertGreaterEqual(point["y"], card["trend_box"]["y"])
                self.assertLessEqual(point["y"], card["trend_box"]["y"] + card["trend_box"]["h"])

        edge = load_fixture("kpi-status-overflow.json")
        edge_layout = layout_kpi_status(edge, {"x": 0.6, "y": 1.4, "w": 12.1, "h": 5.2}, self.tokens)
        self.assertEqual("split", edge_layout["mode"])
        self.assertEqual([6, 1], [len(page["cards"]) for page in edge_layout["pages"]])

    def test_milestone_layout_maps_dates_and_splits_dense_content(self) -> None:
        frame = {"x": 0.7, "y": 1.5, "w": 11.9, "h": 4.9}
        layout = layout_milestone(self.components[1], frame, self.tokens)
        self.assertEqual("single", layout["mode"])
        nodes = layout["pages"][0]["nodes"]
        self.assertEqual(6, len(nodes))
        self.assertTrue(all(frame["x"] <= node["x"] <= frame["x"] + frame["w"] for node in nodes))
        self.assertEqual(sorted(node["x"] for node in nodes), [node["x"] for node in nodes])
        labels = [node["label_box"] for node in nodes]
        for first, second in zip(labels, labels[1:]):
            if first["lane"] == second["lane"]:
                self.assertLessEqual(first["x"] + first["w"], second["x"] + 1e-6)

        edge_layout = layout_milestone(
            load_fixture("milestone-overdense.json"),
            frame,
            self.tokens,
        )
        self.assertEqual("split", edge_layout["mode"])
        self.assertTrue(all(len(page["nodes"]) <= 8 for page in edge_layout["pages"]))

    def test_heat_map_jitters_same_cell_points_and_uses_table_fallback(self) -> None:
        frame = {"x": 0.6, "y": 1.4, "w": 12.1, "h": 5.25}
        layout = layout_risk_heat_map(self.components[2], frame, self.tokens)
        self.assertEqual("single", layout["mode"])
        same_cell = [point for point in layout["pages"][0]["points"] if point["likelihood"] == 4 and point["impact"] == 4]
        self.assertEqual(3, len(same_cell))
        self.assertEqual(3, len({(round(point["x"], 4), round(point["y"], 4)) for point in same_cell}))
        for index, first in enumerate(same_cell):
            for second in same_cell[index + 1 :]:
                self.assertGreaterEqual(
                    math.hypot(first["x"] - second["x"], first["y"] - second["y"]),
                    layout["point_diameter"],
                )
        self.assertTrue(all(point["cell_bounds"]["x"] < point["x"] < point["cell_bounds"]["x2"] for point in same_cell))
        self.assertTrue(all(point["cell_bounds"]["y"] < point["y"] < point["cell_bounds"]["y2"] for point in same_cell))

        dense = layout_risk_heat_map(load_fixture("risk-heat-map-dense.json"), frame, self.tokens)
        self.assertEqual("table", dense["mode"])
        self.assertIn("risk_density_exceeded", dense["reason_codes"])

    def test_raci_layout_splits_role_and_activity_overflow(self) -> None:
        frame = {"x": 0.6, "y": 1.4, "w": 12.1, "h": 5.3}
        layout = layout_raci(self.components[3], frame, self.tokens)
        self.assertEqual("single", layout["mode"])
        self.assertEqual(5, len(layout["pages"][0]["roles"]))
        self.assertEqual(5, len(layout["pages"][0]["rows"]))
        self.assertTrue(qa_component(self.components[3], layout, self.tokens)["passed"])

        edge_layout = layout_raci(load_fixture("raci-overflow.json"), frame, self.tokens)
        self.assertEqual("split", edge_layout["mode"])
        self.assertTrue(all(len(page["roles"]) <= 8 for page in edge_layout["pages"]))
        self.assertTrue(all(len(page["rows"]) <= 10 for page in edge_layout["pages"]))

    def test_geometry_qa_blocks_out_of_bounds_heat_point(self) -> None:
        frame = {"x": 0.6, "y": 1.4, "w": 12.1, "h": 5.25}
        layout = layout_risk_heat_map(self.components[2], frame, self.tokens)
        self.assertTrue(qa_component(self.components[2], layout, self.tokens)["passed"])
        layout["pages"][0]["points"][0]["x"] = layout["pages"][0]["grid"]["x"] - 0.2
        result = qa_component(self.components[2], layout, self.tokens)
        self.assertFalse(result["passed"])
        self.assertIn("risk_point_outside_grid", result["reason_codes"])

    def test_renderer_and_inspection_prove_four_native_editable_slides(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "phase2.pptx"
            manifest = build_phase2_p0_sample_deck(output, tokens=self.tokens, components=self.components)
            report = inspect_phase2_p0_pptx(output)
            prs = Presentation(output)
        self.assertEqual(4, manifest["slide_count"])
        self.assertEqual(4, len(prs.slides))
        self.assertEqual(0, manifest["picture_count"])
        self.assertGreaterEqual(manifest["table_count"], 1)
        self.assertGreater(manifest["native_shape_count"], 90)
        self.assertTrue(manifest["text_bounds"]["passed"])
        self.assertTrue(manifest["text_overflow"]["passed"])
        self.assertTrue(manifest["component_qa"]["passed"])
        self.assertTrue(report["passed"])
        self.assertEqual(0, report["media_file_count"])
        self.assertEqual([], report["missing_component_prefixes"])
        raci_shapes = {shape.name: shape for shape in prs.slides[3].shapes}
        self.assertGreaterEqual(raci_shapes["raci.title"].height, Inches(0.75))
        self.assertLessEqual(
            raci_shapes["raci.title"].top + raci_shapes["raci.title"].height,
            raci_shapes["raci.title_rule"].top,
        )


if __name__ == "__main__":
    unittest.main()
