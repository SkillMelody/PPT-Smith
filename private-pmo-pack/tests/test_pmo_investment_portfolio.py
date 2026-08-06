from __future__ import annotations

import json
import math
import sys
import tempfile
import unittest
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from components.ir import validate_component_ir
from components.layout import angle_in_interval, layout_doughnut, leaders_have_no_crossings
from components.renderer import build_portfolio_doughnut_deck
from components.tokens import resolve_tokens


class PMOInvestmentPortfolioTests(unittest.TestCase):
    def setUp(self) -> None:
        payload = json.loads(
            (ROOT / "examples" / "pmo-investment-portfolio.poster.json").read_text(encoding="utf-8")
        )
        self.ir = payload["components"][0]
        legacy = json.loads((ROOT / "contracts" / "mckinsey-pmo.style.json").read_text(encoding="utf-8"))
        self.tokens = resolve_tokens(
            legacy,
            overrides={"density": {"doughnut": {"min_label_gap_in": 0.88}}},
        )

    def test_five_categories_percentages_and_amounts_reconcile(self) -> None:
        validated = validate_component_ir(self.ir)
        self.assertEqual(5, len(validated["categories"]))
        self.assertAlmostEqual(100.0, sum(item["value"] for item in validated["categories"]), places=6)
        self.assertAlmostEqual(
            validated["budget_total"],
            sum(item["amount"] for item in validated["categories"]),
            delta=validated["amount_tolerance"],
        )
        for item in validated["categories"]:
            expected = validated["budget_total"] * item["value"] / 100.0
            self.assertAlmostEqual(item["amount"], expected, delta=validated["amount_tolerance"])

    def test_four_data_driven_leaders_are_continuous_owned_and_non_crossing(self) -> None:
        layout = layout_doughnut(self.ir, {"x": 0.68, "y": 1.46, "w": 8.15, "h": 4.88}, self.tokens)
        self.assertEqual("annotated", layout["mode"])
        self.assertEqual(4, len(layout["leaders"]))
        self.assertTrue(leaders_have_no_crossings(layout["leaders"]))
        for leader in layout["leaders"]:
            self.assertEqual(4, len(leader["points"]))
            for first, second in zip(leader["points"], leader["points"][1:]):
                self.assertFalse(
                    math.isclose(first["x"], second["x"], abs_tol=1e-9)
                    and math.isclose(first["y"], second["y"], abs_tol=1e-9)
                )
        callout_segments = [segment for segment in layout["segments"] if segment["index"] < 4]
        for segment in callout_segments:
            endpoint = segment["anchor"]
            endpoint_angle = math.atan2(
                endpoint["y"] - layout["center"]["y"],
                endpoint["x"] - layout["center"]["x"],
            )
            self.assertTrue(
                angle_in_interval(endpoint_angle, segment["start_angle"], segment["end_angle"]),
                msg=segment["label"],
            )
            self.assertAlmostEqual(segment["mid_angle"], segment["endpoint_angle"], places=6)

    def test_fifth_category_uses_compact_legend_fallback(self) -> None:
        layout = layout_doughnut(self.ir, {"x": 0.68, "y": 1.46, "w": 8.15, "h": 4.88}, self.tokens)
        self.assertEqual(1, len(layout["legend_fallback"]))
        self.assertEqual("Risk reserve", layout["legend_fallback"][0]["label"])
        self.assertNotIn(
            "Risk reserve",
            {leader["label"] for leader in layout["leaders"]},
        )

    def test_native_render_has_chart_connectors_no_pictures_media_or_text_fit_failures(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "portfolio.pptx"
            build = build_portfolio_doughnut_deck(output, tokens=self.tokens, component=self.ir)
            self.assertTrue(output.exists())
            self.assertEqual(1, build["slide_count"])
            self.assertEqual(1, build["chart_count"])
            self.assertEqual(0, build["picture_count"])
            self.assertGreaterEqual(build["connector_count"], 12)
            self.assertTrue(build["object_bounds"]["passed"], build["object_bounds"])
            self.assertTrue(build["text_bounds"]["passed"], build["text_bounds"])
            self.assertTrue(build["text_overflow"]["passed"], build["text_overflow"])
            self.assertEqual(4, build["components"][0]["leader_count"])
            self.assertEqual(1, build["components"][0]["legend_fallback_count"])

            prs = Presentation(output)
            names = {shape.name for shape in prs.slides[0].shapes}
            self.assertIn("doughnut.native_chart", names)
            self.assertIn("doughnut.legend_fallback.band", names)
            self.assertTrue(all(f"doughnut.anchor.{index}" in names for index in range(1, 5)))
            self.assertEqual(12, len([name for name in names if name.startswith("doughnut.leader.")]))
            with ZipFile(output) as archive:
                media = [
                    name
                    for name in archive.namelist()
                    if name.startswith("ppt/media/") and not name.endswith("/")
                ]
            self.assertEqual([], media)


if __name__ == "__main__":
    unittest.main()
