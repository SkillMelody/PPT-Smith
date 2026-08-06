from __future__ import annotations

import json
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from components.qa import check_text_bounds, check_text_overflow, legend_completeness
from components.renderer import build_phase1_sample_deck
from components.tokens import resolve_tokens


class ComponentQARendererTests(unittest.TestCase):
    def test_legend_and_text_bounds_checks_report_real_failures(self) -> None:
        legend = legend_completeness(
            {"planned", "actual", "health", "milestone", "reference"},
            {"planned", "actual", "milestone", "reference"},
        )
        self.assertFalse(legend["passed"])
        self.assertEqual(["health"], legend["missing"])

        bounds = check_text_bounds(
            [{"name": "outside", "x": 3.8, "y": 1.0, "w": 0.5, "h": 0.4}],
            {"x": 0.0, "y": 0.0, "w": 4.0, "h": 3.0},
        )
        self.assertFalse(bounds["passed"])
        self.assertEqual(["outside"], bounds["outside"])

    def test_text_overflow_uses_density_aware_capacity(self) -> None:
        result = check_text_overflow(
            [{"name": "phase-title", "text": "A very long phase title that cannot fit", "w": 0.9, "h": 0.2, "font_pt": 12}],
            min_font_pt=8.5,
        )
        self.assertFalse(result["passed"])
        self.assertEqual(["phase-title"], result["overflow"])

    def test_renderer_creates_editable_three_slide_native_pptx(self) -> None:
        legacy = json.loads((ROOT / "contracts" / "mckinsey-pmo.style.json").read_text(encoding="utf-8"))
        tokens = resolve_tokens(legacy)
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "phase1.pptx"
            manifest = build_phase1_sample_deck(output, tokens=tokens)
            self.assertTrue(output.exists())
            self.assertEqual(3, manifest["slide_count"])
            self.assertEqual(0, manifest["picture_count"])
            self.assertGreaterEqual(manifest["chart_count"], 1)
            self.assertGreater(manifest["text_shape_count"], 20)
            self.assertGreater(manifest["native_shape_count"], 40)
            self.assertTrue(manifest["text_bounds"]["passed"])
            self.assertTrue(manifest["text_overflow"]["passed"])

            prs = Presentation(output)
            self.assertEqual(3, len(prs.slides))
            self.assertFalse(any(shape.shape_type == 13 for slide in prs.slides for shape in slide.shapes))
            doughnut_shapes = {shape.name: shape for shape in prs.slides[0].shapes}
            metric = doughnut_shapes["doughnut.center_metric"]
            caption = doughnut_shapes["doughnut.center_caption"]
            self.assertGreaterEqual(metric.width, Inches(1.4))
            self.assertGreaterEqual(caption.top - (metric.top + metric.height), Inches(0.08))

            roadmap_shapes = {shape.name: shape for shape in prs.slides[1].shapes}
            chevrons = [roadmap_shapes[f"roadmap.phase.{index}.chevron"] for index in range(1, 5)]
            self.assertEqual(MSO_AUTO_SHAPE_TYPE.PENTAGON, chevrons[0].auto_shape_type)
            self.assertEqual(MSO_AUTO_SHAPE_TYPE.CHEVRON, chevrons[1].auto_shape_type)
            for first, second in zip(chevrons, chevrons[1:]):
                overlap = first.left + first.width - second.left
                self.assertGreaterEqual(overlap, Inches(0.45))

            gantt_shapes = {shape.name: shape for shape in prs.slides[2].shapes}
            week_labels = [shape for name, shape in gantt_shapes.items() if name.startswith("gantt.week_label.")]
            self.assertEqual(manifest["components"][2]["timeline_column_count"], len(week_labels))
            timeline_left = Inches(3.60)
            timeline_right = Inches(3.60 + 9.08)
            for label in week_labels:
                self.assertGreaterEqual(label.left, timeline_left)
                self.assertLessEqual(label.left + label.width, timeline_right)
                self.assertTrue(label.text.startswith("W"))
            self.assertIn("23-30", gantt_shapes["gantt.week_label.13"].text)
            milestones = [
                shape
                for name, shape in gantt_shapes.items()
                if ".milestone." in name and name.startswith("gantt.row.")
            ]
            for milestone in milestones:
                self.assertGreaterEqual(milestone.left, timeline_left)
                self.assertLessEqual(milestone.left + milestone.width, timeline_right)
            legend_labels = {
                shape.text
                for name, shape in gantt_shapes.items()
                if name.startswith("gantt.legend.") and name.endswith(".label")
            }
            self.assertTrue({"On track", "Watch", "Not assessed"}.issubset(legend_labels))


if __name__ == "__main__":
    unittest.main()
