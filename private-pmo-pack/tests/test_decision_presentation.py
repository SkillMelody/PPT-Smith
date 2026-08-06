from __future__ import annotations

import tempfile
import unittest
import json
import re
from dataclasses import replace
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from components.tokens import resolve_tokens
from decision_support.model import AnalysisMode, DecisionType, RouteAction
from decision_support.pipeline import analyze_decision
from decision_support.presentation import build_decision_presentation_ir, render_decision_deck
from decision_support.serialization import load_decision_case

from tests.test_decision_pipeline import _case


def _shape_map(prs: Presentation) -> dict[str, object]:
    return {shape.name: shape for slide in prs.slides for shape in slide.shapes}


class DecisionPresentationTests(unittest.TestCase):
    def test_visual_contract_leads_with_judgment_conclusion_proof_and_acceptance(self) -> None:
        root = Path(__file__).resolve().parents[1]
        contract = (root / "docs" / "decision-support-visual-encoding-contract.md").read_text()
        self.assertIn("commercial production pack", contract.lower())
        page_sections = re.split(r"(?m)^### \d+\. ", contract)[1:]
        self.assertEqual(11, len(page_sections))
        for section in page_sections:
            self.assertIn("**Audience judgment:**", section)
            self.assertIn("**Core conclusion:**", section)
            self.assertIn("**Primary encoding:**", section)
            self.assertIn("**Acceptance criteria:**", section)

    def test_gold_case_runs_through_real_pipeline_and_builds_11_native_pages(self) -> None:
        root = Path(__file__).resolve().parents[1]
        case = load_decision_case(root / "examples" / "pmo-decision-support-gold.json")
        package = analyze_decision(case)
        self.assertEqual(RouteAction.ROUTE, package.route.route_action)
        self.assertEqual(4, len(package.analysis_results))
        with tempfile.TemporaryDirectory() as temp_dir:
            result = render_decision_deck(package, Path(temp_dir) / "gold.pptx", resolve_tokens())
        self.assertEqual(11, result.slide_count)
        self.assertEqual(0, result.picture_count)
        self.assertEqual(0, result.media_count)
        self.assertLessEqual(result.table_count, 2)
        self.assertGreaterEqual(result.connector_count, 45)
        self.assertIn("critical_path_gantt", result.page_types)
        self.assertIn("risk_exposure", result.page_types)

    def test_ir_uses_package_metadata_and_trace_values(self) -> None:
        package = analyze_decision(_case(DecisionType.VALUE, "value_gate", "value-continue.json"))
        pages = build_decision_presentation_ir(package)
        self.assertEqual(5, len(pages))
        expected = {
            "decision_brief", "benefits_dashboard", "business_case_gate",
            "recommendation_trace", "decision_record",
        }
        self.assertEqual(expected, {page["page_type"] for page in pages})
        trace = {item.id: str(item.value) for item in package.analysis_results[0].calculations}
        benefits = next(page for page in pages if page["page_type"] == "benefits_dashboard")
        self.assertEqual(trace["benefit_realization"], benefits["metrics"][0]["value"])
        self.assertEqual("124", benefits["benefit_series"]["actual_to_date"])
        self.assertEqual("120", benefits["benefit_series"]["planned_to_date"])
        for page in pages:
            metadata = page["metadata"]
            self.assertEqual(package.input_hash[:12], metadata["input_hash_short"])
            self.assertTrue(metadata["analysis_run_id"])
            self.assertTrue(metadata["evidence_cutoff"])
            self.assertTrue(metadata["framework_version"])
            self.assertEqual(package.recommendation.status.value, metadata["recommendation_status"])

    def test_readiness_gate_emits_diagnostic_not_recommendation(self) -> None:
        case = _case(DecisionType.VALUE, "value_gate", "value-insufficient.json")
        package = analyze_decision(replace(
            case,
            question=replace(case.question, requested_analysis_mode=AnalysisMode.DIAGNOSTIC),
        ))
        self.assertEqual(RouteAction.DESCRIBE_ONLY, package.route.route_action)
        pages = build_decision_presentation_ir(package)
        self.assertEqual(("decision_brief", "evidence_gap_diagnostic"), tuple(page["page_type"] for page in pages))
        self.assertNotIn("recommendation_trace", {page["page_type"] for page in pages})

    def test_multi_domain_pages_remain_framework_specific(self) -> None:
        value = _case(DecisionType.VALUE, "value_gate", "value-continue.json")
        resource = _case(DecisionType.RESOURCE, "resource_allocation", "resource-portfolio.json")
        schedule = _case(DecisionType.SCHEDULE, "schedule", "schedule-cpm.json", mode=AnalysisMode.DIAGNOSTIC)
        risk = _case(DecisionType.RISK, "risk_treatment", "risk-mitigation.json")
        package = analyze_decision(replace(
            value,
            id="presentation-multi",
            question=replace(value.question, decision_type=DecisionType.MULTI_DOMAIN),
            extensions={
                "value_gate": value.extensions["value_gate"],
                "resource_allocation": resource.extensions["resource_allocation"],
                "schedule": schedule.extensions["schedule"],
                "risk_treatment": risk.extensions["risk_treatment"],
            },
        ))
        pages = build_decision_presentation_ir(package)
        self.assertEqual(11, len(pages))
        self.assertNotIn("opaque_total_score", repr(pages))
        self.assertEqual(
            {"value_gate", "resource_allocation", "schedule_analysis", "risk_treatment"},
            {page["framework_id"] for page in pages if page.get("framework_id")},
        )

    def test_renderer_produces_native_editable_bounded_deck(self) -> None:
        package = analyze_decision(_case(DecisionType.VALUE, "value_gate", "value-continue.json"))
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "decision.pptx"
            result = render_decision_deck(package, output, resolve_tokens())
            prs = Presentation(output)
            names = [shape.name for slide in prs.slides for shape in slide.shapes]
            self.assertEqual(result.slide_count, len(prs.slides))
            self.assertEqual(0, result.picture_count)
            self.assertEqual(0, result.media_count)
            self.assertGreater(result.native_shape_count, 40)
            self.assertTrue(any(name.startswith("decision_brief.") for name in names))
            self.assertTrue(any(name.startswith("benefits_dashboard.") for name in names))
            self.assertTrue(any(name.startswith("business_case_gate.") for name in names))
            self.assertEqual((), result.out_of_bounds_text)
            with ZipFile(output) as archive:
                self.assertFalse(any(name.startswith("ppt/media/") for name in archive.namelist()))

    def test_gold_renderer_uses_bounded_gate_and_record_layouts(self) -> None:
        root = Path(__file__).resolve().parents[1]
        package = analyze_decision(load_decision_case(root / "examples" / "pmo-decision-support-gold.json"))
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "gold-layout.pptx"
            render_decision_deck(package, output, resolve_tokens())
            prs = Presentation(output)

        shapes = {shape.name: shape for slide in prs.slides for shape in slide.shapes}
        self.assertIn("Material: business_case_gate.gate.10", shapes)
        self.assertIn("business_case_gate.native.calculation_trace", shapes)
        self.assertIn("Material: decision_record.timeline_node.risk_treatment", shapes)
        self.assertIn("Material: decision_record.undated.schedule_analysis", shapes)

        deck_text = "\n".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        self.assertIn("100.41", deck_text)
        self.assertNotIn("100.4132231404958677685950413", deck_text)
        self.assertNotIn("['PPTSmith policy", deck_text)
        self.assertNotIn("recommendation_evidence_trace_missing", deck_text)
        self.assertNotIn("0.7781065088757396449704142016", deck_text)
        self.assertIn("flip threshold 0.78", deck_text)

        self.assertLess(
            shapes["decision_record.timeline.2026-08-01"].left,
            shapes["decision_record.timeline.2026-10-30"].left,
        )

    def test_gold_ir_carries_audited_values_required_by_visual_encodings(self) -> None:
        root = Path(__file__).resolve().parents[1]
        package = analyze_decision(load_decision_case(root / "examples" / "pmo-decision-support-gold.json"))
        pages = {page["page_type"]: page for page in build_decision_presentation_ir(package)}

        benefits = pages["benefits_dashboard"]["benefit_series"]
        self.assertEqual(
            {"baseline", "planned_to_date", "actual_to_date", "target", "forecast"},
            set(benefits),
        )
        self.assertEqual("124", benefits["actual_to_date"])
        self.assertEqual("120", benefits["planned_to_date"])

        scenarios = pages["capacity_scenario_allocation"]
        self.assertEqual({"team": 6.0}, scenarios["resource_limits"]["base"])
        self.assertEqual({"team": 8.0}, scenarios["resource_limits"]["relief"])
        self.assertEqual(3.0, scenarios["option_resource_demand"]["platform"]["team"])
        self.assertEqual(2.0, scenarios["option_resource_demand"]["growth"]["team"])

        gantt = pages["critical_path_gantt"]
        self.assertEqual(6, gantt["required_finish"])
        self.assertEqual(6, len(gantt["dependency_edges"]))
        self.assertEqual({"A", "B", "D", "F"}, set(gantt["critical_tasks"]))

        constraints = pages["constraint_board"]
        self.assertEqual(["C", "E", "G"], constraints["near_critical"])
        self.assertEqual({"A", "B", "D", "F"}, set(constraints["binding_tasks"]))

    def test_gold_renderer_encodes_each_judgment_as_relationship_not_card_wall(self) -> None:
        root = Path(__file__).resolve().parents[1]
        package = analyze_decision(load_decision_case(root / "examples" / "pmo-decision-support-gold.json"))
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "semantic-visuals.pptx"
            result = render_decision_deck(package, output, resolve_tokens())
            prs = Presentation(output)
        shapes = _shape_map(prs)

        # 1 Decision funnel: stages progress left-to-right and end in a distinct state.
        brief_stages = [shapes[f"Material: decision_brief.stage.{index}"] for index in range(1, 5)]
        self.assertEqual(sorted(shape.left for shape in brief_stages), [shape.left for shape in brief_stages])
        self.assertIn("decision_brief.connector.framework.1", shapes)
        self.assertIn("Material: decision_brief.terminal_state", shapes)
        # Arrow-stage labels must sit on solid fill, and the framework lanes
        # must not cover the central assessment-stage headline.
        self.assertTrue(
            all(stage.auto_shape_type == MSO_AUTO_SHAPE_TYPE.PENTAGON for stage in brief_stages[:3])
        )
        framework_nodes = [
            shapes[f"Material: decision_brief.framework.{index}"] for index in range(1, 5)
        ]
        self.assertLessEqual(
            brief_stages[2].left + brief_stages[2].width,
            min(node.left for node in framework_nodes),
        )
        self.assertGreaterEqual(
            shapes["decision_brief.stage.4.label"].width,
            Inches(1.85),
        )
        for name in (
            "recommendation_trace.lane.1.readiness",
            "recommendation_trace.lane.2.readiness",
            "recommendation_trace.lane.3.readiness",
            "recommendation_trace.lane.4.readiness",
            "decision_record.detail.value_gate.text",
            "decision_record.detail.resource_allocation.text",
            "decision_record.detail.risk_treatment.text",
        ):
            run = shapes[name].text_frame.paragraphs[0].runs[0]
            self.assertGreaterEqual(run.font.size, Pt(7.2), name)

        # 2 Bullet + bridge: actual extends beyond plan; forecast extends beyond target.
        actual = shapes["Material: benefits_dashboard.bullet.actual"]
        plan = shapes["benefits_dashboard.bullet.plan_marker"]
        target = shapes["benefits_dashboard.bridge.target_marker"]
        forecast = shapes["Material: benefits_dashboard.bridge.forecast_gap"]
        self.assertGreater(actual.left + actual.width, plan.left)
        self.assertGreater(forecast.left + forecast.width, target.left)
        self.assertGreaterEqual(shapes["benefits_dashboard.realization"].height, Inches(0.65))

        # 3 Gate path + comparable NPV bars.
        self.assertTrue(all(f"Material: business_case_gate.gate.{i}" in shapes for i in range(1, 11)))
        self.assertTrue(all(f"business_case_gate.gate.connector.{i}" in shapes for i in range(1, 10)))
        self.assertTrue(all(shapes[f"business_case_gate.gate.{i}.detail"].width >= Inches(1.10) for i in range(1, 11)))
        self.assertGreater(
            shapes["Material: business_case_gate.npv.continue"].width,
            shapes["Material: business_case_gate.npv.pivot"].width,
        )
        self.assertLess(
            shapes["Material: business_case_gate.npv.stop"].left,
            shapes["business_case_gate.npv.zero"].left,
        )

        # 4 Heat/rank/threshold: bar lengths follow fixed-anchor scores.
        ranking = [
            shapes["Material: weighted_decision_matrix.rank.platform"],
            shapes["Material: weighted_decision_matrix.rank.growth"],
            shapes["Material: weighted_decision_matrix.rank.migration"],
            shapes["Material: weighted_decision_matrix.rank.security"],
        ]
        self.assertEqual(sorted((shape.width for shape in ranking), reverse=True), [shape.width for shape in ranking])
        self.assertIn("weighted_decision_matrix.threshold", shapes)
        self.assertTrue(all(f"Material: weighted_decision_matrix.heat.{option}" in shapes for option in ("platform", "growth", "migration", "security")))

        # 5 100% stacks: all normalized bars share a frame, while unused capacity differs.
        frames = [shapes[f"Material: capacity_scenario_allocation.frame.{name}"] for name in ("base", "downside", "relief")]
        self.assertEqual(1, len({shape.width for shape in frames}))
        self.assertIn("Material: capacity_scenario_allocation.segment.downside.unused", shapes)
        self.assertGreater(shapes["Material: capacity_scenario_allocation.segment.downside.unused"].width, 0)
        self.assertIn("capacity_scenario_allocation.ceiling.base", shapes)

        # 6 Gantt: dependencies, finish constraint, and total duration are explicit.
        self.assertTrue(all(f"critical_path_gantt.dependency.{edge}" in shapes for edge in ("ab", "ac", "bd", "ce", "df", "eg")))
        self.assertGreater(shapes["critical_path_gantt.project_finish"].left, shapes["critical_path_gantt.required_finish"].left)
        self.assertIn("critical_path_gantt.total_duration", shapes)

        # 7 Constraint network: severity nodes and directed edges are explicit.
        self.assertTrue(all(f"Material: constraint_board.node.{task}" in shapes for task in ("A", "B", "C", "D", "E", "F", "G")))
        self.assertTrue(all(f"constraint_board.edge.{edge}" in shapes for edge in ("ab", "ac", "bd", "ce", "df", "eg")))
        self.assertIn("constraint_board.required_finish", shapes)

        # 8/9 Monetary relationships use common-scale bars and reference lines.
        self.assertIn("Material: risk_exposure.waterfall.baseline", shapes)
        self.assertIn("Material: risk_exposure.waterfall.reduction", shapes)
        self.assertIn("Material: risk_exposure.waterfall.residual", shapes)
        self.assertIn("risk_exposure.tolerance", shapes)
        self.assertIn("Material: mitigation_roi.bridge.gross_benefit", shapes)
        self.assertIn("Material: mitigation_roi.bridge.cost", shapes)
        self.assertIn("Material: mitigation_roi.bridge.net_benefit", shapes)

        # 10 Swimlanes preserve a connected trace per framework.
        self.assertTrue(all(f"recommendation_trace.connector.{lane}.1" in shapes for lane in range(1, 5)))
        self.assertTrue(all(f"Material: recommendation_trace.lane.{lane}.stage.4" in shapes for lane in range(1, 5)))

        # 11 Review dates map onto a timeline; missing date is kept out of the scale.
        self.assertLess(
            shapes["decision_record.timeline.2026-08-01"].left,
            shapes["decision_record.timeline.2026-08-31"].left,
        )
        self.assertLess(
            shapes["decision_record.timeline.2026-08-31"].left,
            shapes["decision_record.timeline.2026-10-30"].left,
        )
        self.assertIn("Material: decision_record.undated.schedule_analysis", shapes)

        self.assertEqual(0, result.picture_count)
        self.assertEqual(0, result.media_count)
        self.assertGreaterEqual(result.connector_count, 45)


if __name__ == "__main__":
    unittest.main()
