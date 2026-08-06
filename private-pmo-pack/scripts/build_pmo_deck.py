#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from jsonschema import Draft202012Validator
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_INPUT = ROOT / "examples" / "ai-platform-project.pmo.json"
DEFAULT_SCHEMA = ROOT / "schemas" / "pmo-input.schema.json"
DEFAULT_STYLE = ROOT / "contracts" / "mckinsey-pmo.style.json"
DEFAULT_OUTPUT = ROOT / "output" / "ai-platform-project-pmo.pptx"
DEFAULT_MANIFEST = ROOT / "qa" / "build-manifest.json"


def load_json(path: Path) -> dict[str, Any]:
    with path.open("r", encoding="utf-8") as handle:
        data = json.load(handle)
    if not isinstance(data, dict):
        raise ValueError(f"{path} must contain a JSON object")
    return data


def write_json(path: Path, data: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def color(style: dict[str, Any], name: str) -> RGBColor:
    return RGBColor.from_string(style["colors"][name])


def status_color(style: dict[str, Any], status: str) -> RGBColor:
    return RGBColor.from_string(style["status_colors"].get(status, style["status_colors"]["gray"]))


class Deck:
    def __init__(self, data: dict[str, Any], style: dict[str, Any]) -> None:
        self.data = data
        self.style = style
        self.prs = Presentation()
        slide = style["slide"]
        self.prs.slide_width = Inches(slide["width_in"])
        self.prs.slide_height = Inches(slide["height_in"])
        self.blank = self.prs.slide_layouts[6]
        self.slide_records: list[dict[str, Any]] = []

    def save(self, path: Path) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(path)

    def add_slide(self, title: str, page_type: str, page_no: int, kicker: str | None = None) -> Any:
        slide = self.prs.slides.add_slide(self.blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color(self.style, "background")
        if kicker:
            self.text(slide, kicker.upper(), 0.56, 0.25, 4.8, 0.2, size=8.5, color_name="blue", bold=True)
        self.text(slide, title, 0.56, 0.5, 11.95, 0.5, size=self.style["typography"]["title_pt"], color_name="ink", bold=True)
        self.line(slide, 0.56, 1.12, 12.78, 1.12, "line", 0.7)
        self.text(slide, f"PMO CONSULTING PACK / {page_no:02d}", 10.0, 7.12, 2.72, 0.18, size=7.8, color_name="muted", bold=True, align=PP_ALIGN.RIGHT)
        self.slide_records.append({"index": page_no, "page_type": page_type, "title": title})
        return slide

    def text(
        self,
        slide: Any,
        value: str,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        size: float | None = None,
        color_name: str = "text",
        bold: bool = False,
        align: PP_ALIGN = PP_ALIGN.LEFT,
        valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    ) -> Any:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.03)
        frame.margin_right = Inches(0.03)
        frame.margin_top = Inches(0.01)
        frame.margin_bottom = Inches(0.01)
        frame.vertical_anchor = valign
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = str(value)
        run.font.name = self.style["typography"]["font_cjk"]
        run.font.size = Pt(size or self.style["typography"]["body_pt"])
        run.font.bold = bold
        run.font.color.rgb = color(self.style, color_name)
        return box

    def rect(self, slide: Any, x: float, y: float, w: float, h: float, *, fill: str = "white", line: str = "line", radius: bool = False, width: float = 0.7) -> Any:
        kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
        shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color(self.style, fill)
        shape.line.color.rgb = color(self.style, line)
        shape.line.width = Pt(width)
        return shape

    def status_dot(self, slide: Any, x: float, y: float, status: str, size: float = 0.16) -> Any:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        dot.fill.solid()
        dot.fill.fore_color.rgb = status_color(self.style, status)
        dot.line.color.rgb = status_color(self.style, status)
        return dot

    def line(self, slide: Any, x1: float, y1: float, x2: float, y2: float, color_name: str = "line", width: float = 1.0) -> Any:
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        connector.line.color.rgb = color(self.style, color_name)
        connector.line.width = Pt(width)
        return connector

    def card(self, slide: Any, x: float, y: float, w: float, h: float, title: str, body: str, *, accent_status: str | None = None, accent: str = "blue") -> None:
        self.rect(slide, x, y, w, h, fill="white", line="line")
        if accent_status:
            self.status_dot(slide, x + 0.16, y + 0.18, accent_status, 0.15)
            title_x = x + 0.42
        else:
            self.rect(slide, x, y, 0.06, h, fill=accent, line=accent, width=0)
            title_x = x + 0.2
        self.text(slide, title, title_x, y + 0.15, w - (title_x - x) - 0.16, 0.24, size=11.5, color_name="ink", bold=True)
        self.text(slide, body, x + 0.2, y + 0.52, w - 0.36, h - 0.62, size=9.6, color_name="text")

    def table(self, slide: Any, rows: list[list[str]], x: float, y: float, w: float, h: float, *, widths: list[float] | None = None, status_col: int | None = None) -> Any:
        shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
        table = shape.table
        if widths:
            for idx, col_width in enumerate(widths):
                table.columns[idx].width = Inches(col_width)
        for r_idx, row in enumerate(rows):
            for c_idx, value in enumerate(row):
                cell = table.cell(r_idx, c_idx)
                cell.text = str(value)
                fill = "navy" if r_idx == 0 else ("paper" if r_idx % 2 == 0 else "white")
                cell.fill.solid()
                cell.fill.fore_color.rgb = color(self.style, fill)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER if c_idx == status_col else PP_ALIGN.LEFT
                    for run in paragraph.runs:
                        run.font.name = self.style["typography"]["font_cjk"]
                        run.font.size = Pt(8.2 if len(rows) > 6 else 9.2)
                        run.font.bold = r_idx == 0
                        run.font.color.rgb = color(self.style, "white" if r_idx == 0 else "text")
                if status_col is not None and r_idx > 0 and c_idx == status_col:
                    status = str(value).lower()
                    if status in {"green", "yellow", "red", "gray"}:
                        cell.fill.fore_color.rgb = status_color(self.style, status)
        return shape


def render_cover(deck: Deck) -> None:
    data = deck.data
    project = data["project"]
    slide = deck.prs.slides.add_slide(deck.blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color(deck.style, "background")
    deck.text(slide, "PMO CONSULTING PACK", 0.7, 0.55, 4.4, 0.22, size=9, color_name="blue", bold=True)
    deck.text(slide, project["name"], 0.7, 1.15, 7.0, 0.78, size=29, color_name="ink", bold=True)
    deck.text(slide, project["headline"], 0.72, 2.12, 8.2, 0.44, size=14, color_name="muted")
    deck.rect(slide, 0.72, 4.62, 6.0, 0.82, fill="navy", line="navy")
    deck.text(slide, project["decision_ask"], 0.92, 4.88, 5.62, 0.24, size=15, color_name="white", bold=True, align=PP_ALIGN.CENTER)
    labels = [("Dashboard", "status and risks"), ("Roadmap", "phase plan"), ("Decision", "recommended path"), ("Tracker", "owners and dates")]
    for idx, (label, note) in enumerate(labels):
        x = 8.0 + (idx % 2) * 2.28
        y = 1.3 + (idx // 2) * 1.55
        deck.rect(slide, x, y, 1.92, 0.95, fill="paper", line="line")
        deck.text(slide, label, x + 0.12, y + 0.22, 1.68, 0.22, size=12, color_name="ink", bold=True, align=PP_ALIGN.CENTER)
        deck.text(slide, note, x + 0.12, y + 0.55, 1.68, 0.2, size=8.2, color_name="muted", align=PP_ALIGN.CENTER)
    deck.text(slide, f"{project['client']} / {project['period']}", 0.72, 6.86, 4.5, 0.18, size=8.5, color_name="muted", bold=True)
    deck.slide_records.append({"index": 1, "page_type": "cover", "title": project["name"]})


def render_executive_dashboard(deck: Deck) -> None:
    section = deck.data["executive_dashboard"]
    slide = deck.add_slide("项目总体健康度：工程按计划推进，商业证据仍需补齐", "executive_dashboard", 2, "Executive dashboard")
    deck.text(slide, section["summary"], 0.62, 1.34, 9.8, 0.28, size=12, color_name="muted")
    deck.rect(slide, 10.65, 1.28, 1.7, 0.36, fill="paper", line="line", radius=True)
    deck.status_dot(slide, 10.84, 1.38, section["overall_status"], 0.14)
    deck.text(slide, f"Overall: {section['overall_status'].upper()}", 11.08, 1.35, 1.1, 0.18, size=8.6, color_name="ink", bold=True)
    for idx, metric in enumerate(section["metrics"]):
        x = 0.7 + idx * 2.38
        deck.card(slide, x, 1.95, 2.05, 1.35, metric["label"], f"{metric['value']}\n{metric['note']}", accent_status=metric["status"])
    risk_rows = [["Risk", "Impact", "Owner", "Mitigation", "Status"]]
    for risk in section["top_risks"]:
        risk_rows.append([risk["risk"], risk["impact"], risk["owner"], risk["mitigation"], risk["status"]])
    deck.table(slide, risk_rows, 0.72, 3.76, 7.35, 2.18, widths=[1.65, 0.78, 0.82, 3.05, 0.78], status_col=4)
    deck.rect(slide, 8.42, 3.76, 3.9, 2.18, fill="paper", line="line")
    deck.text(slide, "Next actions", 8.68, 4.0, 3.2, 0.22, size=12, color_name="ink", bold=True)
    for idx, action in enumerate(section["next_actions"]):
        y = 4.42 + idx * 0.34
        deck.status_dot(slide, 8.72, y + 0.02, "gray", 0.11)
        deck.text(slide, action, 8.92, y, 3.08, 0.18, size=9.3, color_name="text")


def render_roadmap(deck: Deck) -> None:
    slide = deck.add_slide("路线图不是甘特图：用阶段、目标和门禁表达管理节奏", "roadmap", 3, "Project roadmap")
    phases = deck.data["roadmap"]
    left, top, gap = 0.78, 2.05, 0.18
    width = (11.8 - gap * (len(phases) - 1)) / len(phases)
    for idx, phase in enumerate(phases):
        x = left + idx * (width + gap)
        deck.rect(slide, x, top, width, 2.2, fill="white", line="line")
        deck.status_dot(slide, x + 0.16, top + 0.2, phase["status"], 0.13)
        deck.text(slide, phase["period"], x + 0.38, top + 0.17, width - 0.5, 0.18, size=8.5, color_name="muted", bold=True)
        deck.text(slide, phase["name"], x + 0.16, top + 0.58, width - 0.32, 0.28, size=13, color_name="ink", bold=True, align=PP_ALIGN.CENTER)
        deck.text(slide, phase["objective"], x + 0.18, top + 1.22, width - 0.35, 0.62, size=9.2, color_name="text", align=PP_ALIGN.CENTER)
        if idx < len(phases) - 1:
            deck.line(slide, x + width, top + 1.1, x + width + gap, top + 1.1, "blue", 2.0)
    deck.rect(slide, 1.1, 5.3, 10.9, 0.72, fill="navy", line="navy")
    deck.text(slide, "Management rule: every phase ends with a decision gate, not a feature-complete declaration.", 1.34, 5.54, 10.42, 0.2, size=12.5, color_name="white", bold=True, align=PP_ALIGN.CENTER)


def render_milestones(deck: Deck) -> None:
    slide = deck.add_slide("里程碑页把时间点、证据和状态压到一条水平轴上", "milestones", 4, "Milestones")
    milestones = deck.data["milestones"]
    y_axis = 3.15
    deck.line(slide, 0.95, y_axis, 12.0, y_axis, "line", 1.5)
    step = 10.7 / max(len(milestones) - 1, 1)
    for idx, item in enumerate(milestones):
        x = 0.95 + idx * step
        deck.status_dot(slide, x - 0.08, y_axis - 0.08, item["status"], 0.16)
        deck.text(slide, item["date"], x - 0.45, 2.35, 0.9, 0.2, size=8.5, color_name="muted", bold=True, align=PP_ALIGN.CENTER)
        deck.text(slide, item["name"], x - 0.68, 2.66, 1.36, 0.28, size=10.5, color_name="ink", bold=True, align=PP_ALIGN.CENTER)
        deck.text(slide, item["evidence"], x - 0.72, 3.52, 1.44, 0.48, size=8.2, color_name="text", align=PP_ALIGN.CENTER)
    deck.rect(slide, 1.25, 5.25, 10.35, 0.78, fill="paper", line="line")
    deck.text(slide, "Milestone discipline: dates are only useful when paired with evidence, owner, and decision consequence.", 1.48, 5.51, 9.9, 0.2, size=12, color_name="ink", bold=True, align=PP_ALIGN.CENTER)


def render_raid(deck: Deck) -> None:
    slide = deck.add_slide("RAID 页优先做成可查表格，避免把风险画成装饰图", "raid_log", 5, "RAID log")
    rows = [["Type", "Item", "Impact", "Owner", "Mitigation", "Status"]]
    for item in deck.data["raid_log"]:
        rows.append([item["type"], item["item"], item["impact"], item["owner"], item["mitigation"], item["status"]])
    deck.table(slide, rows, 0.7, 1.55, 11.85, 4.28, widths=[1.05, 2.55, 0.9, 1.0, 4.45, 0.9], status_col=5)
    deck.text(slide, "RAID governance check", 0.75, 6.1, 2.2, 0.2, size=11.5, color_name="ink", bold=True)
    deck.text(slide, "Each row must have a type, owner, mitigation path, and current status. Empty owner or mitigation fields fail PMO QA.", 2.75, 6.1, 8.95, 0.2, size=10, color_name="muted")


def render_workstreams(deck: Deck) -> None:
    slide = deck.add_slide("工作流页用泳道表达责任边界，而不是堆任务清单", "workstreams", 6, "Workstreams")
    streams = deck.data["workstreams"]
    lane_h = 0.9
    for idx, stream in enumerate(streams):
        y = 1.55 + idx * 1.12
        deck.rect(slide, 0.72, y, 2.1, lane_h, fill="navy", line="navy")
        deck.text(slide, stream["name"], 0.9, y + 0.18, 1.55, 0.2, size=11.5, color_name="white", bold=True)
        deck.text(slide, f"Owner: {stream['owner']}", 0.9, y + 0.5, 1.55, 0.16, size=7.8, color_name="white")
        deck.status_dot(slide, 2.48, y + 0.33, stream["status"], 0.14)
        for jdx, deliverable in enumerate(stream["deliverables"]):
            x = 3.15 + jdx * 2.3
            deck.rect(slide, x, y + 0.08, 1.9, 0.62, fill="paper", line="line")
            deck.text(slide, deliverable, x + 0.1, y + 0.27, 1.7, 0.16, size=8.3, color_name="text", align=PP_ALIGN.CENTER)
            if jdx < len(stream["deliverables"]) - 1:
                deck.line(slide, x + 1.9, y + 0.39, x + 2.3, y + 0.39, "line", 1.0)


def render_raci(deck: Deck) -> None:
    slide = deck.add_slide("RACI 是精确查阅页：角色和责任必须保持表格可编辑", "raci_matrix", 7, "RACI matrix")
    raci = deck.data["raci"]
    rows = [["Task", *raci["roles"]]]
    for task in raci["tasks"]:
        rows.append([task["task"], *task["assignments"][: len(raci["roles"])]])
    widths = [3.3, *([1.25] * len(raci["roles"]))]
    deck.table(slide, rows, 0.82, 1.55, 11.0, 4.7, widths=widths)
    legend = [("R", "Responsible"), ("A", "Accountable"), ("C", "Consulted"), ("I", "Informed")]
    for idx, (code, label) in enumerate(legend):
        x = 1.0 + idx * 2.55
        deck.rect(slide, x, 6.35, 0.32, 0.26, fill="navy", line="navy")
        deck.text(slide, code, x, 6.41, 0.32, 0.12, size=8, color_name="white", bold=True, align=PP_ALIGN.CENTER)
        deck.text(slide, label, x + 0.42, 6.38, 1.7, 0.16, size=8.6, color_name="muted")


def render_heat_map(deck: Deck) -> None:
    slide = deck.add_slide("风险热力图是坐标页：影响和可能性必须能被复核", "risk_heat_map", 8, "Risk heat map")
    risks = deck.data["risk_heat_map"]
    start_x, start_y, cell = 1.35, 1.6, 0.72
    for impact in range(5, 0, -1):
        for likelihood in range(1, 6):
            score = impact * likelihood
            fill = "green" if score <= 6 else ("yellow" if score <= 14 else "red")
            x = start_x + (likelihood - 1) * cell
            y = start_y + (5 - impact) * cell
            deck.rect(slide, x, y, cell, cell, fill=fill, line="white")
    deck.text(slide, "Impact", 0.58, 2.9, 0.5, 0.2, size=9, color_name="muted", bold=True)
    deck.text(slide, "Likelihood", 2.3, 5.38, 1.4, 0.2, size=9, color_name="muted", bold=True, align=PP_ALIGN.CENTER)
    for risk in risks:
        x = start_x + (risk["likelihood"] - 1) * cell + 0.18
        y = start_y + (5 - risk["impact"]) * cell + 0.18
        deck.rect(slide, x, y, 0.36, 0.36, fill="white", line="ink", radius=True, width=1)
        deck.text(slide, risk["id"], x, y + 0.1, 0.36, 0.12, size=7.2, color_name="ink", bold=True, align=PP_ALIGN.CENTER)
    rows = [["ID", "Risk", "Owner", "I", "L"]]
    for risk in risks:
        rows.append([risk["id"], risk["label"], risk["owner"], str(risk["impact"]), str(risk["likelihood"])])
    deck.table(slide, rows, 5.45, 1.62, 6.4, 3.85, widths=[0.55, 2.9, 1.05, 0.45, 0.45])
    deck.text(slide, "PMO QA: all risk coordinates must be 1-5 and every point must have an owner.", 5.5, 5.82, 5.8, 0.2, size=9.5, color_name="muted")


def render_decision(deck: Deck) -> None:
    slide = deck.add_slide("决策页必须给出选项、权衡和推荐，而不是只列利弊", "decision_page", 9, "Decision page")
    decision = deck.data["decision"]
    deck.text(slide, decision["question"], 0.72, 1.38, 10.5, 0.28, size=13.5, color_name="ink", bold=True)
    for idx, option in enumerate(decision["options"]):
        x = 0.78 + idx * 4.0
        deck.rect(slide, x, 2.02, 3.55, 2.95, fill="white", line="line")
        deck.text(slide, option["name"], x + 0.18, 2.2, 2.45, 0.24, size=12.5, color_name="ink", bold=True)
        deck.text(slide, f"Score {option['score']}/5", x + 2.62, 2.22, 0.7, 0.16, size=8.5, color_name="blue", bold=True, align=PP_ALIGN.RIGHT)
        deck.text(slide, "Pros", x + 0.2, 2.72, 0.6, 0.16, size=8.5, color_name="green", bold=True)
        deck.text(slide, "\n".join(f"- {item}" for item in option["pros"]), x + 0.2, 2.98, 3.05, 0.55, size=8.6, color_name="text")
        deck.text(slide, "Cons", x + 0.2, 3.78, 0.6, 0.16, size=8.5, color_name="red", bold=True)
        deck.text(slide, "\n".join(f"- {item}" for item in option["cons"]), x + 0.2, 4.04, 3.05, 0.55, size=8.6, color_name="text")
    deck.rect(slide, 1.1, 5.75, 11.0, 0.62, fill="navy", line="navy")
    deck.text(slide, decision["recommendation"], 1.35, 5.95, 10.48, 0.18, size=12.8, color_name="white", bold=True, align=PP_ALIGN.CENTER)


def render_kpis(deck: Deck) -> None:
    slide = deck.add_slide("KPI Dashboard 用大数字和一张原生图表支撑管理判断", "kpi_dashboard", 10, "KPI dashboard")
    kpis = deck.data["kpis"]
    for idx, kpi in enumerate(kpis):
        x = 0.72 + idx * 2.38
        deck.rect(slide, x, 1.55, 2.05, 1.35, fill="white", line="line")
        deck.status_dot(slide, x + 0.16, 1.75, kpi["status"], 0.14)
        deck.text(slide, kpi["label"], x + 0.38, 1.72, 1.4, 0.18, size=8.7, color_name="muted", bold=True)
        deck.text(slide, kpi["value"], x + 0.18, 2.1, 1.7, 0.28, size=19.5, color_name="ink", bold=True, align=PP_ALIGN.CENTER)
        deck.text(slide, kpi["note"], x + 0.18, 2.52, 1.7, 0.16, size=7.8, color_name="muted", align=PP_ALIGN.CENTER)
    chart_data = ChartData()
    chart_data.categories = [kpi["label"] for kpi in kpis]
    chart_data.add_series("Readiness", [_status_score(kpi["status"]) for kpi in kpis])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.05), Inches(3.55), Inches(10.8), Inches(2.35), chart_data).chart
    chart.has_legend = False
    chart.value_axis.maximum_scale = 5
    chart.value_axis.minimum_scale = 0
    chart.value_axis.major_unit = 1
    chart.plots[0].has_data_labels = True


def render_actions(deck: Deck) -> None:
    slide = deck.add_slide("行动项追踪页像 Excel，但必须有清晰 owner、deadline 和 status", "action_tracker", 11, "Action tracker")
    rows = [["Owner", "Action", "Deadline", "Status"]]
    for action in deck.data["actions"]:
        rows.append([action["owner"], action["action"], action["deadline"], action["status"]])
    deck.table(slide, rows, 0.78, 1.55, 11.65, 4.7, widths=[1.2, 6.85, 1.55, 0.95], status_col=3)
    deck.rect(slide, 1.15, 6.35, 10.8, 0.5, fill="paper", line="line")
    deck.text(slide, "Next review: update status before the weekly PMO meeting; overdue red rows require a mitigation owner.", 1.38, 6.52, 10.35, 0.16, size=9.4, color_name="muted", align=PP_ALIGN.CENTER)


def _status_score(status: str) -> int:
    return {"green": 5, "yellow": 3, "red": 1, "gray": 2}.get(status, 2)


def build(data: dict[str, Any], style: dict[str, Any], output: Path, manifest_path: Path) -> dict[str, Any]:
    deck = Deck(data, style)
    render_cover(deck)
    render_executive_dashboard(deck)
    render_roadmap(deck)
    render_milestones(deck)
    render_raid(deck)
    render_workstreams(deck)
    render_raci(deck)
    render_heat_map(deck)
    render_decision(deck)
    render_kpis(deck)
    render_actions(deck)
    deck.save(output)
    manifest = {
        "schema_version": "1.0",
        "pack_id": "pptsmith-pmo-consulting-pack",
        "built_at": datetime.now(timezone.utc).isoformat(),
        "status": "created",
        "builder": "private-python-pptx-pmo-renderer",
        "outputs": {"deck": str(output.relative_to(ROOT)) if output.is_relative_to(ROOT) else str(output)},
        "source": str(DEFAULT_INPUT.relative_to(ROOT)),
        "style": str(DEFAULT_STYLE.relative_to(ROOT)),
        "slide_count": len(deck.prs.slides),
        "slides": deck.slide_records,
        "editable_core_expected": [
            "titles",
            "metric cards",
            "native tables",
            "roadmap phase shapes",
            "milestone nodes",
            "risk heat-map grid and labels",
            "native KPI chart",
            "action tracker rows"
        ],
        "known_limitations": [
            "v0.1 renderer writes directly from PMO domain JSON; the planned PPTSmith PPT IR adapter is not complete yet.",
            "Connector endpoint binding is not claimed for v0.1; relationship-heavy PMO diagrams should move through Diagram IR in v0.2.",
            "Risk heat map is rendered as native shape grid plus text markers, not yet as a first-class public registry component."
        ]
    }
    write_json(manifest_path, manifest)
    return manifest


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build an editable McKinsey-style PMO PPTX from domain JSON.")
    parser.add_argument("--input", type=Path, default=DEFAULT_INPUT)
    parser.add_argument("--schema", type=Path, default=DEFAULT_SCHEMA)
    parser.add_argument("--style", type=Path, default=DEFAULT_STYLE)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--manifest", type=Path, default=DEFAULT_MANIFEST)
    return parser.parse_args(argv)


def main(argv: list[str]) -> int:
    args = parse_args(argv)
    data = load_json(args.input)
    schema = load_json(args.schema)
    style = load_json(args.style)
    Draft202012Validator(schema).validate(data)
    manifest = build(data, style, args.output, args.manifest)
    print(json.dumps({"deck": str(args.output), "manifest": str(args.manifest), "slide_count": manifest["slide_count"]}, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
