#!/usr/bin/env python3
from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "output" / "dashboard-auto-selection-validation.pptx"
MANIFEST = ROOT / "qa" / "selection-manifest.json"
W, H = 13.333, 7.5
FONT = "PingFang SC"
COLORS = {
    "navy": "17324D", "blue": "2D6792", "cyan": "6AA6BF", "ink": "182632",
    "text": "344655", "muted": "6B7A88", "line": "D2DBE2", "paper": "F4F7F9",
    "white": "FFFFFF", "green": "2E7D5B", "green_soft": "DDEDE6",
    "amber": "D99325", "amber_soft": "F7E8CB", "red": "C7463D",
    "red_soft": "F4DDDA", "gray": "9AA7B2", "gray_soft": "E8EDF1",
}


def rgb(name: str) -> RGBColor:
    return RGBColor.from_string(COLORS[name])


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(W)
        self.prs.slide_height = Inches(H)
        self.blank = self.prs.slide_layouts[6]

    def base(self, title: str, page: int, family: str) -> Any:
        slide = self.prs.slides.add_slide(self.blank)
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb("white")
        self.text(slide, family.upper(), 0.58, 0.20, 5.2, 0.18, 8.2, "blue", True)
        self.text(slide, title, 0.58, 0.47, 12.0, 0.54, 20.5, "ink", True)
        self.rect(slide, 0.58, 1.18, 12.18, 0.012, "line", "line")
        self.text(slide, "AUTO-SELECTION VALIDATION", 0.58, 7.14, 3.0, 0.14, 7.2, "muted", True)
        self.text(slide, f"PMO PACK / {page:02d}", 11.10, 7.14, 1.65, 0.14, 7.2, "muted", True, PP_ALIGN.RIGHT)
        return slide

    def text(self, slide: Any, value: str, x: float, y: float, w: float, h: float,
             size: float = 11, color: str = "text", bold: bool = False,
             align: PP_ALIGN = PP_ALIGN.LEFT) -> Any:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = shape.text_frame; frame.clear(); frame.word_wrap = True
        frame.margin_left = Inches(0.02); frame.margin_right = Inches(0.02)
        frame.margin_top = Inches(0.01); frame.margin_bottom = Inches(0.01)
        p = frame.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = str(value); run.font.name = FONT
        run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = rgb(color)
        return shape

    def rect(self, slide: Any, x: float, y: float, w: float, h: float,
             fill: str = "white", line: str = "line", width: float = 0.7) -> Any:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
        shape.line.color.rgb = rgb(line); shape.line.width = Pt(width)
        return shape

    def pill(self, slide: Any, x: float, y: float, label: str, status: str, w: float = 0.82) -> None:
        soft = {"green": "green_soft", "amber": "amber_soft", "red": "red_soft", "gray": "gray_soft"}[status]
        self.rect(slide, x, y, w, 0.27, soft, soft, 0)
        self.text(slide, label, x, y + 0.04, w, 0.15, 8.0, status, True, PP_ALIGN.CENTER)

    def takeaway(self, slide: Any, value: str) -> None:
        self.rect(slide, 0.76, 6.35, 11.82, 0.48, "navy", "navy", 0)
        self.text(slide, value, 0.96, 6.47, 11.42, 0.20, 11.5, "white", True, PP_ALIGN.CENTER)

    def table(self, slide: Any, rows: list[list[str]], x: float, y: float, w: float, h: float,
              widths: list[float], center: set[int] | None = None, font_size: float = 9.3) -> Any:
        shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
        table = shape.table; center = center or set()
        for idx, width in enumerate(widths): table.columns[idx].width = Inches(width)
        for r, row in enumerate(rows):
            for c, value in enumerate(row):
                cell = table.cell(r, c); cell.text = value; cell.fill.solid()
                cell.fill.fore_color.rgb = rgb("navy" if r == 0 else ("paper" if r % 2 == 0 else "white"))
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER if c in center else PP_ALIGN.LEFT
                    for run in p.runs:
                        run.font.name = FONT; run.font.size = Pt(font_size); run.font.bold = r == 0
                        run.font.color.rgb = rgb("white" if r == 0 else "text")
        return shape


def cover(d: Deck) -> None:
    s = d.prs.slides.add_slide(d.blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb("white")
    d.rect(s, 0, 0, 4.18, H, "navy", "navy", 0)
    d.text(s, "PMO DASHBOARD\nPAGE FAMILIES", 0.62, 0.82, 3.05, 1.10, 24, "white", True)
    d.text(s, "文章证据自动选图验收", 0.64, 2.30, 2.95, 0.34, 15, "cyan", True)
    d.text(s, "不是模板轮播，而是根据管理问题、数据形状与决策任务选择页面表达。", 0.64, 4.24, 2.95, 0.90, 14, "white", True)
    families = ["单项目健康", "项目组合", "时间与资源", "PMO 控制塔", "Executive One-Pager"]
    d.text(s, "5", 5.00, 1.04, 1.1, 0.70, 34, "navy", True, PP_ALIGN.CENTER)
    d.text(s, "Dashboard 页面族", 4.72, 1.66, 1.65, 0.22, 10.5, "muted", True, PP_ALIGN.CENTER)
    for i, family in enumerate(families):
        y = 2.34 + i * 0.76
        d.text(s, f"0{i+1}", 5.02, y, 0.46, 0.22, 10, "blue", True)
        d.rect(s, 5.60, y - 0.08, 6.30, 0.48, "paper", "line")
        d.text(s, family, 5.84, y + 0.04, 5.80, 0.22, 12.5, "ink", True)
    d.text(s, "v0.2 / PRIVATE PRODUCTION PACK", 0.64, 6.88, 3.0, 0.16, 7.8, "gray", True)


def single_project(d: Deck) -> None:
    s = d.base("项目总体为 Amber：数据依赖导致关键路径晚 10 天，预算仍在控制范围", 2, "single project health dashboard")
    info = [("项目", "企业内容平台"), ("Sponsor", "COO"), ("周期", "W1-W24"), ("阶段", "Pilot")]
    for i, (k, v) in enumerate(info):
        x = 0.76 + i * 2.15
        d.text(s, k, x, 1.42, 0.74, 0.18, 8.5, "muted", True)
        d.text(s, v, x, 1.68, 1.85, 0.22, 11.5, "ink", True)
    d.pill(s, 10.52, 1.50, "AMBER", "amber", 1.20)
    metrics = [("进度", "62%", "-6pp", "amber"), ("预算", "6.8M", "+6%", "amber"), ("质量", "78%", "-12pp", "red"), ("价值", "17.4M", "67%目标", "green")]
    for i, (name, value, note, status) in enumerate(metrics):
        x = 0.76 + i * 2.28
        d.rect(s, x, 2.16, 2.04, 1.18, "white", "line")
        d.pill(s, x + 0.14, 2.34, name, status, 0.68)
        d.text(s, value, x + 0.16, 2.68, 1.72, 0.38, 18, "ink", True, PP_ALIGN.CENTER)
        d.text(s, note, x + 0.16, 3.05, 1.72, 0.16, 8.5, status, True, PP_ALIGN.CENTER)
    d.rect(s, 10.18, 2.16, 2.10, 1.18, "red_soft", "red_soft", 0)
    d.text(s, "关键例外", 10.42, 2.38, 1.62, 0.20, 10, "red", True)
    d.text(s, "主数据 owner 未确认", 10.42, 2.75, 1.62, 0.40, 10.5, "ink", True)
    rows = [["管理事项", "影响", "Owner", "下一步", "截止"], ["数据覆盖 71%", "门禁延迟", "CDO", "确认 3 个数据域", "07/30"], ["采用率 64%", "价值受限", "COO", "主管责任制", "08/02"]]
    d.table(s, rows, 0.76, 3.76, 8.25, 1.75, [2.05, 1.45, 1.05, 2.60, 1.10], {2, 4}, 9.5)
    d.text(s, "交付时间线", 9.42, 3.72, 2.00, 0.22, 11.5, "ink", True)
    phases = [("Build", 0.82, "green"), ("Pilot", 0.54, "amber"), ("Gate", 0.12, "gray")]
    for i, (name, pct, status) in enumerate(phases):
        y = 4.16 + i * 0.48
        d.text(s, name, 9.42, y, 0.62, 0.18, 9.5, "text", True)
        d.rect(s, 10.16, y + 0.02, 1.88, 0.18, "gray_soft", "gray_soft", 0)
        d.rect(s, 10.16, y + 0.02, 1.88 * pct, 0.18, status, status, 0)
    d.takeaway(s, "管理动作：不追加总预算；把非关键功能资源转向数据覆盖和采用率。")


def portfolio(d: Deck) -> None:
    s = d.base("六个项目中两项占用 57% 预算却处于 Red/Amber，应优先重配资源", 3, "portfolio dashboard")
    rows = [["项目", "优先级", "健康度", "完成率", "预算", "已花费", "建议"], ["AI 内容平台", "P0", "Amber", "62%", "12.0M", "6.8M", "保留"], ["数据中台", "P0", "Red", "48%", "8.5M", "5.2M", "增配"], ["CRM 升级", "P1", "Green", "74%", "5.2M", "3.4M", "维持"], ["渠道自动化", "P1", "Amber", "55%", "4.8M", "2.9M", "收窄"], ["知识门户", "P2", "Green", "83%", "2.6M", "2.1M", "收尾"], ["实验项目", "P2", "Gray", "21%", "3.0M", "0.8M", "暂停"]]
    shape = d.table(s, rows, 0.72, 1.46, 8.82, 4.45, [2.12, 0.88, 1.05, 1.05, 1.15, 1.20, 1.37], {1, 2, 3, 4, 5, 6}, 8.9)
    for r, fill in [(1, "amber_soft"), (2, "red_soft"), (4, "amber_soft")]: shape.table.cell(r, 2).fill.fore_color.rgb = rgb(fill)
    data = ChartData(); data.categories = ["AI平台", "数据中台", "CRM", "渠道", "门户", "实验"]
    data.add_series("预算", [12, 8.5, 5.2, 4.8, 2.6, 3.0])
    chart = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(9.82), Inches(1.52), Inches(2.45), Inches(2.10), data).chart
    chart.has_legend = False; chart.plots[0].hole_size = 58
    d.text(s, "预算组合", 10.26, 1.40, 1.55, 0.22, 11.5, "ink", True, PP_ALIGN.CENTER)
    d.text(s, "57%", 10.55, 2.38, 0.96, 0.30, 19, "navy", True, PP_ALIGN.CENTER)
    d.text(s, "集中于前两项", 10.30, 2.76, 1.47, 0.18, 8.5, "muted", True, PP_ALIGN.CENTER)
    d.rect(s, 9.82, 3.95, 2.45, 1.48, "paper", "line")
    d.text(s, "资源优先级", 10.06, 4.18, 1.95, 0.22, 11.5, "ink", True)
    d.text(s, "1. 数据中台增配\n2. 渠道项目收窄\n3. 实验项目暂停", 10.06, 4.58, 1.95, 0.65, 10.2, "text", True)
    d.takeaway(s, "管理动作：从实验项目和渠道项目释放 1.4M，优先保护数据中台关键路径。")


def resources(d: Deck) -> None:
    s = d.base("数据工程负荷达到 118%，而设计资源仅 63%；当前问题是结构错配，不是总人数不足", 4, "time and resource dashboard")
    teams = [("数据工程", 118, "red"), ("平台工程", 92, "amber"), ("产品", 84, "green"), ("运营", 76, "green"), ("设计", 63, "gray")]
    d.text(s, "团队利用率", 0.78, 1.42, 1.8, 0.22, 12, "ink", True)
    for i, (team, value, status) in enumerate(teams):
        y = 1.88 + i * 0.62
        d.text(s, team, 0.78, y, 1.18, 0.20, 10.5, "text", True)
        d.rect(s, 2.08, y + 0.02, 4.25, 0.22, "gray_soft", "gray_soft", 0)
        d.rect(s, 2.08, y + 0.02, 4.25 * min(value, 120) / 120, 0.22, status, status, 0)
        d.text(s, f"{value}%", 6.50, y - 0.01, 0.62, 0.22, 10.5, status, True, PP_ALIGN.RIGHT)
    d.rect(s, 7.52, 1.45, 4.72, 2.00, "paper", "line")
    d.text(s, "未来两周关键任务", 7.82, 1.73, 4.10, 0.22, 12, "ink", True)
    tasks = [("主数据接入", "Data", "07/30", "red"), ("规则引擎 UAT", "Eng", "08/05", "amber"), ("主管试点", "Ops", "08/08", "amber")]
    for i, (name, owner, due, status) in enumerate(tasks):
        y = 2.18 + i * 0.39
        d.pill(s, 7.82, y, f"0{i+1}", status, 0.38)
        d.text(s, name, 8.34, y + 0.03, 1.82, 0.17, 9.5, "text", True)
        d.text(s, owner, 10.32, y + 0.03, 0.72, 0.17, 8.8, "muted", True)
        d.text(s, due, 11.16, y + 0.03, 0.72, 0.17, 8.8, status, True, PP_ALIGN.RIGHT)
    d.text(s, "预算消耗", 7.54, 3.86, 1.2, 0.22, 12, "ink", True)
    d.rect(s, 7.54, 4.28, 4.70, 0.30, "gray_soft", "gray_soft", 0); d.rect(s, 7.54, 4.28, 2.66, 0.30, "blue", "blue", 0)
    d.text(s, "56.7% 已花费", 7.54, 4.72, 2.10, 0.20, 10.5, "blue", True)
    d.text(s, "43.3% 剩余", 10.14, 4.72, 2.10, 0.20, 10.5, "muted", True, PP_ALIGN.RIGHT)
    d.rect(s, 0.78, 5.38, 11.46, 0.48, "red_soft", "red_soft", 0)
    d.text(s, "缺口：Lead Data Engineer 2 FTE；可从设计与实验项目各调 1 FTE，不增加总编制。", 1.02, 5.51, 10.98, 0.20, 11, "red", True, PP_ALIGN.CENTER)
    d.takeaway(s, "管理动作：重配 2 FTE 到数据工程，未来四周冻结非关键设计增强。")


def control_tower(d: Deck) -> None:
    s = d.base("PMO 控制塔显示 12 个项目中 3 个需要升级，但只有 2 项需要管委会本周拍板", 5, "pmo control tower")
    states = [("正常", 6, "green"), ("关注", 3, "amber"), ("阻断", 2, "red"), ("完成", 1, "gray")]
    for i, (label, value, status) in enumerate(states):
        x = 0.76 + i * 1.62
        d.rect(s, x, 1.44, 1.42, 0.92, "white", "line")
        d.text(s, str(value), x, 1.62, 1.42, 0.42, 20, status, True, PP_ALIGN.CENTER)
        d.text(s, label, x, 2.04, 1.42, 0.18, 9.5, "muted", True, PP_ALIGN.CENTER)
    d.rect(s, 7.48, 1.44, 4.75, 0.92, "paper", "line")
    d.text(s, "组合 KPI", 7.76, 1.66, 1.0, 0.22, 11, "ink", True)
    d.text(s, "准时交付 76%", 8.96, 1.66, 1.25, 0.22, 11, "amber", True)
    d.text(s, "预算偏差 +3.4%", 10.42, 1.66, 1.55, 0.22, 11, "amber", True)
    rows = [["治理维度", "状态", "关键事实", "Owner"], ["里程碑", "Amber", "4 个节点两周内到期", "PMO"], ["风险与合规", "Red", "数据留存策略待批", "CISO"], ["预算", "Green", "EAC 低于批准预算", "CFO"], ["Stakeholder", "Amber", "2 位 Sponsor 决策待办", "COO"]]
    shape = d.table(s, rows, 0.76, 2.75, 7.25, 2.70, [1.65, 1.10, 3.40, 1.10], {1, 3}, 9.7)
    shape.table.cell(2, 1).fill.fore_color.rgb = rgb("red_soft")
    d.rect(s, 8.36, 2.75, 3.87, 2.70, "paper", "line")
    d.text(s, "本周决策队列", 8.66, 3.02, 3.26, 0.24, 12, "ink", True)
    decisions = [("D1", "批准 180 天数据留存", "CISO", "red"), ("D2", "重配 1.4M 到数据中台", "CFO", "amber")]
    for i, (code, text, owner, status) in enumerate(decisions):
        y = 3.52 + i * 0.78
        d.pill(s, 8.66, y, code, status, 0.48)
        d.text(s, text, 9.28, y + 0.03, 2.02, 0.30, 10.5, "ink", True)
        d.text(s, owner, 11.34, y + 0.03, 0.58, 0.18, 8.5, status, True, PP_ALIGN.RIGHT)
    d.takeaway(s, "管理动作：控制塔只升级例外和决策，不把 12 个项目的全部细节搬进管委会。")


def one_pager(d: Deck) -> None:
    s = d.base("本周完成核心构建，但数据与采用率仍阻碍 Scale Gate；下周只处理三项例外", 6, "executive status one-pager")
    d.rect(s, 0.76, 1.42, 7.35, 1.00, "paper", "line")
    d.text(s, "EXECUTIVE SUMMARY", 1.00, 1.66, 1.80, 0.20, 9.5, "blue", True)
    d.text(s, "工程闭环已运行；项目按预算推进，但数据覆盖与采用率未达到规模化门禁。", 2.88, 1.62, 4.92, 0.40, 12, "ink", True)
    d.rect(s, 8.44, 1.42, 3.80, 1.00, "navy", "navy", 0)
    d.text(s, "OVERALL STATUS", 8.72, 1.66, 1.50, 0.20, 9.5, "cyan", True)
    d.text(s, "AMBER", 10.36, 1.58, 1.48, 0.30, 18, "white", True, PP_ALIGN.RIGHT)
    panels = [("上周完成", "核心平台 Build 完成\n规则引擎进入 UAT", "green"), ("本周目标", "确认主数据 owner\n启动主管试点", "blue"), ("Top Issues", "数据覆盖 71%\n关键路径晚 10 天", "red"), ("Top Risks", "采用率仅 64%\n收益证据缺口 8.6M", "amber")]
    for i, (title, body, status) in enumerate(panels):
        x = 0.76 + i * 3.03
        d.rect(s, x, 2.82, 2.72, 1.40, "white", "line")
        d.rect(s, x, 2.82, 0.06, 1.40, status, status, 0)
        d.text(s, title, x + 0.22, 3.06, 2.25, 0.22, 11, "ink", True)
        d.text(s, body, x + 0.22, 3.48, 2.25, 0.48, 10.5, "text")
    d.text(s, "关键里程碑", 0.78, 4.66, 1.40, 0.22, 11.5, "ink", True)
    milestones = [("Scope lock", "07/20", "green"), ("Data gate", "07/30", "red"), ("Pilot review", "08/12", "amber"), ("Scale gate", "08/25", "gray")]
    d.rect(s, 2.30, 4.88, 9.60, 0.03, "line", "line", 0)
    for i, (name, date, status) in enumerate(milestones):
        x = 2.30 + i * 3.18
        d.rect(s, x, 4.76, 0.18, 0.28, status, status, 0)
        d.text(s, name, x - 0.40, 5.20, 1.00, 0.20, 9.3, "ink", True, PP_ALIGN.CENTER)
        d.text(s, date, x - 0.40, 5.52, 1.00, 0.18, 8.5, status, True, PP_ALIGN.CENTER)
    d.takeaway(s, "管理动作：确认数据 owner、资源重配和 4 周验证门禁；不批准新增范围。")


def main() -> None:
    manifest = json.loads(MANIFEST.read_text(encoding="utf-8"))
    expected = ["single_project_health_dashboard", "portfolio_dashboard", "time_resource_dashboard", "pmo_control_tower", "executive_status_one_pager"]
    actual = [page["page_family"] for page in manifest["pages"]]
    if not manifest["passed"] or actual != expected:
        raise SystemExit(f"selection manifest failed: {actual}")
    d = Deck(); cover(d); single_project(d); portfolio(d); resources(d); control_tower(d); one_pager(d)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True); d.prs.save(OUTPUT); print(OUTPUT)


if __name__ == "__main__":
    main()
