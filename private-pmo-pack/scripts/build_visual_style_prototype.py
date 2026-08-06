#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_STYLE = ROOT / "contracts" / "mckinsey-pmo.style.json"
DEFAULT_OUTPUT = ROOT / "output" / "pmo-reference-visual-styles.pptx"
DEFAULT_MANIFEST = ROOT / "qa" / "visual-style-prototype-manifest.json"


def load_json(path: Path) -> dict[str, Any]:
    data = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(data, dict):
        raise ValueError(f"{path} must contain a JSON object")
    return data


def write_json(path: Path, data: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def pct_x(start: float, width: float, value: float, min_value: float, max_value: float) -> float:
    return start + width * ((value - min_value) / (max_value - min_value))


class VisualDeck:
    def __init__(self, style: dict[str, Any]) -> None:
        self.style = style
        self.prs = Presentation()
        slide_spec = style["slide"]
        self.prs.slide_width = Inches(slide_spec["width_in"])
        self.prs.slide_height = Inches(slide_spec["height_in"])
        self.blank = self.prs.slide_layouts[6]
        self.records: list[dict[str, Any]] = []

    def color(self, name: str) -> RGBColor:
        return rgb(self.style["colors"][name])

    def status_color(self, status: str) -> RGBColor:
        return rgb(self.style["status_colors"].get(status, self.style["status_colors"]["gray"]))

    def add_slide(self, title: str, page_type: str, page_no: int, kicker: str) -> Any:
        slide = self.prs.slides.add_slide(self.blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.color("background")
        self.text(slide, kicker.upper(), 0.56, 0.25, 4.5, 0.2, size=8.5, color_name="blue", bold=True)
        self.text(slide, title, 0.56, 0.5, 12.0, 0.5, size=self.style["typography"]["title_pt"], color_name="ink", bold=True)
        self.line(slide, 0.56, 1.12, 12.78, 1.12, "line", 0.7)
        self.text(slide, f"VISUAL STYLE PROTOTYPE / {page_no:02d}", 9.42, 7.12, 3.3, 0.18, size=7.8, color_name="muted", bold=True, align=PP_ALIGN.RIGHT)
        self.records.append({"index": page_no, "page_type": page_type, "title": title})
        return slide

    def text(
        self,
        slide: Any,
        value: str,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        size: float | None = None,
        color_name: str = "text",
        bold: bool = False,
        align: PP_ALIGN = PP_ALIGN.LEFT,
        valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    ) -> Any:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        box.name = "Text: " + str(value)[:48]
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.03)
        frame.margin_right = Inches(0.03)
        frame.margin_top = Inches(0.01)
        frame.margin_bottom = Inches(0.01)
        frame.vertical_anchor = valign
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = str(value)
        run.font.name = self.style["typography"]["font_cjk"]
        run.font.size = Pt(size or self.style["typography"]["body_pt"])
        run.font.bold = bold
        run.font.color.rgb = self.color(color_name)
        return box

    def rect(
        self,
        slide: Any,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        fill: str = "white",
        line: str = "line",
        radius: bool = False,
        width: float = 0.7,
        name: str = "Material: native block",
    ) -> Any:
        kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
        shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.name = name
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color(fill)
        shape.line.color.rgb = self.color(line)
        shape.line.width = Pt(width)
        return shape

    def oval(
        self,
        slide: Any,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        fill: str = "white",
        line: str = "line",
        width: float = 0.7,
        name: str = "Material: native oval",
    ) -> Any:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.name = name
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color(fill)
        shape.line.color.rgb = self.color(line)
        shape.line.width = Pt(width)
        return shape

    def line(self, slide: Any, x1: float, y1: float, x2: float, y2: float, color_name: str = "line", width: float = 1.0, name: str = "Connector: native line") -> Any:
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        connector.name = name
        connector.line.color.rgb = self.color(color_name)
        connector.line.width = Pt(width)
        return connector

    def rule(self, slide: Any, x: float, y: float, w: float, h: float, color_name: str = "line", name: str = "Material: native rule") -> Any:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.name = name
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color(color_name)
        shape.line.color.rgb = self.color(color_name)
        shape.line.width = Pt(0)
        return shape

    def arc_line(self, slide: Any, x: float, y: float, w: float, h: float, color_name: str, *, width_pt: float = 28, name: str = "Material: gauge arc") -> Any:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.name = name
        shape.adjustments[0] = 90
        shape.adjustments[1] = 270
        shape.fill.background()
        shape.line.color.rgb = self.color(color_name)
        shape.line.width = Pt(width_pt)
        return shape

    def label_chip(self, slide: Any, x: float, y: float, label: str, value: str, color_name: str) -> None:
        self.oval(slide, x, y + 0.05, 0.12, 0.12, fill=color_name, line=color_name, width=0, name="Material: legend dot")
        self.text(slide, label, x + 0.18, y, 1.2, 0.18, size=8.2, color_name="muted", bold=True)
        if value:
            self.text(slide, value, x + 1.0, y, 0.6, 0.18, size=8.2, color_name="ink", bold=True, align=PP_ALIGN.RIGHT)

    def add_doughnut_chart(
        self,
        slide: Any,
        x: float,
        y: float,
        w: float,
        h: float,
        labels: list[str],
        values: list[float],
        colors: list[str],
        *,
        name: str,
        show_data_labels: bool = False,
    ) -> Any:
        chart_data = ChartData()
        chart_data.categories = labels
        chart_data.add_series("Share", values)
        chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y), Inches(w), Inches(h), chart_data)
        chart_shape.name = name
        chart = chart_shape.chart
        chart.has_title = False
        chart.has_legend = False
        chart.plots[0].has_data_labels = show_data_labels
        if show_data_labels:
            chart.plots[0].data_labels.number_format = "0%"
            chart.plots[0].data_labels.font.size = Pt(8)
            chart.plots[0].data_labels.font.color.rgb = self.color("ink")
        series = chart.series[0]
        for idx, color_name in enumerate(colors):
            point = series.points[idx]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = self.color(color_name)
            point.format.line.color.rgb = self.color("white")
            point.format.line.width = Pt(0.8)
        return chart_shape

    def save(self, path: Path) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(path)


def render_half_gauge(deck: VisualDeck) -> None:
    slide = deck.add_slide("半圆仪表盘把状态、缺口和行动优先级压在同一视线内", "half_gauge_dashboard", 1, "Dashboard gauge")
    deck.arc_line(slide, 1.1, 1.95, 3.95, 3.2, "soft", width_pt=33, name="Material: gauge base arc")
    deck.arc_line(slide, 1.1, 1.95, 3.95, 3.2, "green", width_pt=25, name="Material: gauge status arc")
    deck.text(slide, "74", 2.5, 3.22, 1.05, 0.38, size=31, color_name="ink", bold=True, align=PP_ALIGN.CENTER)
    deck.text(slide, "Readiness index", 2.2, 3.72, 1.66, 0.18, size=8.8, color_name="muted", bold=True, align=PP_ALIGN.CENTER)
    deck.label_chip(slide, 1.25, 4.8, "Risk zone", "0-55", "red")
    deck.label_chip(slide, 3.12, 4.8, "Watch", "56-70", "yellow")
    deck.label_chip(slide, 4.78, 4.8, "Healthy", "71-100", "green")
    for idx, (label, value, status) in enumerate([("Governance", "82", "green"), ("Evidence", "67", "yellow"), ("Adoption", "58", "red")]):
        x = 6.55 + idx * 1.9
        deck.rect(slide, x, 1.74, 1.45, 2.55, fill="paper", line="line")
        deck.rect(slide, x + 0.35, 3.2 - value_to_height(int(value)), 0.75, value_to_height(int(value)), fill=status, line=status, width=0, name="Material: gauge mini bar")
        deck.text(slide, value, x + 0.25, 1.98, 0.95, 0.28, size=17, color_name="ink", bold=True, align=PP_ALIGN.CENTER)
        deck.text(slide, label, x + 0.14, 3.62, 1.17, 0.18, size=8.5, color_name="muted", bold=True, align=PP_ALIGN.CENTER)
    deck.rect(slide, 0.82, 5.62, 11.4, 0.56, fill="navy", line="navy")
    deck.text(slide, "Design rule: half ring for executive status; evidence stays in editable labels and micro-bars.", 1.08, 5.8, 10.88, 0.2, size=9.7, color_name="white", bold=True, align=PP_ALIGN.CENTER)


def value_to_height(value: int) -> float:
    return 0.35 + (value / 100) * 1.15


def render_basic_doughnut(deck: VisualDeck) -> None:
    slide = deck.add_slide("基础甜甜圈图适合表达预算或资源构成，中心只保留一个结论", "basic_doughnut", 2, "Composition chart")
    labels = ["Build", "Data", "Change", "Reserve"]
    values = [0.38, 0.24, 0.21, 0.17]
    colors = ["blue", "green", "yellow", "orange"]
    deck.add_doughnut_chart(slide, 1.0, 1.58, 4.7, 4.15, labels, values, colors, name="Native chart: editable budget doughnut", show_data_labels=True)
    deck.text(slide, "$4.8M", 2.78, 3.12, 1.15, 0.34, size=22, color_name="ink", bold=True, align=PP_ALIGN.CENTER)
    deck.text(slide, "FY budget", 2.78, 3.5, 1.15, 0.16, size=8.8, color_name="muted", bold=True, align=PP_ALIGN.CENTER)
    for idx, (label, value, color_name) in enumerate(zip(labels, ["38%", "24%", "21%", "17%"], colors)):
        y = 1.72 + idx * 0.62
        deck.label_chip(slide, 6.7, y, label, value, color_name)
    deck.rect(slide, 6.55, 4.82, 4.95, 0.74, fill="paper", line="line")
    deck.text(slide, "Editable-core: native chart data, center value, legend, and interpretation remain separate.", 6.78, 5.06, 4.5, 0.22, size=8.6, color_name="muted", align=PP_ALIGN.CENTER)


def render_callout_doughnut(deck: VisualDeck) -> None:
    slide = deck.add_slide("带外部标注的环形图适合让高管先看类别，再看原因", "callout_doughnut", 3, "Annotated composition")
    labels = ["Core platform", "AI workflow", "Data migration", "Enablement", "Contingency"]
    values = [35, 23, 18, 14, 10]
    colors = ["navy", "blue", "green", "yellow", "orange"]
    deck.add_doughnut_chart(slide, 3.8, 1.42, 4.25, 4.25, labels, values, colors, name="Native chart: annotated ring")
    deck.text(slide, "100%", 5.22, 3.08, 1.35, 0.36, size=20, color_name="ink", bold=True, align=PP_ALIGN.CENTER)
    deck.text(slide, "scope split", 5.18, 3.45, 1.45, 0.15, size=8.2, color_name="muted", bold=True, align=PP_ALIGN.CENTER)
    callouts = [
        ("Core platform", "35%", "navy", 1.28, 1.74, 4.35, 2.35),
        ("AI workflow", "23%", "blue", 8.18, 1.74, 7.12, 2.32),
        ("Data migration", "18%", "green", 8.58, 4.82, 7.12, 4.52),
        ("Enablement", "14%", "yellow", 1.22, 4.72, 4.55, 4.35),
        ("Contingency", "10%", "orange", 1.65, 3.16, 4.25, 3.35),
    ]
    for label, value, color_name, tx, ty, lx, ly in callouts:
        deck.oval(slide, tx, ty + 0.07, 0.13, 0.13, fill=color_name, line=color_name, width=0, name="Material: callout marker")
        deck.text(slide, label, tx + 0.22, ty, 1.62, 0.18, size=8.8, color_name="ink", bold=True)
        deck.text(slide, value, tx + 1.42, ty, 0.52, 0.18, size=8.8, color_name="muted", bold=True, align=PP_ALIGN.RIGHT)
        deck.line(slide, tx + 0.08, ty + 0.14, lx, ly, "line", 0.85, name="Connector: callout leader")


def render_budget_kpi_rings(deck: VisualDeck) -> None:
    slide = deck.add_slide("预算 KPI 圆环组把金额、消耗率和风险灯号分开展示", "budget_kpi_rings", 4, "KPI ring group")
    metrics = [
        ("Run-rate", "$1.8M", "42%", ["blue", "soft"], [42, 58], "green"),
        ("Committed", "$2.7M", "63%", ["green", "soft"], [63, 37], "yellow"),
        ("At risk", "$0.6M", "14%", ["red", "soft"], [14, 86], "red"),
    ]
    for idx, (label, amount, pct, colors, values, status) in enumerate(metrics):
        x = 1.0 + idx * 4.0
        deck.rect(slide, x, 1.56, 3.25, 4.38, fill="white", line="line")
        deck.add_doughnut_chart(slide, x + 0.54, 2.05, 2.12, 2.12, [label, "remaining"], values, colors, name=f"Native chart: {label} KPI ring")
        deck.text(slide, pct, x + 1.12, 2.86, 0.96, 0.26, size=16.5, color_name="ink", bold=True, align=PP_ALIGN.CENTER)
        deck.oval(slide, x + 0.22, 1.84, 0.14, 0.14, fill=status, line=status, width=0, name="Material: KPI status dot")
        deck.text(slide, label, x + 0.45, 1.8, 1.6, 0.18, size=9.4, color_name="muted", bold=True)
        deck.text(slide, amount, x + 0.38, 4.42, 2.45, 0.3, size=21, color_name="ink", bold=True, align=PP_ALIGN.CENTER)
        deck.text(slide, "native chart + center labels", x + 0.34, 4.88, 2.54, 0.18, size=7.2, color_name="muted", align=PP_ALIGN.CENTER)
    deck.text(slide, "Use three rings only when the measures are peers; otherwise promote the most decision-critical KPI.", 1.3, 6.28, 10.7, 0.18, size=9.5, color_name="muted", align=PP_ALIGN.CENTER)


def render_gantt(deck: VisualDeck) -> None:
    slide = deck.add_slide("项目时间轴甘特图用月份网格、阶段条和门禁线控制复杂度", "gantt_timeline", 5, "Roadmap gantt")
    months = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug"]
    left, top, month_w, row_h = 2.15, 1.62, 1.18, 0.54
    deck.rect(slide, 0.75, top, 10.98, 0.42, fill="navy", line="navy")
    deck.text(slide, "Workstream", 0.95, top + 0.13, 0.9, 0.12, size=8, color_name="white", bold=True)
    for idx, month in enumerate(months):
        x = left + idx * month_w
        deck.text(slide, month, x + 0.28, top + 0.1, 0.58, 0.18, size=7.7, color_name="white", bold=True, align=PP_ALIGN.CENTER)
        deck.rule(slide, x, top + 0.42, 0.006, 4.23, "line", name="Material: month grid")
    deck.rule(slide, left + len(months) * month_w, top + 0.42, 0.006, 4.23, "line", name="Material: month grid")
    rows = [
        ("Discovery", 0, 2.2, "blue"),
        ("Platform build", 1.1, 4.6, "green"),
        ("Data migration", 2.4, 3.9, "yellow"),
        ("Pilot", 4.4, 2.2, "orange"),
        ("Scale-up", 5.8, 2.0, "navy"),
    ]
    for idx, (label, start, width, color_name) in enumerate(rows):
        y = top + 0.72 + idx * row_h
        deck.text(slide, label, 0.92, y + 0.1, 1.04, 0.16, size=8.4, color_name="ink", bold=True)
        deck.rect(slide, 0.75, y, 10.98, row_h, fill="background" if idx % 2 == 0 else "paper", line="line", width=0.35, name="Material: gantt row")
        deck.rect(slide, left + start * month_w, y + 0.12, width * month_w, 0.25, fill=color_name, line=color_name, radius=True, width=0, name="Material: gantt phase bar")
    gates = [(2.8, "Gate 1"), (5.6, "Gate 2")]
    for pos, label in gates:
        x = left + pos * month_w
        deck.rule(slide, x, top + 0.48, 0.014, 5.3, "red", name="Material: decision gate")
        deck.text(slide, label, x - 0.38, 5.92, 0.76, 0.16, size=7.7, color_name="red", bold=True, align=PP_ALIGN.CENTER)
    for idx, (label, _, _, color_name) in enumerate(rows[:4]):
        deck.label_chip(slide, 2.4 + idx * 2.05, 6.35, label, "", color_name)


def render_milestone_line(deck: VisualDeck) -> None:
    slide = deck.add_slide("项目里程碑折线图把进度曲线和阶段标签分层表达", "milestone_line", 6, "Milestone trend")
    left, top, width, height = 1.12, 1.62, 9.85, 3.45
    chart_data = ChartData()
    chart_data.categories = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul"]
    chart_data.add_series("Readiness", [18, 28, 43, 54, 69, 78, 88])
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(left), Inches(top), Inches(width), Inches(height), chart_data)
    chart_shape.name = "Native chart: milestone readiness line"
    chart = chart_shape.chart
    chart.has_legend = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 100
    chart.value_axis.major_unit = 25
    chart.value_axis.tick_labels.font.size = Pt(7.5)
    chart.category_axis.tick_labels.font.size = Pt(7.5)
    series = chart.series[0]
    series.format.line.color.rgb = deck.color("blue")
    series.format.line.width = Pt(2)
    for idx, label in [(1, "M1"), (3, "M3"), (5, "M5"), (6, "M6")]:
        x = left + idx * (width / 6)
        y = top + height - height * [18, 28, 43, 54, 69, 78, 88][idx] / 100
        deck.text(slide, label, x - 0.18, y - 0.34, 0.42, 0.18, size=7.2, color_name="ink", bold=True, align=PP_ALIGN.CENTER)
    phases = [("Design", 0, 2, "blue"), ("Build", 2, 4, "green"), ("Pilot", 4, 5.6, "yellow"), ("Scale", 5.6, 6.0, "navy")]
    for label, start, end, color_name in phases:
        x = left + start * (width / 6)
        w = (end - start) * (width / 6)
        deck.rect(slide, x, 5.98, w, 0.3, fill=color_name, line=color_name, width=0, name="Material: phase label band")
        deck.text(slide, label, x + 0.03, 6.06, max(w - 0.06, 0.24), 0.16, size=6.8, color_name="white", bold=True, align=PP_ALIGN.CENTER)
    deck.rect(slide, 11.05, 1.85, 1.35, 0.98, fill="paper", line="line")
    deck.text(slide, "88%", 11.25, 2.12, 0.92, 0.26, size=18, color_name="ink", bold=True, align=PP_ALIGN.CENTER)
    deck.text(slide, "target readiness", 11.18, 2.48, 1.05, 0.14, size=7.5, color_name="muted", align=PP_ALIGN.CENTER)


def count_package_objects(pptx: Path) -> dict[str, int]:
    counts = {"slides": 0, "pictures": 0, "shapes": 0, "graphic_frames": 0, "connectors": 0, "charts": 0}
    with zipfile.ZipFile(pptx) as package:
        slide_names = sorted(name for name in package.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml"))
        counts["slides"] = len(slide_names)
        for name in slide_names:
            xml = package.read(name).decode("utf-8", errors="ignore")
            counts["pictures"] += xml.count("<p:pic>")
            counts["shapes"] += xml.count("<p:sp>")
            counts["graphic_frames"] += xml.count("<p:graphicFrame>")
            counts["connectors"] += xml.count("<p:cxnSp>")
        counts["charts"] = len([name for name in package.namelist() if name.startswith("ppt/charts/chart") and name.endswith(".xml")])
    return counts


def build(style: dict[str, Any], output: Path, manifest_path: Path) -> dict[str, Any]:
    deck = VisualDeck(style)
    render_half_gauge(deck)
    render_basic_doughnut(deck)
    render_callout_doughnut(deck)
    render_budget_kpi_rings(deck)
    render_gantt(deck)
    render_milestone_line(deck)
    deck.save(output)
    counts = count_package_objects(output)
    manifest = {
        "schema_version": "1.0",
        "pack_id": "pptsmith-pmo-consulting-pack",
        "prototype_id": "pmo-reference-visual-styles",
        "built_at": datetime.now(timezone.utc).isoformat(),
        "status": "created",
        "builder": "private-python-pptx-visual-style-renderer",
        "outputs": {"deck": str(output.relative_to(ROOT)) if output.is_relative_to(ROOT) else str(output)},
        "style": str(DEFAULT_STYLE.relative_to(ROOT)),
        "slide_count": len(deck.prs.slides),
        "slides": deck.records,
        "native_object_counts": counts,
        "editable_core_expected": [
            "native doughnut charts for ring and composition data",
            "editable center values and callout labels",
            "native leader lines and markers",
            "native Gantt grid, phase bars, gate lines, and legend",
            "native milestone axis, trend segments, markers, and phase bands"
        ],
        "known_limitations": [
            "Half-gauge uses editable native ARC shapes rather than a data-bound PowerPoint chart.",
            "Doughnut arc thickness and start angle remain governed by PowerPoint chart behavior; fine-grained segment geometry is a v0.2 component concern.",
            "Milestone trend uses a native line chart plus separate native phase bands; phase bands are not data-bound to the chart series."
        ]
    }
    write_json(manifest_path, manifest)
    return manifest


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build six editable PMO visual style prototype slides.")
    parser.add_argument("--style", type=Path, default=DEFAULT_STYLE)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--manifest", type=Path, default=DEFAULT_MANIFEST)
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    manifest = build(load_json(args.style), args.output, args.manifest)
    print(json.dumps({"deck": str(args.output), "manifest": str(args.manifest), "slide_count": manifest["slide_count"]}, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
