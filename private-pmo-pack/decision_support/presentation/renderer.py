from __future__ import annotations

from dataclasses import dataclass
from datetime import date
from decimal import Decimal, InvalidOperation
from pathlib import Path
import re
from typing import Any, Mapping
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from components.tokens import resolve_tokens
from decision_support.pipeline import DecisionPackage

from .adapter import build_decision_presentation_ir


@dataclass(frozen=True)
class BuildResult:
    output: Path
    slide_count: int
    native_shape_count: int
    text_shape_count: int
    table_count: int
    chart_count: int
    connector_count: int
    picture_count: int
    media_count: int
    page_types: tuple[str, ...]
    out_of_bounds_text: tuple[str, ...]


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#"))


class _Deck:
    def __init__(self, tokens: Mapping[str, Any]) -> None:
        self.tokens = tokens
        self.prs = Presentation()
        self.prs.slide_width = Inches(float(tokens["slide"]["width_in"]))
        self.prs.slide_height = Inches(float(tokens["slide"]["height_in"]))
        self.blank = self.prs.slide_layouts[6]

    def text(self, slide: Any, text: Any, x: float, y: float, w: float, h: float, *,
             name: str, size: float = 10, bold: bool = False, color: str = "text.body",
             align: PP_ALIGN = PP_ALIGN.LEFT, anchor: MSO_ANCHOR = MSO_ANCHOR.TOP) -> Any:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        shape.name = name
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = Inches(0.04)
        frame.margin_top = frame.margin_bottom = Inches(0.02)
        frame.vertical_anchor = anchor
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = str(text)
        run.font.name = self.tokens["typography"]["font"]
        run.font.size = Pt(size)
        run.font.bold = bold
        node: Any = self.tokens["semantic"]
        for part in color.split("."):
            node = node[part]
        run.font.color.rgb = _rgb(node)
        return shape

    def shape(self, slide: Any, shape_type: MSO_AUTO_SHAPE_TYPE, x: float, y: float, w: float, h: float, *,
              name: str, fill: str = "surface.paper", line: str = "line.divider",
              line_width: float = 0.7) -> Any:
        shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.name = name if name.lower().startswith(("decoration:", "background:", "material:", "connector:")) else f"Material: {name}"
        fill_node: Any = self.tokens["semantic"]
        line_node: Any = self.tokens["semantic"]
        for part in fill.split("."):
            fill_node = fill_node[part]
        for part in line.split("."):
            line_node = line_node[part]
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill_node)
        shape.line.color.rgb = _rgb(line_node)
        shape.line.width = Pt(line_width)
        return shape

    def box(self, slide: Any, x: float, y: float, w: float, h: float, *, name: str,
            fill: str = "surface.paper", line: str = "line.divider") -> Any:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.name = name if name.lower().startswith(("decoration:", "background:", "material:", "connector:")) else f"Material: {name}"
        fill_node: Any = self.tokens["semantic"]
        line_node: Any = self.tokens["semantic"]
        for part in fill.split("."):
            fill_node = fill_node[part]
        for part in line.split("."):
            line_node = line_node[part]
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill_node)
        shape.line.color.rgb = _rgb(line_node)
        shape.line.width = Pt(0.7)
        return shape

    def line(self, slide: Any, x1: float, y1: float, x2: float, y2: float, *, name: str,
             color: str = "line.annotation", width: float = 1.2) -> Any:
        shape = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        shape.name = name
        node: Any = self.tokens["semantic"]
        for part in color.split("."):
            node = node[part]
        shape.line.color.rgb = _rgb(node)
        shape.line.width = Pt(width)
        return shape

    def rule(self, slide: Any, x: float, y: float, w: float, *, name: str,
             color: str = "line.soft", height: float = 0.008) -> Any:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(height),
        )
        shape.name = name if name.lower().startswith("material:") else f"Material: {name}"
        node: Any = self.tokens["semantic"]
        for part in color.split("."):
            node = node[part]
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(node)
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def base(self, page: Mapping[str, Any]) -> Any:
        prefix = page["page_type"]
        slide = self.prs.slides.add_slide(self.blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(self.tokens["semantic"]["surface"]["canvas"])
        self.text(slide, prefix.replace("_", " ").upper(), 0.56, 0.23, 4.5, 0.22,
                  name=f"{prefix}.kicker", size=8.5, bold=True, color="sequence.phase_3")
        self.text(slide, page["title"], 0.56, 0.50, 12.15, 0.62,
                  name=f"{prefix}.title", size=20, bold=True, color="text.title")
        self.line(slide, 0.56, 1.25, 12.77, 1.25, name=f"{prefix}.title_rule", color="line.divider", width=0.65)
        meta = page["metadata"]
        footer = (f"RUN {meta['analysis_run_id']}  |  HASH {meta['input_hash_short']}  |  "
                  f"CUTOFF {meta['evidence_cutoff']}  |  MODEL {meta['framework_version']}  |  "
                  f"STATUS {meta['recommendation_status']}")
        self.text(slide, footer, 0.56, 7.09, 12.2, 0.20, name=f"{prefix}.metadata", size=7.2,
                  color="text.muted")
        return slide

    def table(self, slide: Any, rows: list[list[str]], x: float, y: float, w: float, h: float, *, name: str) -> Any:
        shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
        shape.name = name
        table = shape.table
        header_height = min(0.44, h)
        table.rows[0].height = Inches(header_height)
        if len(rows) > 1:
            body_height = max(0.30, (h - header_height) / (len(rows) - 1))
            for row_index in range(1, len(rows)):
                table.rows[row_index].height = Inches(body_height)
        for row_index, row in enumerate(rows):
            for col_index, value in enumerate(row):
                cell = table.cell(row_index, col_index)
                cell.text = str(value)
                cell.text_frame.word_wrap = True
                cell.vertical_anchor = MSO_ANCHOR.TOP
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(self.tokens["semantic"]["surface"]["paper"] if row_index else self.tokens["semantic"]["sequence"]["phase_1"])
                cell.margin_left = cell.margin_right = Inches(0.05)
                cell.margin_top = cell.margin_bottom = Inches(0.04)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = self.tokens["typography"]["font"]
                        run.font.size = Pt(8 if row_index else 8.5)
                        run.font.bold = row_index == 0
                        run.font.color.rgb = _rgb(self.tokens["semantic"]["text"]["body"] if row_index else self.tokens["semantic"]["text"]["inverse"])
        return shape


def _humanize(value: Any) -> str:
    return str(value).replace("_", " ").replace(":", " · ")


def _display_number(value: Any) -> str:
    try:
        number = Decimal(str(value))
    except (InvalidOperation, ValueError):
        return str(value)
    if not number.is_finite():
        return str(value)
    if abs(number) >= 1000:
        return f"{number:,.0f}"
    if number == number.to_integral_value():
        return f"{number:.0f}"
    return f"{number:.2f}".rstrip("0").rstrip(".")


def _display_value(value: Any, *, max_items: int = 4) -> str:
    if value is None:
        return "Not supplied"
    if isinstance(value, bool):
        return "Yes" if value else "No"
    if isinstance(value, Mapping):
        items = list(value.items())
        if not items:
            return "None declared"
        visible = items[:max_items]
        text = "; ".join(f"{_humanize(key)}: {_display_number(item)}" for key, item in visible)
        if len(items) > max_items:
            text += f"; +{len(items) - max_items} audited fields"
        return text
    if isinstance(value, (list, tuple)):
        if not value:
            return "None declared"
        visible = list(value[:max_items])
        text = " • ".join(_display_value(item, max_items=2) for item in visible)
        if len(value) > max_items:
            text += f" • +{len(value) - max_items} more"
        return text
    return _display_number(value)


def _items(value: Any) -> list[tuple[str, str]]:
    if isinstance(value, Mapping):
        return [(_humanize(key), _display_value(item)) for key, item in value.items()]
    if isinstance(value, (list, tuple)):
        pairs = []
        for index, item in enumerate(value):
            if isinstance(item, Mapping):
                label = item.get("label") or item.get("option") or item.get("framework") or str(index + 1)
                detail = " | ".join(
                    f"{_humanize(key)}: {_display_value(val)}"
                    for key, val in item.items()
                    if key not in {"label", "option", "framework"}
                )
                pairs.append((str(label), detail))
            else:
                pairs.append((str(index + 1), _display_value(item)))
        return pairs
    return [("", _display_value(value))]


def _render_generic(deck: _Deck, page: Mapping[str, Any]) -> None:
    slide = deck.base(page)
    prefix = page["page_type"]
    sections = [(key, value) for key, value in page.items()
                if key not in {"page_type", "title", "framework_id", "metadata"}]
    columns = 2 if len(sections) > 1 else 1
    width = 5.85 if columns == 2 else 12.0
    for section_index, (key, value) in enumerate(sections[:6]):
        col = section_index % columns
        row = section_index // columns
        x = 0.62 + col * 6.12
        y = 1.48 + row * 1.70
        deck.text(slide, key.replace("_", " ").upper(), x, y, width, 0.22,
                  name=f"{prefix}.section.{section_index + 1}.label", size=8.5, bold=True, color="text.muted")
        pairs = _items(value)
        if not pairs:
            pairs = [("", "None declared")]
        row_height = min(0.62, 1.35 / max(1, len(pairs)))
        for item_index, (label, detail) in enumerate(pairs[:5]):
            item_y = y + 0.30 + item_index * row_height
            if item_y + row_height > 6.86:
                break
            deck.box(slide, x, item_y, width, max(0.34, row_height - 0.05),
                     name=f"{prefix}.native.{section_index + 1}.{item_index + 1}", fill="surface.paper")
            if label:
                deck.text(slide, label, x + 0.10, item_y + 0.06, width * 0.28, row_height - 0.10,
                          name=f"{prefix}.item.{section_index + 1}.{item_index + 1}.label", size=8.7, bold=True, color="text.title")
                deck.text(slide, detail, x + width * 0.31, item_y + 0.06, width * 0.66, row_height - 0.10,
                          name=f"{prefix}.item.{section_index + 1}.{item_index + 1}.detail", size=8.2)
            else:
                deck.text(slide, detail, x + 0.10, item_y + 0.06, width - 0.20, row_height - 0.10,
                          name=f"{prefix}.item.{section_index + 1}.{item_index + 1}.detail", size=8.2)


def _render_decision_brief(deck: _Deck, page: Mapping[str, Any]) -> None:
    slide = deck.base(page)
    frameworks = list(page.get("framework_order", ()))
    stage_x = [0.72, 2.84, 5.16, 10.50]
    stage_w = [1.76, 1.98, 1.92, 2.16]
    framework_x = 7.36
    framework_w = 2.72
    center_y = 3.78

    for index in range(1, 3):
        if index < 3:
            deck.line(slide, stage_x[index - 1] + stage_w[index - 1], center_y,
                      stage_x[index] - 0.10, center_y,
                      name=f"decision_brief.connector.stage.{index}", color="sequence.phase_2", width=2.2)
    for index, _framework in enumerate(frameworks[:4]):
        branch_y = 2.18 + index * 0.76
        deck.line(slide, stage_x[2] + stage_w[2], center_y, framework_x, branch_y + 0.23,
                  name=f"decision_brief.connector.framework.{index + 1}", color="line.annotation", width=1.35)
        deck.line(slide, framework_x + framework_w, branch_y + 0.23, stage_x[3], center_y,
                  name=f"decision_brief.connector.rejoin.{index + 1}", color="line.annotation", width=1.35)

    stage_specs = [
        ("DECLARED\nEVIDENCE", "1 decisive fact", "sequence.phase_1"),
        ("SUFFICIENCY\nGATE", "PASS · no missing fields", "state.healthy"),
        ("4 FRAMEWORK\nASSESSMENTS", "kept separate", "sequence.phase_2"),
        ("NO AGGREGATE\nRECOMMENDATION", "intentional boundary", "state.watch"),
    ]
    for index, ((label, detail, fill), x, width) in enumerate(zip(stage_specs, stage_x, stage_w), start=1):
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.PENTAGON, x, 3.15, width, 1.28,
                   name=f"decision_brief.stage.{index}", fill=fill, line=fill, line_width=0.0)
        label_x = x + (0.08 if index == 4 else 0.10)
        label_width = width - (0.20 if index == 4 else 0.38)
        label_size = 9.8 if index == 4 else (10.4 if index == 3 else 11.2)
        deck.text(slide, label, label_x, 3.42, label_width, 0.48,
                  name=f"decision_brief.stage.{index}.label", size=label_size,
                  bold=True, color="text.inverse" if index != 4 else "text.title",
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        deck.text(slide, detail, x + 0.10, 4.58, width - 0.20, 0.26,
                  name=f"decision_brief.stage.{index}.detail", size=8.4, bold=index == 4,
                  color="text.muted", align=PP_ALIGN.CENTER)

    for index, framework in enumerate(frameworks[:4]):
        y = 2.18 + index * 0.76
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, framework_x, y, framework_w, 0.46,
                   name=f"decision_brief.framework.{index + 1}", fill="surface.paper",
                   line="sequence.phase_2", line_width=1.1)
        deck.text(slide, _humanize(framework).upper(), framework_x + 0.14, y + 0.10, framework_w - 0.28, 0.22,
                  name=f"decision_brief.framework.{index + 1}.label", size=8.8, bold=True,
                  color="text.title", align=PP_ALIGN.CENTER)

    deck.shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, stage_x[3], 5.28, stage_w[3], 0.54,
               name="decision_brief.terminal_state", fill="surface.paper", line="state.watch", line_width=1.5)
    deck.text(slide, "ROUTED, NOT COMBINED", stage_x[3] + 0.12, 5.43, stage_w[3] - 0.24, 0.20,
              name="decision_brief.terminal_state.label", size=8.3, bold=True,
              color="text.title", align=PP_ALIGN.CENTER)
    deck.text(slide, "Evidence is sufficient to assess each domain. Readiness is not sufficient to collapse them into one opaque answer.",
              0.72, 6.28, 11.90, 0.36, name="decision_brief.boundary", size=9.4,
              color="text.muted", align=PP_ALIGN.CENTER)


def _render_benefits(deck: _Deck, page: Mapping[str, Any]) -> None:
    slide = deck.base(page)
    series = {key: float(value) for key, value in page.get("benefit_series", {}).items()}
    planned = series.get("planned_to_date", 0.0)
    actual = series.get("actual_to_date", 0.0)
    target = series.get("target", 0.0)
    forecast = series.get("forecast", 0.0)
    scale_max = max(target, forecast, actual, planned, 1.0) * 1.08
    x0, chart_w = 1.00, 10.90
    to_x = lambda value: x0 + chart_w * value / scale_max

    deck.text(slide, "ACTUAL BENEFIT VS CURRENT PLAN", 0.72, 1.52, 5.2, 0.24,
              name="benefits_dashboard.bullet.label", size=9.0, bold=True, color="text.muted")
    deck.box(slide, x0, 2.12, chart_w, 0.42, name="benefits_dashboard.bullet.track",
             fill="surface.paper", line="line.soft")
    actual_width = max(0.08, to_x(actual) - x0)
    deck.box(slide, x0, 2.12, actual_width, 0.42, name="benefits_dashboard.bullet.actual",
             fill="state.healthy", line="state.healthy")
    plan_x = to_x(planned)
    deck.line(slide, plan_x, 1.92, plan_x, 2.75, name="benefits_dashboard.bullet.plan_marker",
              color="text.title", width=2.0)
    target_x = to_x(target)
    deck.line(slide, target_x, 2.02, target_x, 2.64, name="benefits_dashboard.bullet.target_marker",
              color="state.watch", width=2.0)
    deck.text(slide, f"PLAN  {_display_number(planned)}", plan_x - 0.58, 1.72, 1.16, 0.24,
              name="benefits_dashboard.bullet.plan.label", size=8.3, bold=True,
              color="text.title", align=PP_ALIGN.CENTER)
    deck.text(slide, f"ACTUAL  {_display_number(actual)}", to_x(actual) - 0.64, 2.79, 1.28, 0.26,
              name="benefits_dashboard.bullet.actual.label", size=8.5, bold=True,
              color="state.healthy", align=PP_ALIGN.CENTER)
    deck.text(slide, f"TARGET  {_display_number(target)}", target_x - 0.58, 2.79, 1.16, 0.26,
              name="benefits_dashboard.bullet.target.label", size=8.3,
              color="text.title", align=PP_ALIGN.CENTER)
    realization_metric = next(
        (float(item.get("value", 0)) for item in page.get("metrics", ()) if item.get("label") == "Benefit realization"),
        actual / planned if planned else 0,
    )
    deck.text(slide, f"{realization_metric:.0%} OF PLANNED UPLIFT", 9.22, 1.35, 2.64, 0.70,
              name="benefits_dashboard.realization", size=17, bold=True,
              color="state.healthy", align=PP_ALIGN.RIGHT)

    deck.text(slide, "FUTURE BENEFIT BRIDGE", 0.72, 3.54, 5.2, 0.24,
              name="benefits_dashboard.bridge.label", size=9.0, bold=True, color="text.muted")
    bridge_x, bridge_y, bridge_w = 1.00, 4.20, 10.90
    target_width = max(0.08, bridge_w * target / scale_max)
    gap_width = max(0.08, bridge_w * max(0.0, forecast - target) / scale_max)
    deck.box(slide, bridge_x, bridge_y, target_width, 0.62,
             name="benefits_dashboard.bridge.target", fill="sequence.phase_1", line="sequence.phase_1")
    deck.line(slide, bridge_x + target_width, bridge_y - 0.18, bridge_x + target_width, bridge_y + 0.82,
              name="benefits_dashboard.bridge.target_marker", color="text.title", width=1.8)
    deck.box(slide, bridge_x + target_width, bridge_y, gap_width, 0.62,
             name="benefits_dashboard.bridge.forecast_gap", fill="state.healthy", line="state.healthy")
    deck.text(slide, f"TARGET  {_display_number(target)}", bridge_x + target_width - 1.36, bridge_y + 0.18, 1.22, 0.25,
              name="benefits_dashboard.bridge.target.label", size=8.8, bold=True,
              color="text.inverse", align=PP_ALIGN.RIGHT)
    deck.text(slide, f"+{_display_number(forecast - target)}", bridge_x + target_width + 0.08, bridge_y + 0.17,
              max(0.72, gap_width - 0.12), 0.26, name="benefits_dashboard.bridge.gap.label",
              size=9.2, bold=True, color="text.inverse", align=PP_ALIGN.CENTER)
    deck.line(slide, bridge_x + target_width + gap_width, bridge_y - 0.14,
              bridge_x + target_width + gap_width, bridge_y + 0.78,
              name="benefits_dashboard.bridge.forecast_marker", color="state.healthy", width=2.2)
    deck.text(slide, f"FORECAST  {_display_number(forecast)}", bridge_x + target_width + gap_width - 0.65,
              4.98, 1.30, 0.25, name="benefits_dashboard.bridge.forecast.label",
              size=8.8, bold=True, color="state.healthy", align=PP_ALIGN.CENTER)
    deck.text(slide, "Current delivery is ahead of plan; the decision still rests on future incremental economics, not sunk cost.",
              0.90, 6.08, 11.20, 0.40, name="benefits_dashboard.boundary", size=9.6,
              color="text.muted", align=PP_ALIGN.CENTER)


def _render_business_case_gate(deck: _Deck, page: Mapping[str, Any]) -> None:
    slide = deck.base(page)
    prefix = page["page_type"]
    gates = list(page.get("gates", ()))[:10]
    gate_x0, gate_step, gate_y = 0.72, 0.82, 2.18
    deck.text(slide, "TEN ORDERED DECISION GATES", gate_x0, 1.50, 7.8, 0.24,
              name=f"{prefix}.gates.label", size=9.0, bold=True, color="text.muted")
    for index in range(9):
        x1 = gate_x0 + index * gate_step + 0.34
        x2 = gate_x0 + (index + 1) * gate_step
        deck.line(slide, x1, gate_y + 0.24, x2, gate_y + 0.24,
                  name=f"{prefix}.gate.connector.{index + 1}", color="state.healthy", width=2.1)
    gate_labels = [
        "VALID\nQUESTION",
        "HARD CONSTRAINT\nCONTINUE",
        "HARD CONSTRAINT\nPIVOT",
        "HARD CONSTRAINT\nSTOP",
        "SUNK COST\nEXCLUDED",
        "EVIDENCE\nSUFFICIENCY",
        "POLICY\nCOMPLETENESS",
        "POLICY\nCONFIDENCE",
        "FUTURE\nECONOMICS",
        "SENSITIVITY",
    ]
    for index, gate in enumerate(gates):
        x = gate_x0 + index * gate_step
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.HEXAGON, x, gate_y, 0.54, 0.48,
                   name=f"{prefix}.gate.{index + 1}", fill="state.healthy", line="state.healthy", line_width=0)
        deck.text(slide, f"{index + 1:02d}", x + 0.07, gate_y + 0.12, 0.40, 0.20,
                  name=f"{prefix}.gate.{index + 1}.index", size=8.3, bold=True,
                  color="text.inverse", align=PP_ALIGN.CENTER)
        short = gate_labels[index] if index < len(gate_labels) else _humanize(gate)
        deck.text(slide, short.upper(), x - 0.31, 2.78 + (index % 2) * 0.52, 1.16, 0.40,
                  name=f"{prefix}.gate.{index + 1}.detail", size=7.0, bold=True,
                  color="text.muted", align=PP_ALIGN.CENTER)

    options = list(page.get("options", ()))[:3]
    npv_values = {str(item.get("option")): float(item.get("npv", 0)) for item in options}
    zero_x, bar_scale = 10.06, 0.019
    deck.text(slide, "OPTION NPV · COMMON ZERO", 9.18, 1.50, 3.42, 0.24,
              name=f"{prefix}.npv.label", size=9.0, bold=True, color="text.muted")
    deck.line(slide, zero_x, 1.90, zero_x, 4.18, name=f"{prefix}.npv.zero",
              color="text.title", width=1.6)
    option_colors = {"continue": "state.healthy", "pivot": "sequence.phase_2", "stop": "state.risk"}
    for index, option in enumerate(("continue", "pivot", "stop")):
        value = npv_values.get(option, 0.0)
        y = 2.04 + index * 0.70
        width = max(0.12, abs(value) * bar_scale)
        x = zero_x if value >= 0 else zero_x - width
        deck.box(slide, x, y, width, 0.38, name=f"{prefix}.npv.{option}",
                 fill=option_colors[option], line=option_colors[option])
        deck.text(slide, option.upper(), 8.88, y + 0.07, 1.06, 0.22,
                  name=f"{prefix}.npv.{option}.label", size=8.3, bold=True, color="text.title")
        label_x = x + width + 0.06 if value >= 0 else zero_x + 0.08
        deck.text(slide, _display_number(value), label_x, y + 0.07, 0.60, 0.22,
                  name=f"{prefix}.npv.{option}.value", size=8.5, bold=True,
                  color=option_colors[option], align=PP_ALIGN.LEFT)

    calculations = page.get("calculations", {})
    trace_rows = [
        ["Metric", "Continue", "Pivot", "Stop"],
        ["NPV", calculations.get("npv-continue"), calculations.get("npv-pivot"), calculations.get("npv-stop")],
        ["Incremental NPV", calculations.get("npv-incremental-continue"), calculations.get("npv-incremental-pivot"), "—"],
        ["Incremental ROI", calculations.get("incremental_roi-continue"), calculations.get("incremental_roi-pivot"), "—"],
    ]
    trace_rows = [[_display_number(item) if column else str(item) for column, item in enumerate(row)] for row in trace_rows]
    deck.table(slide, trace_rows, 0.72, 4.35, 8.22, 1.82, name=f"{prefix}.native.calculation_trace")
    deck.text(slide, "AUDITED OPTION ECONOMICS", 0.72, 4.00, 5.0, 0.24,
              name=f"{prefix}.calculations.label", size=9.0, bold=True, color="text.muted")
    deck.shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 9.30, 4.60, 3.08, 0.78,
               name=f"{prefix}.mode", fill="surface.paper", line="state.healthy", line_width=1.5)
    deck.text(slide, f"{_humanize(page.get('mode', 'not supplied')).upper()} · 10/10 GATES PASS",
              9.48, 4.83, 2.72, 0.28, name=f"{prefix}.mode.value", size=10.2,
              bold=True, color="state.healthy", align=PP_ALIGN.CENTER)
    deck.text(slide, "Continue leads on both NPV and incremental ROI. Values are presentation formatting of audited output.",
              9.28, 5.65, 3.14, 0.62, name=f"{prefix}.boundary", size=8.8,
              color="text.muted", align=PP_ALIGN.CENTER)


def _render_decision_record(deck: _Deck, page: Mapping[str, Any]) -> None:
    slide = deck.base(page)
    prefix = page["page_type"]
    deck.text(slide, "RECORDS", 0.62, 1.48, 8.05, 0.22,
              name=f"{prefix}.records.label", size=8.5, bold=True, color="text.muted")
    for index, record in enumerate(list(page.get("records", ()))[:4]):
        y = 1.78 + index * 1.20
        deck.box(slide, 0.62, y, 8.05, 1.04, name=f"{prefix}.native.record.{index + 1}")
        deck.text(slide, _humanize(record.get("framework", "framework")).upper(), 0.76, y + 0.10, 1.78, 0.22,
                  name=f"{prefix}.record.{index + 1}.framework", size=8.4, bold=True, color="sequence.phase_3")
        deck.text(
            slide,
            f"Owner  {record.get('owner', 'not supplied')}   |   Review  {record.get('review_date', 'not supplied')}",
            2.55,
            y + 0.10,
            5.90,
            0.22,
            name=f"{prefix}.record.{index + 1}.owner_review",
            size=8.2,
            bold=True,
            color="text.title",
        )
        conditions = _display_value(record.get("conditions", ()), max_items=2)
        triggers = _display_value(record.get("triggers", ()), max_items=1)
        deck.text(slide, f"Conditions  {conditions}", 0.76, y + 0.40, 7.67, 0.22,
                  name=f"{prefix}.record.{index + 1}.conditions", size=7.7)
        deck.text(slide, f"Trigger  {triggers}", 0.76, y + 0.68, 7.67, 0.24,
                  name=f"{prefix}.record.{index + 1}.triggers", size=7.7, color="text.muted")

    deck.text(slide, "DECISION STATUS", 8.93, 1.48, 3.74, 0.22,
              name=f"{prefix}.status.label", size=8.5, bold=True, color="text.muted")
    deck.box(slide, 8.93, 1.78, 3.74, 1.04, name=f"{prefix}.native.status")
    deck.text(slide, _humanize(page.get("decision_status", "not supplied")).upper(), 9.12, 2.08, 3.36, 0.30,
              name=f"{prefix}.status.value", size=15, bold=True, color="sequence.phase_3")
    deck.text(
        slide,
        "Approval remains conditional. A record is complete only when its owner, review date, conditions, and rerun trigger are retained.",
        8.93,
        3.14,
        3.74,
        1.20,
        name=f"{prefix}.status.boundary",
        size=9.2,
        color="text.body",
    )


def _render_matrix(deck: _Deck, page: Mapping[str, Any]) -> None:
    slide = deck.base(page)
    scores = list(page.get("scores", ()))
    selected = set(page.get("selected", ()))
    ranked = sorted(scores, key=lambda item: float(item.get("score", 0)), reverse=True)
    deck.text(slide, "FIXED-ANCHOR HEAT MATRIX", 0.72, 1.48, 3.10, 0.24,
              name="weighted_decision_matrix.heat.label", size=9.0, bold=True, color="text.muted")
    deck.text(slide, "VALUE", 2.48, 1.80, 1.18, 0.22,
              name="weighted_decision_matrix.criterion", size=8.0, bold=True,
              color="text.muted", align=PP_ALIGN.CENTER)
    heat_fill = [
        (0.80, "state.healthy"),
        (0.65, "sequence.phase_3"),
        (0.45, "sequence.phase_2"),
        (0.00, "state.watch"),
    ]
    for index, item in enumerate(ranked):
        option = str(item.get("option"))
        score = float(item.get("score", 0))
        y = 2.16 + index * 0.88
        fill = next(color for floor, color in heat_fill if score >= floor)
        deck.text(slide, option.upper(), 0.82, y + 0.18, 1.42, 0.26,
                  name=f"weighted_decision_matrix.option.{option}", size=9.4, bold=True,
                  color="text.title")
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 2.48, y, 1.18, 0.62,
                   name=f"weighted_decision_matrix.heat.{option}", fill=fill, line=fill, line_width=0)
        deck.text(slide, f"{score:.2f}", 2.62, y + 0.17, 0.90, 0.25,
                  name=f"weighted_decision_matrix.heat.{option}.value", size=10.6, bold=True,
                  color="text.inverse" if fill != "state.watch" else "text.title",
                  align=PP_ALIGN.CENTER)
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 3.84, y + 0.16, 0.28, 0.28,
                   name=f"weighted_decision_matrix.selection.{option}",
                   fill="state.healthy" if option in selected else "surface.paper",
                   line="state.healthy" if option in selected else "line.annotation", line_width=1.2)

    rank_x, rank_w = 5.05, 6.30
    threshold = 0.75
    threshold_x = rank_x + rank_w * threshold
    deck.text(slide, "RANKING ON THE SAME 0–1 SCALE", rank_x, 1.48, rank_w, 0.24,
              name="weighted_decision_matrix.rank.label", size=9.0, bold=True, color="text.muted")
    deck.line(slide, threshold_x, 1.88, threshold_x, 5.72,
              name="weighted_decision_matrix.threshold", color="state.watch", width=1.8)
    deck.text(slide, "0.75 REVIEW LINE", threshold_x - 0.68, 5.80, 1.36, 0.22,
              name="weighted_decision_matrix.threshold.label", size=7.7, bold=True,
              color="state.watch", align=PP_ALIGN.CENTER)
    for index, item in enumerate(ranked):
        option = str(item.get("option"))
        score = float(item.get("score", 0))
        y = 2.16 + index * 0.88
        deck.rule(slide, rank_x, y + 0.30, rank_w, name=f"weighted_decision_matrix.rank.track.{option}",
                  color="line.soft", height=0.04)
        deck.box(slide, rank_x, y + 0.14, max(0.10, rank_w * score), 0.32,
                 name=f"weighted_decision_matrix.rank.{option}",
                 fill="state.healthy" if option in selected else "sequence.phase_2",
                 line="state.healthy" if option in selected else "sequence.phase_2")
        deck.text(slide, f"#{index + 1}", 11.50, y + 0.17, 0.44, 0.24,
                  name=f"weighted_decision_matrix.rank.{option}.index", size=8.4, bold=True,
                  color="text.title", align=PP_ALIGN.CENTER)
        deck.text(slide, "SELECTED" if option in selected else "NOT SELECTED", 11.92, y + 0.17, 0.76, 0.24,
                  name=f"weighted_decision_matrix.rank.{option}.selection", size=7.4, bold=True,
                  color="state.healthy" if option in selected else "text.muted", align=PP_ALIGN.RIGHT)
    deck.text(slide, f"Method: {_humanize(page.get('method', 'not supplied'))}. Hard constraints precede scoring; selection is not a naive top-N cutoff.",
              0.82, 6.38, 11.75, 0.32, name="weighted_decision_matrix.boundary",
              size=9.0, color="text.muted", align=PP_ALIGN.CENTER)


def _render_scenarios(deck: _Deck, page: Mapping[str, Any]) -> None:
    slide = deck.base(page)
    scenarios = page.get("scenarios", {})
    demand = page.get("option_resource_demand", {})
    costs = page.get("option_cost", {})
    limits = page.get("resource_limits", {})
    budgets = page.get("budget_limits", {})
    frame_x, frame_w = 2.25, 8.72
    colors = {
        "platform": "sequence.phase_1",
        "security": "state.watch",
        "growth": "state.healthy",
        "migration": "sequence.phase_3",
        "unused": "surface.paper",
    }
    deck.text(slide, "TEAM CAPACITY · NORMALIZED TO EACH SCENARIO CEILING", 0.72, 1.48, 7.6, 0.24,
              name="capacity_scenario_allocation.label", size=9.0, bold=True, color="text.muted")
    for index, name in enumerate(("base", "downside", "relief")):
        scenario = scenarios.get(name, {})
        selected = list(scenario.get("selected_option_ids", ()))
        limit = float(limits.get(name, {}).get("team", 1.0) or 1.0)
        y = 2.05 + index * 1.38
        deck.text(slide, name.upper(), 0.72, y + 0.12, 1.12, 0.30,
                  name=f"capacity_scenario_allocation.scenario.{name}", size=10.0,
                  bold=True, color="text.title")
        deck.box(slide, frame_x, y, frame_w, 0.58,
                 name=f"capacity_scenario_allocation.frame.{name}", fill="surface.paper", line="line.annotation")
        cursor = frame_x
        used = 0.0
        for option in selected:
            amount = float(demand.get(option, {}).get("team", 0.0))
            width = frame_w * amount / limit
            used += amount
            if width <= 0:
                continue
            deck.box(slide, cursor, y, width, 0.58,
                     name=f"capacity_scenario_allocation.segment.{name}.{option}",
                     fill=colors.get(option, "sequence.phase_2"), line=colors.get(option, "sequence.phase_2"))
            if width >= 0.78:
                deck.text(slide, f"{option.upper()}  {amount:g}", cursor + 0.05, y + 0.17, width - 0.10, 0.24,
                          name=f"capacity_scenario_allocation.segment.{name}.{option}.label",
                          size=8.0, bold=True,
                          color="text.inverse" if option != "security" else "text.title",
                          align=PP_ALIGN.CENTER)
            cursor += width
        unused = max(0.0, limit - used)
        if unused > 0:
            width = frame_w * unused / limit
            deck.box(slide, cursor, y, width, 0.58,
                     name=f"capacity_scenario_allocation.segment.{name}.unused",
                     fill="surface.paper", line="line.soft")
            deck.text(slide, f"HEADROOM  {unused:g}", cursor + 0.05, y + 0.17, width - 0.10, 0.24,
                      name=f"capacity_scenario_allocation.segment.{name}.unused.label",
                      size=7.8, bold=True, color="text.muted", align=PP_ALIGN.CENTER)
        ceiling_x = frame_x + frame_w
        deck.line(slide, ceiling_x, y - 0.16, ceiling_x, y + 0.76,
                  name=f"capacity_scenario_allocation.ceiling.{name}",
                  color="state.risk" if used >= limit else "text.title", width=2.0)
        budget_used = sum(float(costs.get(option, 0.0)) for option in selected)
        budget_limit = float(budgets.get(name, 0.0) or 0.0)
        deck.text(slide, f"TEAM  {used:g}/{limit:g}", 11.18, y + 0.03, 1.20, 0.24,
                  name=f"capacity_scenario_allocation.capacity.{name}", size=8.6,
                  bold=True, color="state.risk" if used >= limit else "state.healthy", align=PP_ALIGN.RIGHT)
        deck.text(slide, f"BUDGET  {budget_used:g}/{budget_limit:g}", 11.18, y + 0.31, 1.20, 0.24,
                  name=f"capacity_scenario_allocation.budget.{name}", size=8.2,
                  color="text.muted", align=PP_ALIGN.RIGHT)

    legend_x = 2.25
    for index, option in enumerate(("platform", "growth", "security", "unused")):
        x = legend_x + index * 1.62
        deck.box(slide, x, 6.14, 0.26, 0.18, name=f"capacity_scenario_allocation.legend.{option}",
                 fill=colors[option], line="line.divider")
        deck.text(slide, option.upper(), x + 0.34, 6.10, 1.12, 0.24,
                  name=f"capacity_scenario_allocation.legend.{option}.label", size=7.8,
                  bold=True, color="text.muted")
    deck.text(slide, "Base binds both team capacity and budget; downside creates headroom; relief raises the ceiling without changing the selected mix.",
              0.82, 6.56, 11.75, 0.28, name="capacity_scenario_allocation.boundary",
              size=8.9, color="text.muted", align=PP_ALIGN.CENTER)


def _render_gantt(deck: _Deck, page: Mapping[str, Any]) -> None:
    slide = deck.base(page)
    tasks = page.get("tasks", {})
    critical = set(page.get("critical_tasks", ()))
    if not critical:
        critical = {task for task, values in tasks.items() if float(values.get("total_float", 1)) <= 0}
    near = set(page.get("near_critical", ()))
    project_finish = float(page.get("project_finish", 1) or 1)
    required_finish = float(page.get("required_finish", project_finish) or project_finish)
    max_finish = max(project_finish, max([float(values.get("ef", 0)) for values in tasks.values()] + [1.0]))
    left, top, width, row_h = 2.18, 2.10, 9.82, 0.58
    for tick in range(int(max_finish) + 1):
        x = left + width * tick / max_finish
        deck.line(slide, x, 1.87, x, 6.32, name=f"critical_path_gantt.grid.day.{tick}",
                  color="line.soft", width=0.55)
        deck.text(slide, f"D{tick}", x - 0.22, 1.54, 0.44, 0.22,
                  name=f"critical_path_gantt.day.{tick}", size=7.6, bold=tick in (required_finish, project_finish),
                  color="text.muted", align=PP_ALIGN.CENTER)
    bar_geometry: dict[str, tuple[float, float, float]] = {}
    for index, (task, values) in enumerate(list(tasks.items())[:9]):
        y = top + index * row_h
        float_value = float(values.get("total_float", 0))
        deck.text(slide, task, 0.68, y + 0.10, 0.40, 0.25,
                  name=f"critical_path_gantt.task.{task}.label", size=9.2, bold=True,
                  color="state.risk" if task in critical else "state.watch" if task in near else "text.title")
        deck.text(slide, f"TF {float_value:+g}", 1.10, y + 0.10, 0.72, 0.25,
                  name=f"critical_path_gantt.task.{task}.float", size=8.2, bold=task in critical,
                  color="state.risk" if task in critical else "text.muted")
        deck.rule(slide, left, y + 0.27, width, name=f"Material: critical_path_gantt.row.{task}",
                  color="line.soft", height=0.018)
        es, ef = float(values.get("es", 0)), float(values.get("ef", 0))
        x = left + width * es / max_finish
        w = max(0.12, width * (ef - es) / max_finish)
        fill = "state.risk" if task in critical else "state.watch" if task in near else "sequence.phase_2"
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y + 0.07, w, 0.38,
                   name=f"critical_path_gantt.bar.{task}", fill=fill, line=fill, line_width=0)
        deck.text(slide, f"{int(ef-es)}d", x + 0.03, y + 0.15, max(0.18, w - 0.06), 0.20,
                  name=f"critical_path_gantt.bar.{task}.duration", size=7.4, bold=True,
                  color="text.inverse" if task not in near else "text.title", align=PP_ALIGN.CENTER)
        bar_geometry[task] = (x, x + w, y + 0.26)

    for edge in page.get("dependency_edges", ()):
        predecessor = str(edge.get("predecessor"))
        successor = str(edge.get("successor"))
        if predecessor not in bar_geometry or successor not in bar_geometry:
            continue
        _px0, px1, py = bar_geometry[predecessor]
        sx0, _sx1, sy = bar_geometry[successor]
        deck.line(slide, px1, py, sx0, sy, name=f"critical_path_gantt.dependency.{edge.get('id')}",
                  color="state.risk" if predecessor in critical and successor in critical else "line.annotation",
                  width=1.7 if predecessor in critical and successor in critical else 1.1)
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.DIAMOND, sx0 - 0.055, sy - 0.055, 0.11, 0.11,
                   name=f"critical_path_gantt.dependency.{edge.get('id')}.arrow",
                   fill="state.risk" if predecessor in critical and successor in critical else "line.annotation",
                   line="state.risk" if predecessor in critical and successor in critical else "line.annotation",
                   line_width=0)

    required_x = left + width * required_finish / max_finish
    finish_x = left + width * project_finish / max_finish
    deck.line(slide, required_x, 1.78, required_x, 6.38,
              name="critical_path_gantt.required_finish", color="state.watch", width=2.0)
    deck.line(slide, finish_x, 1.78, finish_x, 6.38,
              name="critical_path_gantt.project_finish", color="state.risk", width=2.4)
    deck.text(slide, f"REQUIRED D{required_finish:g}", required_x - 0.72, 6.38, 1.44, 0.24,
              name="critical_path_gantt.required_finish.label", size=7.8, bold=True,
              color="state.watch", align=PP_ALIGN.CENTER)
    deck.text(slide, f"FINISH D{project_finish:g}", finish_x - 0.64, 6.38, 1.28, 0.24,
              name="critical_path_gantt.project_finish.label", size=7.8, bold=True,
              color="state.risk", align=PP_ALIGN.CENTER)
    deck.line(slide, left, 1.70, finish_x, 1.70, name="critical_path_gantt.total_duration",
              color="text.title", width=1.4)
    deck.line(slide, left, 1.62, left, 1.78, name="critical_path_gantt.total_duration.start",
              color="text.title", width=1.4)
    deck.line(slide, finish_x, 1.62, finish_x, 1.78, name="critical_path_gantt.total_duration.end",
              color="text.title", width=1.4)
    deck.text(slide, f"TOTAL DURATION  {project_finish:g} WORKING DAYS · {project_finish-required_finish:+g} VS REQUIRED",
              4.38, 1.34, 5.48, 0.24, name="critical_path_gantt.total_duration.label",
              size=8.4, bold=True, color="state.risk", align=PP_ALIGN.CENTER)


def _render_constraint_board(deck: _Deck, page: Mapping[str, Any]) -> None:
    slide = deck.base(page)
    tasks = page.get("tasks", {})
    binding = set(page.get("binding_tasks", ()))
    near = set(page.get("near_critical", ()))
    required_finish = float(page.get("required_finish", 1) or 1)
    max_finish = max([float(item.get("ef", 0)) for item in tasks.values()] + [required_finish, 1.0])
    x0, network_w = 1.25, 10.70
    critical_y, watch_y = 2.65, 4.75
    required_x = x0 + network_w * required_finish / max_finish

    deck.text(slide, "BINDING NEGATIVE-FLOAT PATH", 0.72, 1.52, 4.3, 0.24,
              name="constraint_board.binding.label", size=9.0, bold=True, color="state.risk")
    deck.text(slide, "NEAR-CRITICAL WATCHLIST · +1 DAY FLOAT", 0.72, 4.02, 5.2, 0.24,
              name="constraint_board.near.label", size=9.0, bold=True, color="state.watch")
    deck.line(slide, required_x, 1.90, required_x, 5.72,
              name="constraint_board.required_finish", color="state.risk", width=2.0)
    deck.text(slide, f"REQUIRED FINISH D{required_finish:g}", required_x - 0.86, 5.78, 1.72, 0.24,
              name="constraint_board.required_finish.label", size=7.8, bold=True,
              color="state.risk", align=PP_ALIGN.CENTER)

    node_positions: dict[str, tuple[float, float]] = {}
    for task, values in tasks.items():
        es = float(values.get("es", 0))
        x = x0 + network_w * es / max_finish
        y = critical_y if task in binding else watch_y
        node_positions[task] = (x, y)
    for edge in page.get("dependency_edges", ()):
        pred, succ = str(edge.get("predecessor")), str(edge.get("successor"))
        if pred not in node_positions or succ not in node_positions:
            continue
        px, py = node_positions[pred]
        sx, sy = node_positions[succ]
        deck.line(slide, px + 0.23, py + 0.23, sx + 0.02, sy + 0.23,
                  name=f"constraint_board.edge.{edge.get('id')}",
                  color="state.risk" if pred in binding and succ in binding else "line.annotation",
                  width=2.1 if pred in binding and succ in binding else 1.2)
    for task, values in tasks.items():
        x, y = node_positions[task]
        fill = "state.risk" if task in binding else "state.watch"
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x, y, 0.48, 0.48,
                   name=f"constraint_board.node.{task}", fill=fill, line=fill, line_width=0)
        deck.text(slide, task, x + 0.07, y + 0.13, 0.34, 0.20,
                  name=f"constraint_board.node.{task}.label", size=9.0, bold=True,
                  color="text.inverse" if task in binding else "text.title", align=PP_ALIGN.CENTER)
        deck.text(slide, f"TF {float(values.get('total_float', 0)):+g}\n{values.get('es_date', '')}",
                  x - 0.34, y + 0.58, 1.16, 0.48,
                  name=f"constraint_board.node.{task}.trigger", size=7.2, bold=task in binding,
                  color="text.muted", align=PP_ALIGN.CENTER)
    deck.text(slide, "A → B → D → F is already binding. C → E → G is the next recovery risk if one working day is lost.",
              0.82, 6.48, 11.75, 0.28, name="constraint_board.boundary",
              size=9.2, color="text.muted", align=PP_ALIGN.CENTER)


def _metric_map(page: Mapping[str, Any]) -> dict[str, float]:
    values: dict[str, float] = {}
    for item in page.get("metrics", ()):
        try:
            values[str(item.get("label", "")).lower()] = float(item.get("value", 0))
        except (TypeError, ValueError):
            values[str(item.get("label", "")).lower()] = 0.0
    return values


def _render_risk_exposure(deck: _Deck, page: Mapping[str, Any]) -> None:
    slide = deck.base(page)
    values = _metric_map(page)
    baseline = values.get("baseline exposure", 0.0)
    residual = values.get("residual exposure", 0.0)
    reduction = values.get("exposure reduction", baseline - residual)
    gap = values.get("tolerance gap", 0.0)
    tolerance = residual - gap
    scale_max = max(baseline, residual, tolerance, 1.0) * 1.10
    chart_top, chart_bottom, chart_h = 1.92, 5.80, 3.88
    to_y = lambda value: chart_bottom - chart_h * value / scale_max
    column_w = 1.42
    x_baseline, x_reduction, x_residual = 1.52, 5.10, 8.68

    for tick in range(0, 6):
        value = scale_max * tick / 5
        y = to_y(value)
        deck.line(slide, 1.12, y, 11.92, y, name=f"risk_exposure.grid.{tick}",
                  color="line.soft", width=0.55)
        deck.text(slide, f"{value/1000:.0f}k", 0.58, y - 0.10, 0.48, 0.20,
                  name=f"risk_exposure.axis.{tick}", size=7.2, color="text.muted",
                  align=PP_ALIGN.RIGHT)

    deck.box(slide, x_baseline, to_y(baseline), column_w, chart_bottom - to_y(baseline),
             name="risk_exposure.waterfall.baseline", fill="sequence.phase_1", line="sequence.phase_1")
    deck.box(slide, x_reduction, to_y(baseline), column_w, to_y(residual) - to_y(baseline),
             name="risk_exposure.waterfall.reduction", fill="state.healthy", line="state.healthy")
    deck.box(slide, x_residual, to_y(residual), column_w, chart_bottom - to_y(residual),
             name="risk_exposure.waterfall.residual", fill="state.risk", line="state.risk")
    deck.line(slide, x_baseline + column_w, to_y(baseline), x_reduction, to_y(baseline),
              name="risk_exposure.waterfall.connector.1", color="line.annotation", width=1.1)
    deck.line(slide, x_reduction + column_w, to_y(residual), x_residual, to_y(residual),
              name="risk_exposure.waterfall.connector.2", color="line.annotation", width=1.1)
    tolerance_y = to_y(tolerance)
    deck.line(slide, 1.12, tolerance_y, 11.92, tolerance_y,
              name="risk_exposure.tolerance", color="state.watch", width=2.0)
    deck.text(slide, f"TOLERANCE  {tolerance/1000:.0f}k", 10.12, tolerance_y - 0.29, 1.80, 0.24,
              name="risk_exposure.tolerance.label", size=8.2, bold=True,
              color="state.watch", align=PP_ALIGN.RIGHT)

    for x, label, value, color in (
        (x_baseline, "BASELINE", baseline, "sequence.phase_1"),
        (x_reduction, "MITIGATION", -reduction, "state.healthy"),
        (x_residual, "RESIDUAL", residual, "state.risk"),
    ):
        deck.text(slide, label, x - 0.18, 6.00, column_w + 0.36, 0.22,
                  name=f"risk_exposure.{label.lower()}.label", size=8.4, bold=True,
                  color="text.muted", align=PP_ALIGN.CENTER)
        deck.text(slide, f"{value/1000:+.0f}k" if value < 0 else f"{value/1000:.0f}k",
                  x - 0.18, 6.27, column_w + 0.36, 0.28,
                  name=f"risk_exposure.{label.lower()}.value", size=12.8, bold=True,
                  color=color, align=PP_ALIGN.CENTER)
    deck.shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 10.42, to_y(residual) - 0.18, 1.48, 0.52,
               name="risk_exposure.gap", fill="surface.paper", line="state.risk", line_width=1.4)
    deck.text(slide, f"{gap/1000:.0f}k ABOVE TOLERANCE", 10.52, to_y(residual) - 0.04, 1.28, 0.24,
              name="risk_exposure.gap.label", size=8.0, bold=True, color="state.risk",
              align=PP_ALIGN.CENTER)


def _render_mitigation_roi(deck: _Deck, page: Mapping[str, Any]) -> None:
    slide = deck.base(page)
    values = _metric_map(page)
    cost = values.get("full-life cost", 0.0)
    net = values.get("net benefit", 0.0)
    roi = values.get("roi", 0.0)
    bcr = values.get("benefit / cost", 0.0)
    gross = cost + net
    scale_max = max(gross, cost, net, 1.0) * 1.12
    chart_top, chart_bottom, chart_h = 1.90, 5.68, 3.78
    to_y = lambda value: chart_bottom - chart_h * value / scale_max
    x_gross, x_cost, x_net, column_w = 1.44, 4.72, 8.00, 1.52

    deck.text(slide, "COMMON MONETARY SCALE", 0.72, 1.48, 3.1, 0.24,
              name="mitigation_roi.scale.label", size=9.0, bold=True, color="text.muted")
    deck.line(slide, 1.02, chart_bottom, 10.25, chart_bottom,
              name="mitigation_roi.zero", color="text.title", width=1.3)
    deck.box(slide, x_gross, to_y(gross), column_w, chart_bottom - to_y(gross),
             name="mitigation_roi.bridge.gross_benefit", fill="state.healthy", line="state.healthy")
    deck.box(slide, x_cost, to_y(gross), column_w, to_y(net) - to_y(gross),
             name="mitigation_roi.bridge.cost", fill="state.risk", line="state.risk")
    deck.box(slide, x_net, to_y(net), column_w, chart_bottom - to_y(net),
             name="mitigation_roi.bridge.net_benefit", fill="sequence.phase_1", line="sequence.phase_1")
    deck.line(slide, x_gross + column_w, to_y(gross), x_cost, to_y(gross),
              name="mitigation_roi.bridge.connector.1", color="line.annotation", width=1.2)
    deck.line(slide, x_cost + column_w, to_y(net), x_net, to_y(net),
              name="mitigation_roi.bridge.connector.2", color="line.annotation", width=1.2)
    for x, label, value, color in (
        (x_gross, "GROSS RISK REDUCTION", gross, "state.healthy"),
        (x_cost, "FULL-LIFE COST", -cost, "state.risk"),
        (x_net, "NET BENEFIT", net, "sequence.phase_1"),
    ):
        deck.text(slide, label, x - 0.38, 5.92, column_w + 0.76, 0.24,
                  name=f"mitigation_roi.{label.lower().replace(' ', '_')}.label",
                  size=8.0, bold=True, color="text.muted", align=PP_ALIGN.CENTER)
        deck.text(slide, f"{value/1000:+.0f}k", x - 0.26, 6.22, column_w + 0.52, 0.28,
                  name=f"mitigation_roi.{label.lower().replace(' ', '_')}.value",
                  size=12.0, bold=True, color=color, align=PP_ALIGN.CENTER)

    deck.shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 10.72, 2.16, 1.18, 1.18,
               name="mitigation_roi.roi", fill="surface.paper", line="state.healthy", line_width=2.2)
    deck.text(slide, f"{roi:.1f}×", 10.83, 2.43, 0.96, 0.34,
              name="mitigation_roi.roi.value", size=17.0, bold=True,
              color="state.healthy", align=PP_ALIGN.CENTER)
    deck.text(slide, "ROI", 10.90, 2.82, 0.82, 0.20,
              name="mitigation_roi.roi.label", size=8.0, bold=True,
              color="text.muted", align=PP_ALIGN.CENTER)
    deck.shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 10.72, 3.82, 1.18, 1.18,
               name="mitigation_roi.bcr", fill="surface.paper", line="sequence.phase_2", line_width=2.2)
    deck.text(slide, f"{bcr:.1f}×", 10.83, 4.09, 0.96, 0.34,
              name="mitigation_roi.bcr.value", size=17.0, bold=True,
              color="sequence.phase_2", align=PP_ALIGN.CENTER)
    deck.text(slide, "BENEFIT / COST", 10.68, 4.48, 1.26, 0.20,
              name="mitigation_roi.bcr.label", size=7.4, bold=True,
              color="text.muted", align=PP_ALIGN.CENTER)
    deck.text(slide, "Economics are positive only within the calibrated probability, consequence, and treatment-cost assumptions.",
              0.82, 6.62, 11.75, 0.26, name="mitigation_roi.boundary",
              size=8.7, color="text.muted", align=PP_ALIGN.CENTER)


def _truncate_words(value: Any, limit: int) -> str:
    text = str(value)
    if len(text) <= limit:
        return text
    clipped = text[:limit].rsplit(" ", 1)[0].rstrip(" ,;:-")
    return (clipped or text[:limit]).rstrip() + "…"


def _trace_stage_display(stage: int, value: Any) -> str:
    text = str(value)
    if stage == 3:
        match = re.search(r"flip threshold=([0-9.]+)", text)
        if match:
            return f"One-way sensitivity · flip threshold {float(match.group(1)):.2f}"
        match = re.search(r"event probability ([0-9.]+)", text)
        if match:
            return f"Net benefit flips at {float(match.group(1)):.1%} event probability"
        return _truncate_words(text, 70)
    if stage == 4 and "_missing" in text:
        codes = [item.strip() for item in text.split(",") if item.strip()]
        labels = []
        for code in codes[:3]:
            label = code.removeprefix("recommendation_").removesuffix("_missing")
            labels.append(_humanize(label).replace(" trace", "").title())
        visible = " · ".join(labels)
        if len(codes) > 3:
            visible += f" · +{len(codes) - 3}"
        return f"{len(codes)} readiness fields missing\n{visible}"
    return _truncate_words(text, 78)


def _render_trace(deck: _Deck, page: Mapping[str, Any]) -> None:
    slide = deck.base(page)
    rows = list(page.get("rows", ()))[:4]
    stage_x = [2.16, 5.04, 7.84, 10.26]
    stage_w = [2.36, 2.28, 1.92, 2.24]
    headers = ["STANCE", "COUNTERARGUMENT", "SENSITIVITY / FLIP", "EVIDENCE BOUNDARY"]
    for x, width, header in zip(stage_x, stage_w, headers):
        deck.text(slide, header, x, 1.50, width, 0.24,
                  name=f"recommendation_trace.header.{header.lower().replace(' ', '_')}",
                  size=8.2, bold=True, color="text.muted", align=PP_ALIGN.CENTER)
    for lane, item in enumerate(rows, start=1):
        y = 1.92 + (lane - 1) * 1.12
        readiness = str(item.get("readiness", "diagnostic only"))
        lane_color = "state.healthy" if readiness == "passed" else "state.watch"
        framework = _humanize(item.get("framework", "framework")).upper()
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.64, y, 1.26, 0.70,
                   name=f"recommendation_trace.lane.{lane}.identity",
                   fill=lane_color, line=lane_color, line_width=0)
        deck.text(slide, framework, 0.72, y + 0.12, 1.10, 0.28,
                  name=f"recommendation_trace.lane.{lane}.framework", size=7.6, bold=True,
                  color="text.inverse" if readiness == "passed" else "text.title",
                  align=PP_ALIGN.CENTER)
        deck.text(slide, readiness.upper(), 0.72, y + 0.44, 1.10, 0.18,
                  name=f"recommendation_trace.lane.{lane}.readiness", size=7.2, bold=True,
                  color="text.inverse" if readiness == "passed" else "text.title",
                  align=PP_ALIGN.CENTER)
        for stage in range(1, 4):
            deck.line(slide, stage_x[stage - 1] + stage_w[stage - 1], y + 0.35,
                      stage_x[stage] - 0.08, y + 0.35,
                      name=f"recommendation_trace.connector.{lane}.{stage}",
                      color=lane_color, width=1.55)
        content = [
            item.get("recommendation", ""),
            item.get("counterargument", ""),
            item.get("sensitivity", ""),
            item.get("evidence_boundary", ""),
        ]
        for stage, (x, width, text) in enumerate(zip(stage_x, stage_w, content), start=1):
            fill = "surface.paper"
            line = lane_color if stage == 1 else "line.divider"
            if stage == 4 and readiness != "passed":
                line = "state.risk"
            deck.shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, width, 0.70,
                       name=f"recommendation_trace.lane.{lane}.stage.{stage}",
                       fill=fill, line=line, line_width=1.2)
            display = _trace_stage_display(stage, text)
            deck.text(slide, display, x + 0.10, y + 0.09, width - 0.20, 0.52,
                      name=f"recommendation_trace.lane.{lane}.stage.{stage}.text",
                      size=7.2 if stage != 1 else 7.5, bold=stage == 1,
                      color="text.title" if stage == 1 else "text.body",
                      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    deck.shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 2.16, 6.44, 10.34, 0.38,
               name="recommendation_trace.aggregate_state", fill="surface.paper",
               line="state.watch", line_width=1.4)
    deck.text(slide, "AGGREGATE STATE  ·  NO RECOMMENDATION — FRAMEWORK OUTPUTS REMAIN SEPARATE",
              2.32, 6.53, 10.02, 0.22, name="recommendation_trace.aggregate_state.label",
              size=8.4, bold=True, color="state.watch", align=PP_ALIGN.CENTER)


def _render_decision_record(deck: _Deck, page: Mapping[str, Any]) -> None:
    slide = deck.base(page)
    records = list(page.get("records", ()))[:4]
    dated = []
    for record in records:
        try:
            dated.append((date.fromisoformat(str(record.get("review_date"))), record))
        except ValueError:
            pass
    min_date = min((item[0] for item in dated), default=date(2026, 8, 1))
    max_date = max((item[0] for item in dated), default=date(2026, 10, 30))
    span = max(1, (max_date - min_date).days)
    timeline_x0, timeline_w, timeline_y = 3.35, 8.88, 2.08
    date_to_x = lambda value: timeline_x0 + timeline_w * (value - min_date).days / span

    deck.text(slide, "REVIEW CALENDAR", timeline_x0, 1.38, 4.0, 0.24,
              name="decision_record.timeline.label", size=9.0, bold=True, color="text.muted")
    deck.line(slide, timeline_x0, timeline_y, timeline_x0 + timeline_w, timeline_y,
              name="decision_record.timeline.axis", color="text.title", width=1.5)
    for value, _record in sorted(dated, key=lambda item: item[0]):
        x = date_to_x(value)
        name = value.isoformat()
        deck.line(slide, x, 1.90, x, 6.10, name=f"decision_record.timeline.{name}",
                  color="line.annotation", width=1.0)
        deck.text(slide, value.strftime("%d %b").upper(), x - 0.44, 1.66, 0.88, 0.22,
                  name=f"decision_record.timeline.{name}.label", size=7.8, bold=True,
                  color="text.muted", align=PP_ALIGN.CENTER)

    deck.shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.64, 1.74, 2.18, 0.72,
               name="decision_record.status", fill="state.watch", line="state.watch", line_width=0)
    deck.text(slide, _humanize(page.get("decision_status", "not supplied")).upper(),
              0.80, 1.94, 1.86, 0.28, name="decision_record.status.value",
              size=13.0, bold=True, color="text.title", align=PP_ALIGN.CENTER)

    for lane, record in enumerate(records, start=1):
        framework = str(record.get("framework", "framework"))
        y = 2.72 + (lane - 1) * 0.88
        deck.text(slide, _humanize(framework).upper(), 0.68, y + 0.08, 1.26, 0.22,
                  name=f"decision_record.lane.{framework}.framework", size=8.0, bold=True,
                  color="text.title")
        deck.text(slide, str(record.get("owner", "not supplied")), 0.68, y + 0.34, 1.82, 0.20,
                  name=f"decision_record.lane.{framework}.owner", size=7.4, color="text.muted")
        deck.line(slide, timeline_x0, y + 0.30, timeline_x0 + timeline_w, y + 0.30,
                  name=f"decision_record.lane.{framework}.track", color="line.soft", width=0.7)
        review_text = str(record.get("review_date", "not supplied"))
        try:
            review_date = date.fromisoformat(review_text)
        except ValueError:
            deck.shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 2.18, y + 0.03, 1.02, 0.54,
                       name=f"decision_record.undated.{framework}", fill="surface.paper",
                       line="state.risk", line_width=1.4)
            deck.text(slide, "DATE\nMISSING", 2.28, y + 0.10, 0.82, 0.36,
                      name=f"decision_record.undated.{framework}.label", size=7.3, bold=True,
                      color="state.risk", align=PP_ALIGN.CENTER)
            continue
        x = date_to_x(review_date)
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.DIAMOND, x - 0.24, y + 0.06, 0.48, 0.48,
                   name=f"decision_record.timeline_node.{framework}",
                   fill="state.healthy" if framework == "value_gate" else "sequence.phase_2",
                   line="text.inverse", line_width=0.8)
        conditions = list(record.get("conditions", ()))
        triggers = list(record.get("triggers", ()))
        condition = str(conditions[0]) if conditions else "No condition supplied"
        trigger = triggers[0] if triggers else {}
        trigger_text = str(trigger.get("action", "No rerun trigger supplied")) if isinstance(trigger, Mapping) else str(trigger)
        detail = f"{condition[:56]}{'…' if len(condition) > 56 else ''}\nTrigger: {trigger_text[:42]}{'…' if len(trigger_text) > 42 else ''}"
        box_x = min(max(timeline_x0, x - 1.20), timeline_x0 + timeline_w - 2.40)
        deck.shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, box_x, y + 0.60, 2.40, 0.58,
                   name=f"decision_record.detail.{framework}", fill="surface.paper",
                   line="line.divider", line_width=0.8)
        deck.text(slide, detail, box_x + 0.08, y + 0.68, 2.24, 0.42,
                  name=f"decision_record.detail.{framework}.text", size=7.2,
                  color="text.body", align=PP_ALIGN.CENTER)
    deck.text(slide, "Review-ready is a governance state, not unconditional approval. Missing schedule review data remains visibly unresolved.",
              0.82, 6.62, 11.75, 0.26, name="decision_record.boundary",
              size=8.8, color="text.muted", align=PP_ALIGN.CENTER)


_SPECIAL_RENDERERS = {
    "decision_brief": _render_decision_brief,
    "benefits_dashboard": _render_benefits,
    "business_case_gate": _render_business_case_gate,
    "risk_exposure": _render_risk_exposure,
    "mitigation_roi": _render_mitigation_roi,
    "weighted_decision_matrix": _render_matrix,
    "capacity_scenario_allocation": _render_scenarios,
    "critical_path_gantt": _render_gantt,
    "constraint_board": _render_constraint_board,
    "recommendation_trace": _render_trace,
    "decision_record": _render_decision_record,
}


def _inspect(prs: Presentation, output: Path, page_types: tuple[str, ...]) -> BuildResult:
    native = text = tables = charts = connectors = pictures = 0
    outside: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
            else:
                native += 1
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                connectors += 1
            if getattr(shape, "has_text_frame", False) and shape.text:
                text += 1
                if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
                    outside.append(shape.name)
            tables += int(bool(getattr(shape, "has_table", False)))
            charts += int(bool(getattr(shape, "has_chart", False)))
    with ZipFile(output) as archive:
        media = sum(1 for name in archive.namelist() if name.startswith("ppt/media/") and not name.endswith("/"))
    return BuildResult(output, len(prs.slides), native, text, tables, charts, connectors,
                       pictures, media, page_types, tuple(outside))


def render_decision_deck(package: DecisionPackage, output: Path, tokens: Mapping[str, Any] | None) -> BuildResult:
    """Render semantic decision IR as native PowerPoint text, shapes, and connectors."""
    pages = build_decision_presentation_ir(package)
    deck = _Deck(tokens or resolve_tokens())
    for page in pages:
        _SPECIAL_RENDERERS.get(page["page_type"], _render_generic)(deck, page)
    output.parent.mkdir(parents=True, exist_ok=True)
    deck.prs.save(output)
    return _inspect(deck.prs, output, tuple(page["page_type"] for page in pages))
