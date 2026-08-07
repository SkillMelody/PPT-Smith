from __future__ import annotations

import importlib.metadata
import math
from pathlib import Path
from typing import Any

from .base import BuildInspection, BuildPlan, BuildResult, BuilderCapability, SupportLevel


class PythonPptxAdapter:
    name = "python_pptx"

    def probe(self) -> BuilderCapability:
        version = _package_version()
        available = version is not None
        return BuilderCapability(
            name=self.name,
            available=available,
            version=version,
            command="python" if available else None,
            features={
                "native_text": True,
                "native_table": True,
                "native_chart": True,
                "native_connector": True,
                "svg_embed": False,
                "speaker_notes": False,
                "theme": "partial",
                "slide_master": "partial",
                "readback": True,
            },
            components={
                "causal_chain": "full",
                "heat_matrix": "full",
                "layered_architecture": "full",
                "drill_down_stair": "full",
                "phase_roadmap": "full",
                "product_ui_overview": "full",
                "interaction_storyboard": "full",
            } if available else {},
            warnings=[] if available else ["python-pptx package is not importable."],
        )

    def supports(self, component_type: str, delivery_route: str) -> SupportLevel:
        if (
            component_type
            in {
                "causal_chain",
                "heat_matrix",
                "layered_architecture",
                "drill_down_stair",
                "phase_roadmap",
                "product_ui_overview",
                "interaction_storyboard",
            }
            and delivery_route == "native_diagram"
        ):
            return "full"
        if delivery_route in {"native_ppt", "native_table", "native_chart"}:
            return "full"
        if delivery_route in {"generated_image", "background_image", "hybrid_overlay"}:
            return "visual_only"
        return "unsupported"

    def plan(self, ppt_ir: dict[str, Any], style_contract: dict[str, Any], delivery_plan: dict[str, Any]) -> BuildPlan:
        object_plans = {
            (obj.get("slide_id"), obj.get("object_id")): obj
            for slide in delivery_plan.get("slides", []) or []
            for obj in slide.get("objects", []) or []
            if isinstance(obj, dict)
        }
        slides = []
        for slide in ppt_ir.get("slides", []) or []:
            planned_objects = []
            for obj in slide.get("objects", []) or []:
                planned = dict(obj)
                planned["delivery_plan"] = object_plans.get((slide.get("id"), obj.get("id")), {})
                planned_objects.append(planned)
            slide_plan = dict(slide)
            slide_plan["objects"] = planned_objects
            slide_plan["style_contract"] = style_contract
            slides.append(slide_plan)
        warnings = _unsupported_style_warnings(style_contract)
        return BuildPlan(builder=self.name, slides=slides, style_contract=dict(style_contract), warnings=warnings)

    def build(self, build_plan: BuildPlan, output_dir: Path, *, visual_plan: dict[str, Any] | None = None) -> BuildResult:
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
            from pptx.util import Inches, Pt
        except Exception as exc:
            return BuildResult(builder=self.name, status="failed", errors=[{"code": "BUILDER_RUNTIME_IMPORT_FAILED", "message": str(exc)}])

        output_dir.mkdir(parents=True, exist_ok=True)
        deck_path = output_dir / "deck.pptx"
        style = _resolve_style(build_plan.style_contract)
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        blank_layout = presentation.slide_layouts[6]

        # Index visual plan slides by id for layout hints
        vp_slides: dict[str, dict[str, Any]] = {}
        if isinstance(visual_plan, dict):
            for vs in visual_plan.get("slides", []) or []:
                if isinstance(vs, dict) and vs.get("slide_id"):
                    vp_slides[vs["slide_id"]] = vs

        object_results: list[dict[str, Any]] = []
        fallbacks: list[dict[str, Any]] = []
        warnings: list[str] = list(build_plan.warnings)

        for slide_number, slide_plan in enumerate(build_plan.slides, start=1):
            slide = presentation.slides.add_slide(blank_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor.from_string(style["background"])
            # --- visual-plan-driven layout hints ---
            vp = vp_slides.get(slide_plan.get("id", "")) or {}
            vp_layout = vp.get("layout") if isinstance(vp.get("layout"), dict) else {}
            vp_density = vp.get("density") if isinstance(vp.get("density"), dict) else {}
            primary_zone_ratio = float(vp_layout.get("primary_zone_ratio", 0.56))
            density_level = str(vp_density.get("level", "medium"))
            page_archetype = str(vp.get("page_archetype", ""))
            if slide_plan.get("slide_role") == "cover":
                _add_cover(slide, slide_plan, style, vp)
                if not (slide_plan.get("objects") or []):
                    continue
            elif slide_plan.get("slide_role") == "closing":
                _add_closing(slide, slide_plan, style, slide_number)
                if not (slide_plan.get("objects") or []):
                    continue
            else:
                title = str(slide_plan.get("title") or "Untitled")
                cjk_chars = sum(1 for c in title if '\u4e00' <= c <= '\u9fff')
                # Allow up to 2 lines; shrink if text overflows
                title_h = 0.82 if len(title) > 28 else 0.62
                title_size = style["title_size"] if len(title) <= 28 else style["title_size"] - 2.0
                _add_textbox(
                    slide,
                    title,
                    x=style["margin_left"],
                    y=style["margin_top"],
                    w=12.2,
                    h=title_h,
                    font_size=title_size,
                    bold=True,
                    color=style["primary"],
                    font_name=style["font_name"],
                    shrink=True,
                    min_size=max(14.0, style["body_size"]),
                )
            body_parts: list[str] = []
            seen_text = {
                str(slide_plan.get("title") or "").strip()
            }
            for part in (slide_plan.get("judgment"), slide_plan.get("message")):
                if not isinstance(part, str):
                    continue
                normalized = part.strip()
                if not normalized or normalized in seen_text:
                    continue
                body_parts.append(normalized)
                seen_text.add(normalized)
            body = "\n".join(body_parts)
            if body and slide_plan.get("slide_role") != "cover":
                _add_textbox(slide, body, x=style["margin_left"], y=1.05, w=11.7, h=0.45, font_size=style["body_size"], color=style["text_secondary"], font_name=style["font_name"], shrink=True, min_size=max(8.0, style["body_size"] - 3))
            if slide_plan.get("slide_role") != "cover":
                _add_footer(slide, slide_plan, style, slide_number)
            objects = slide_plan.get("objects", []) or []
            if not objects:
                continue
            # --- visual-plan-driven layout: dynamic position & height ---
            title_zone_h = 0.58 if len(title) <= 28 else 0.78
            content_top = style["margin_top"] + title_zone_h
            if body and slide_plan.get("slide_role") != "cover":
                content_top += 0.4  # body text zone
            footer_h = 0.35 if style["footer_enabled"] else 0.0
            available_h = 7.5 - content_top - style["margin_bottom"] - footer_h - 0.3
            if vp:
                # Visual-plan-driven: use primary_zone_ratio and density from the plan
                content_height = available_h * primary_zone_ratio
                dens_scale = {"high": 1.0, "medium": 0.92, "low": 0.82}.get(density_level, 0.88)
                content_height *= dens_scale
                y = content_top + 0.15
            else:
                # No visual plan: use classic hardcoded layout
                y = 1.85
                content_height = 4.65
            card_width = max(2.6, min(3.6, 10.8 / max(len(objects), 1)))
            for index, obj in enumerate(objects):
                route = ((obj.get("delivery_plan") or {}).get("selected_route") or obj.get("delivery_preferences", {}).get("preferred_route") or "native_ppt")
                component_type = str(obj.get("component_type") or obj.get("type") or "")
                renderer = _semantic_renderer(component_type)
                if renderer is not None and route == "native_diagram":
                    rendered = renderer(
                        slide,
                        obj,
                        slide_plan.get("page_design_intent") if isinstance(slide_plan.get("page_design_intent"), dict) else {},
                        slide_plan.get("style_contract") if isinstance(slide_plan.get("style_contract"), dict) else {},
                    )
                    warnings.extend(rendered.get("warnings", []))
                    object_results.append({
                        "slide_id": slide_plan.get("id"),
                        "object_id": obj.get("id"),
                        "component_type": component_type,
                        "planned_route": route,
                        "actual_route": rendered["actual_route"],
                        "status": "created",
                        "object_names": rendered.get("object_names", []),
                    })
                    continue
                actual_route = route
                component_type = obj.get("component_type") or obj.get("type")
                is_chart = route == "native_chart" and obj.get("type") == "chart"
                is_process = route == "native_diagram" and component_type == "process"
                if route in {"hybrid_overlay", "svg_component", "generated_image", "background_image"}:
                    actual_route = "native_ppt"
                    fallbacks.append({
                        "slide_id": slide_plan.get("id"),
                        "object_id": obj.get("id"),
                        "component_type": component_type,
                        "planned_route": route,
                        "actual_route": actual_route,
                        "reason_codes": ["PYTHON_PPTX_MINIMAL_NATIVE_FALLBACK"],
                        "editable_core_preserved": obj.get("editability") != "native_required",
                    })
                single_object = len(objects) == 1
                x = (style["margin_left"] if single_object else style["margin_left"] + index * (card_width + style["card_gap"]))
                object_width = 11.7 if single_object else card_width
                if route == "native_table" or obj.get("type") == "table":
                    _add_table(
                        slide,
                        obj,
                        x=x,
                        y=y,
                        w=object_width,
                        h=content_height if single_object else 2.3,
                        style=style,
                    )
                elif is_chart:
                    _add_chart(
                        slide,
                        obj,
                        x=x,
                        y=y,
                        w=object_width,
                        h=content_height if single_object else 3.4,
                        style=style,
                    )
                elif is_process:
                    process_height = 2.2 if single_object else 1.35
                    process_y = (
                        y + (content_height - process_height) / 2
                        if single_object
                        else y
                    )
                    _add_process(
                        slide,
                        obj,
                        x=x,
                        y=process_y,
                        w=object_width,
                        h=process_height,
                        style=style,
                    )
                elif component_type == "metric_card" or (isinstance(obj.get("content"), dict) and obj["content"].get("cards")):
                    _add_metric_wall(
                        slide,
                        obj,
                        x=x,
                        y=y,
                        w=object_width,
                        h=content_height if single_object else 1.35,
                        style=style,
                    )
                else:
                    _add_card(
                        slide,
                        obj,
                        x=x,
                        y=y,
                        w=object_width,
                        h=content_height if single_object else 1.35,
                        style=style,
                    )
                object_results.append(
                    {
                        "slide_id": slide_plan.get("id"),
                        "object_id": obj.get("id"),
                        "component_type": obj.get("component_type") or obj.get("type"),
                        "planned_route": route,
                        "actual_route": actual_route,
                        "status": "created",
                    }
                )

        # --- quality disclaimer (final slide) ---
        _add_disclaimer_page(presentation, style)

        presentation.save(deck_path)
        return BuildResult(builder=self.name, status="created", pptx=str(deck_path), object_results=object_results, fallbacks=fallbacks, warnings=warnings)

    def inspect_output(self, pptx_path: Path) -> BuildInspection:
        try:
            from ppt_qa.package_inspector import inspect_package
        except Exception as exc:
            return BuildInspection(builder=self.name, status="failed", errors=[{"code": "BUILDER_INSPECT_IMPORT_FAILED", "message": str(exc)}])
        inspection = inspect_package(pptx_path)
        return BuildInspection(
            builder=self.name,
            status="passed" if inspection.status == "passed" else "failed",
            issues=[issue.__dict__ for issue in inspection.issues],
        )


def _package_version() -> str | None:
    try:
        return importlib.metadata.version("python-pptx")
    except importlib.metadata.PackageNotFoundError:
        return None


def _semantic_renderer(component_type: str) -> Any:
    try:
        from scripts.visual_narrative.renderers import get_renderer
    except ImportError:
        return None
    return get_renderer(component_type)


def _text_metrics(text: str, font_size_pt: float) -> tuple[float, float]:
    """Estimate rendered text extent: (width_inches, lines_if_wrapped_at_12in)"""
    if not text:
        return 0.0, 0.0
    cjk = sum(1 for c in text if '\u4e00' <= c <= '\u9fff' or '\u3000' <= c <= '\u303f' or '\uff00' <= c <= '\uffef')
    latin = len(text) - cjk
    width_pt = (cjk * 1.0 + latin * 0.55) * font_size_pt
    return width_pt / 72.0, font_size_pt * 1.35 / 72.0


def _fit_text_size(text: str, max_w_in: float, max_h_in: float, start_pt: float, min_pt: float = 8.0) -> float:
    """Return the largest font size ≤ start_pt that fits within max_w × max_h."""
    if not text or max_w_in <= 0 or max_h_in <= 0:
        return max(start_pt, min_pt)
    size = start_pt
    while size > min_pt:
        est_w, line_h = _text_metrics(text, size)
        lines = max(1, round(est_w / max_w_in))
        if lines * line_h <= max_h_in:
            return size
        size -= 1.0
    return max(min_pt, size)


def _add_textbox(slide: Any, text: str, *, x: float, y: float, w: float, h: float, font_size: float, color: str, bold: bool = False, font_name: str | None = None, shrink: bool = False, min_size: float = 8.0) -> Any:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.util import Inches, Pt

    if shrink:
        font_size = _fit_text_size(str(text), w, h, font_size, min_size)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = "PingFang SC" if any("\u4e00" <= char <= "\u9fff" for char in str(text)) else font_name
    run.font.color.rgb = RGBColor.from_string(color)
    return box


def _add_cover(slide: Any, slide_plan: dict[str, Any], style: dict[str, Any], vp: dict[str, Any] | None = None) -> None:
    """Render cover slide with style determined by visual plan or data."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches

    title = str(slide_plan.get("title") or "Untitled")
    subtitle = str(slide_plan.get("message") or "").strip()
    if subtitle == title:
        subtitle = ""
    # Determine cover style from visual plan, page_design_intent, or fallback to title length
    route = (vp or {}).get("professional_route", "")
    archetype = (vp or {}).get("page_archetype", "")
    # navy_hero: dark bg for decision/investor decks; minimal_light for reports/research
    if route in {"strategic_decision", "investor_commercial"} or archetype == "navy_hero":
        _cover_navy_hero(slide, title, subtitle, style)
    else:
        _cover_minimal_light(slide, title, subtitle, style)


def _cover_navy_hero(slide: Any, title: str, subtitle: str, style: dict[str, Any]) -> None:
    """Dark navy background cover — McKinsey board-deck style."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches

    # Full-slide dark background
    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.name = "Decoration:cover:bg"
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(style["primary"])
    bg.line.fill.background()
    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(0.08), Inches(2.8)
    )
    bar.name = "Decoration:cover:bar"
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor.from_string((style.get("data_series", ["2251FF"])[0] if style.get("data_series") else "2251FF").lstrip("#"))
    bar.line.fill.background()
    # Title
    _add_textbox(slide, title, x=1.1, y=2.2, w=11.2, h=1.6,
        font_size=min(38.0, style["title_size"] * 1.6), bold=True,
        color=style["background"], font_name=style["font_name"], shrink=True, min_size=22)
    if subtitle:
        _add_textbox(slide, subtitle, x=1.1, y=4.0, w=10.0, h=0.7,
            font_size=max(14.0, style["body_size"]), color=style.get("text_secondary", "8899AA"),
            font_name=style["font_name"])
    # Meta line
    _add_textbox(slide, "—— MeowClaw PPT Smith 自动生成 ——",
        x=0.8, y=6.6, w=8.0, h=0.3,
        font_size=style["footnote_size"], color=style.get("text_secondary", "8899AA"), font_name=style["font_name"])


def _cover_minimal_light(slide: Any, title: str, subtitle: str, style: dict[str, Any]) -> None:
    """Clean light-background cover — research report / internal style."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(style["margin_left"]), Inches(1.45), Inches(0.12), Inches(3.75),
    )
    accent.name = "Decoration:cover:accent"
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor.from_string(style["primary"])
    accent.line.fill.background()
    _add_textbox(slide, title, x=style["margin_left"] + 0.42, y=2.0, w=10.8, h=1.65,
        font_size=max(32.0, style["title_size"] * 1.55), bold=True,
        color=style["primary"], font_name=style["font_name"], shrink=True, min_size=20)
    if subtitle:
        _add_textbox(slide, subtitle, x=style["margin_left"] + 0.42, y=4.0, w=9.8, h=0.75,
            font_size=max(14.0, style["body_size"]), color=style["text_secondary"], font_name=style["font_name"])
    rule = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(style["margin_left"] + 0.42), Inches(5.25), Inches(4.0), Inches(0.04),
    )
    rule.name = "Decoration:cover:rule"
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor.from_string(style["border"])
    rule.line.fill.background()


def _add_closing(slide: Any, slide_plan: dict[str, Any], style: dict[str, Any], slide_number: int) -> None:
    """Render professional closing/thank-you slide."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    title = str(slide_plan.get("title") or "Thank You / 谢谢")
    message = str(slide_plan.get("message") or "").strip()
    if message == title:
        message = ""
    # Thank you in large text
    _add_textbox(slide, title, x=style["margin_left"], y=2.0, w=12.2, h=1.2,
        font_size=max(36.0, style["title_size"] * 1.5), bold=True,
        color=style["primary"], font_name=style["font_name"], shrink=True, min_size=22)
    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(style["margin_left"]), Inches(3.4), Inches(2.5), Inches(0.06),
    )
    bar.fill.solid()
    accent_color = style["primary"]
    ds = style.get("data_series")
    if isinstance(ds, list) and ds and isinstance(ds[0], str) and len(ds[0].lstrip("#")) >= 6:
        accent_color = ds[0]
    bar.fill.fore_color.rgb = RGBColor.from_string(accent_color.lstrip("#"))
    bar.line.fill.background()
    if message:
        _add_textbox(slide, message, x=style["margin_left"], y=3.8, w=11.0, h=0.6,
            font_size=style["body_size"] + 2, color=style["text_secondary"], font_name=style["font_name"])
    # Contact / feedback line
    _add_textbox(slide, "如有问题或建议，欢迎通过 GitHub Issues 或公众号留言反馈",
        x=style["margin_left"], y=5.0, w=11.0, h=0.35,
        font_size=style["body_size"], color=style["text_secondary"], font_name=style["font_name"])
    # Footer
    _add_footer(slide, slide_plan, style, slide_number)


def _add_footer(
    slide: Any,
    slide_plan: dict[str, Any],
    style: dict[str, Any],
    slide_number: int,
) -> None:
    if not style["footer_enabled"]:
        return
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    y = 7.5 - style["margin_bottom"] - style["footer_height"]
    divider = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(style["margin_left"]),
        Inches(y),
        Inches(12.2),
        Inches(0.012),
    )
    divider.name = "Decoration:footer:divider"
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor.from_string(style["border"])
    divider.line.fill.background()

    refs = [
        item
        for item in slide_plan.get("source_refs", []) or []
        if isinstance(item, dict)
    ]
    source_parts = []
    for item in refs:
        source_id = str(item.get("source_id") or "").strip()
        locator = str(item.get("locator") or "").strip()
        combined = " · ".join(part for part in (source_id, locator) if part)
        if combined and combined not in source_parts:
            source_parts.append(combined)
    source_text = "来源：" + "；".join(source_parts) if source_parts else "来源：内部综合"
    source = _add_textbox(
        slide,
        source_text,
        x=style["margin_left"],
        y=y + 0.055,
        w=10.8,
        h=style["footer_height"],
        font_size=style["footnote_size"],
        color=style["text_secondary"],
        font_name=style["font_name"],
    )
    source.name = "source_note"
    page = _add_textbox(
        slide,
        f"{slide_number:02d}",
        x=11.9,
        y=y + 0.055,
        w=0.85,
        h=style["footer_height"],
        font_size=style["footnote_size"],
        color=style["text_secondary"],
        font_name=style["font_name"],
    )
    page.name = "footer_page"
    page.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def _metric_entries(obj: dict[str, Any]) -> list[dict[str, str]]:
    content = obj.get("content")
    if not isinstance(content, dict):
        return []
    # content.cards[] — shape/card-based IR (most common)
    cards = content.get("cards")
    if isinstance(cards, list):
        entries: list[dict[str, str]] = []
        for index, item in enumerate(cards, start=1):
            if not isinstance(item, dict):
                continue
            value = item.get("value")
            title = item.get("title") or item.get("label")
            detail = item.get("detail") or item.get("description")
            if value is not None:
                entries.append({
                    "label": str(title or f"指标 {index:02d}"),
                    "value": (
                        f"{item.get('prefix') or ''}"
                        f"{value}"
                        f"{item.get('suffix') or ''}"
                    ),
                })
            elif title:
                # Action card: title only, no metric value
                num = item.get("num")
                prefix = f"{num}. " if num is not None else ""
                display_value = prefix + str(title)
                detail_text = str(detail) if detail else ""
                entries.append({
                    "label": detail_text[:60] if detail_text else display_value,
                    "value": display_value,
                })
        if entries:
            return entries
    # content.metrics[] — metric_wall IR
    metrics = content.get("metrics")
    if isinstance(metrics, list):
        entries = []
        for index, item in enumerate(metrics, start=1):
            if not isinstance(item, dict):
                continue
            value = item.get("value")
            if value is None:
                continue
            entries.append(
                {
                    "label": str(item.get("label") or f"指标 {index:02d}"),
                    "value": (
                        f"{item.get('prefix') or ''}"
                        f"{value}"
                        f"{item.get('suffix') or ''}"
                    ),
                }
            )
        if entries:
            return entries
    raw_values = content.get("values")
    labels = content.get("labels")
    if isinstance(raw_values, list):
        return [
            {
                "label": (
                    str(labels[index])
                    if isinstance(labels, list) and index < len(labels)
                    else f"指标 {index + 1:02d}"
                ),
                "value": str(value),
            }
            for index, value in enumerate(raw_values)
        ]
    quartiles = content.get("quartiles")
    if isinstance(quartiles, list):
        quartile_labels = ["25 分位", "中位数", "75 分位"]
        return [
            {
                "label": quartile_labels[index] if index < len(quartile_labels) else f"分位 {index + 1}",
                "value": f"{value}%",
            }
            for index, value in enumerate(quartiles)
        ]
    return []


def _add_metric_wall(
    slide: Any,
    obj: dict[str, Any],
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    style: dict[str, Any],
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches

    entries = _metric_entries(obj)
    if not entries:
        _add_card(slide, obj, x=x, y=y, w=w, h=h, style=style)
        return
    columns = 3 if len(entries) <= 6 else 4
    rows = math.ceil(len(entries) / columns)
    gap = style["card_gap"]
    wall_h = min(h, rows * 2.3 + gap * (rows - 1))
    wall_y = y + max(0.0, (h - wall_h) / 2)
    card_w = (w - gap * (columns - 1)) / columns
    card_h = (wall_h - gap * (rows - 1)) / rows
    for index, entry in enumerate(entries):
        row, column = divmod(index, columns)
        left = x + column * (card_w + gap)
        top = wall_y + row * (card_h + gap)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(card_w),
            Inches(card_h),
        )
        card.name = f"Material:metric_wall:card:{index + 1:02d}"
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string(style["surface"])
        card.line.color.rgb = RGBColor.from_string(style["border"])
        value_box = _add_textbox(
            slide,
            entry["value"],
            x=left + 0.22,
            y=top + 0.38,
            w=card_w - 0.44,
            h=max(0.58, card_h * 0.42),
            font_size=30,
            bold=True,
            color=style["primary"],
            font_name=style["font_name"],
        )
        value_box.name = f"Component:metric_wall:value:{index + 1:02d}"
        value_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        value_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        label_box = _add_textbox(
            slide,
            entry["label"],
            x=left + 0.22,
            y=top + card_h * 0.64,
            w=card_w - 0.44,
            h=max(0.35, card_h * 0.2),
            font_size=max(11, style["table_size"]),
            color=style["text_secondary"],
            font_name=style["font_name"],
        )
        label_box.name = f"Component:metric_wall:label:{index + 1:02d}"
        label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def _object_text(obj: dict[str, Any]) -> str:
    content = obj.get("content")
    if isinstance(content, str) and content.strip():
        return content.strip()
    if isinstance(content, dict):
        if all(isinstance(content.get(key), str) and content[key].strip() for key in ("title", "claim")):
            return f"{content['title'].strip()}\n{content['claim'].strip()}"
        for key in ("title", "label", "claim", "text", "value"):
            value = content.get(key)
            if isinstance(value, str) and value.strip():
                return value.strip()
    component = obj.get("component_type") or obj.get("semantic_role") or obj.get("type")
    return str(component or "PPT object").replace("_", " ").title()


def _add_card(
    slide: Any,
    obj: dict[str, Any],
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    style: dict[str, Any],
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.name = f"Component:card:{obj.get('id') or 'primary'}"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(style["surface"])
    shape.line.color.rgb = RGBColor.from_string(style["border"])
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.margin_left = Inches(0.42)
    text_frame.margin_right = Inches(0.42)
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = _object_text(obj)
    run.font.size = Pt(
        max(20.0, style["body_size"] * 1.6)
        if h >= 3.0
        else style["body_size"]
    )
    run.font.name = style["font_name"]
    run.font.color.rgb = RGBColor.from_string(style["text_primary"])


def _table_data(obj: dict[str, Any]) -> list[list[str]]:
    content = obj.get("content")
    if isinstance(content, dict):
        rows = content.get("rows")
        if isinstance(rows, list) and rows:
            return [[str(cell) for cell in row] for row in rows if isinstance(row, list)]
        headers = content.get("headers")
        body = content.get("body")
        if isinstance(headers, list) and isinstance(body, list):
            return [[str(cell) for cell in headers], *[[str(cell) for cell in row] for row in body if isinstance(row, list)]]
    return [["Item", "Status"], [_object_text(obj), "Editable"]]


def _chart_data(obj: dict[str, Any]) -> tuple[list[str], list[tuple[str, list[float]]]]:
    content = obj.get("content")
    if not isinstance(content, dict):
        raise ValueError(f"Chart {obj.get('id')} requires mapping content")
    categories = content.get("categories")
    raw_series = content.get("series")
    if not isinstance(categories, list) or not categories or not isinstance(raw_series, list) or not raw_series:
        raise ValueError(f"Chart {obj.get('id')} requires non-empty categories and series")
    parsed_series: list[tuple[str, list[float]]] = []
    for series in raw_series:
        if not isinstance(series, dict) or not isinstance(series.get("values"), list):
            raise ValueError(f"Chart {obj.get('id')} has invalid series")
        values = series["values"]
        if len(values) != len(categories) or not all(isinstance(value, (int, float)) for value in values):
            raise ValueError(f"Chart {obj.get('id')} series values must align with categories")
        parsed_series.append((str(series.get("name") or "Series"), [float(value) for value in values]))
    return [str(category) for category in categories], parsed_series


def _add_chart(slide: Any, obj: dict[str, Any], *, x: float, y: float, w: float, h: float, style: dict[str, Any]) -> None:
    from pptx.chart.data import ChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
    from pptx.util import Inches, Pt

    categories, series = _chart_data(obj)
    data = ChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    # Choose horizontal bars for many categories, vertical columns for few
    chart_type = XL_CHART_TYPE.BAR_CLUSTERED if len(categories) > 5 else XL_CHART_TYPE.COLUMN_CLUSTERED
    chart_frame = slide.shapes.add_chart(
        chart_type, Inches(x), Inches(y), Inches(max(w, 5.8)), Inches(h), data
    )
    chart = chart_frame.chart
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.name = style["font_name"]
        chart.legend.font.size = Pt(style["table_size"])
    chart.value_axis.has_major_gridlines = True
    # Scale axis labels: narrower when many categories
    axis_label_size = max(7.0, style["table_size"] - (len(categories) - 5) * 0.8 if len(categories) > 5 else style["table_size"])
    chart.category_axis.tick_labels.font.name = style["font_name"]
    chart.category_axis.tick_labels.font.size = Pt(axis_label_size)
    # --- data series colour injection ---
    data_series_colors = style.get("data_series")
    if isinstance(data_series_colors, list) and data_series_colors:
        try:
            plot = chart.plots[0]
            for i, s in enumerate(plot.series):
                color_hex = data_series_colors[i % len(data_series_colors)]
                if isinstance(color_hex, str):
                    s.format.fill.solid()
                    s.format.fill.fore_color.rgb = RGBColor.from_string(color_hex)
        except Exception:
            pass
    # --- data labels ---
    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(max(8, style["table_size"] - 1))
        data_labels.font.name = style["font_name"]
        data_labels.number_format = "0"
        data_labels.show_value = True
        if chart_type == XL_CHART_TYPE.BAR_CLUSTERED:
            data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    except Exception:
        pass


def _process_labels(obj: dict[str, Any]) -> list[str]:
    content = obj.get("content")
    nodes = content.get("nodes") if isinstance(content, dict) else None
    if not isinstance(nodes, list) or not nodes:
        raise ValueError(f"Process {obj.get('id')} requires non-empty nodes")
    labels = [str(node.get("label")) for node in nodes if isinstance(node, dict) and node.get("label") is not None]
    if len(labels) != len(nodes):
        raise ValueError(f"Process {obj.get('id')} requires a label for every node")
    return labels


def _add_routed_connector(
    slide: Any,
    start: Any,
    end: Any,
    *,
    route: str = "straight",
    feedback: bool = False,
    color: str = "2D3340",
    width_pt: float = 1.5,
) -> Any:
    """Add one native, endpoint-bound arrow using the requested route."""
    from lxml import etree
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn
    from pptx.util import Pt

    connector_type = {
        "straight": MSO_CONNECTOR.STRAIGHT,
        "orthogonal": MSO_CONNECTOR.ELBOW,
        "curved": MSO_CONNECTOR.CURVE,
    }.get(route)
    if connector_type is None:
        raise ValueError(f"Unsupported connector route: {route}")

    start_x = start.left + start.width
    start_y = start.top + start.height // 2
    end_x = end.left
    end_y = end.top + end.height // 2
    if feedback:
        start_x = start.left + start.width // 2
        start_y = start.top + start.height
        end_x = end.left + end.width // 2
        end_y = end.top

    connector = slide.shapes.add_connector(
        connector_type,
        start_x,
        start_y,
        end_x,
        end_y,
    )
    connector.name = f"Connector:{route}"
    connector.line.color.rgb = RGBColor.from_string(color)
    connector.line.width = Pt(width_pt)
    connector.begin_connect(start, 3 if not feedback else 2)
    connector.end_connect(end, 1 if not feedback else 0)
    line_xml = connector.line._get_or_add_ln()
    tail_end = etree.Element(qn("a:tailEnd"))
    tail_end.set("type", "triangle")
    line_xml.append(tail_end)
    return connector


def _add_process(slide: Any, obj: dict[str, Any], *, x: float, y: float, w: float, h: float, style: dict[str, Any]) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt

    labels = _process_labels(obj)
    total_width = max(w, 5.8)
    gap = 0.28
    node_width = (total_width - gap * (len(labels) - 1)) / len(labels)
    shapes = []
    for index, label in enumerate(labels):
        left = x + index * (node_width + gap)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(y), Inches(node_width), Inches(h))
        shape.name = f"Component:process:{obj.get('id') or 'primary'}:node:{index + 1:02d}"
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(style["surface"])
        shape.line.color.rgb = RGBColor.from_string(style["border"])
        shape.text = label
        shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name = style["font_name"]
            run.font.size = Pt(
                max(18.0, style["body_size"] * 1.4)
                if h >= 2.0
                else style["body_size"]
            )
            run.font.color.rgb = RGBColor.from_string(style["text_primary"])
        shapes.append(shape)
    route = str(obj.get("connector_route") or "straight")
    for before, after in zip(shapes, shapes[1:]):
        _add_routed_connector(
            slide,
            before,
            after,
            route=route,
            color=style["primary"],
        )


def _add_table(slide: Any, obj: dict[str, Any], *, x: float, y: float, w: float, h: float, style: dict[str, Any]) -> None:
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt

    data = _table_data(obj)
    rows = max(len(data), 1)
    cols = max(len(row) for row in data) if data else 1
    compact = rows <= 3
    table_height = min(h, max(1.8, rows * 0.8)) if compact else h
    table_y = y + max(0.0, (h - table_height) / 2)
    table_shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(x),
        Inches(table_y),
        Inches(w),
        Inches(table_height),
    )
    table_shape.name = f"Component:table:{obj.get('id') or 'primary'}"
    table = table_shape.table
    # Proportional column widths based on content length
    col_max_len: list[int] = [0] * cols
    for row_data in data:
        for ci in range(min(cols, len(row_data))):
            col_max_len[ci] = max(col_max_len[ci], len(str(row_data[ci])))
    total_len = sum(col_max_len) or 1
    col_widths = [max(0.12, (cl / total_len) * w) for cl in col_max_len]
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)
    # Stripe colours
    surface2 = style.get("surface_2", "EBEDF0")
    border_color = style.get("border", "D1D5DB")
    for row_index, row_data in enumerate(data):
        is_header = row_index == 0
        for col_index in range(cols):
            cell = table.cell(row_index, col_index)
            cell.text = row_data[col_index] if col_index < len(row_data) else ""
            # Fill
            cell.fill.solid()
            if is_header:
                cell.fill.fore_color.rgb = RGBColor.from_string(style["primary"])
            elif row_index % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor.from_string(surface2)
            else:
                cell.fill.fore_color.rgb = RGBColor.from_string(style["background"])
            # Thin inner border via XML (python-pptx doesn't expose per-cell borders easily)
            try:
                tc_pr = cell._tc.get_or_add_tcPr()
                for border_tag, border_attr in [("a:lnL", "w"), ("a:lnR", "w"), ("a:lnT", "w"), ("a:lnB", "w")]:
                    existing = tc_pr.find(qn(border_tag))
                    if existing is None:
                        from lxml import etree
                        el = etree.SubElement(tc_pr, qn(border_tag))
                        el.set(border_attr, "6350")
                        solid = etree.SubElement(el, qn("a:solidFill"))
                        srgb = etree.SubElement(solid, qn("a:srgbClr"))
                        srgb.set("val", border_color.lstrip("#"))
            except Exception:
                pass
            # Typography — shrink long cell text to fit
            cell_text = row_data[col_index] if col_index < len(row_data) else ""
            cell_w_in = col_widths[col_index] if col_index < len(col_widths) else w / cols
            text_len = len(str(cell_text))
            base_size = max(13.0, style["table_size"]) if compact else style["table_size"]
            # Scale down for overflowing cells, but keep compact tables crisp
            if not compact and (text_len > 20 or (text_len > 12 and cell_w_in < 1.5)):
                cell_size = max(8.0, base_size - 2.0)
            elif not compact and text_len > 35:
                cell_size = max(7.0, base_size - 3.0)
            else:
                cell_size = base_size
            for paragraph in cell.text_frame.paragraphs:
                paragraph.space_before = Pt(2)
                paragraph.space_after = Pt(2)
                for run in paragraph.runs:
                    run.font.size = Pt(cell_size)
                    run.font.name = style["font_name"]
                    run.font.color.rgb = RGBColor.from_string(
                        style["background"] if is_header else style["text_primary"]
                    )


def _hex(value: Any, default: str) -> str:
    if isinstance(value, str) and len(value.lstrip("#")) == 6:
        return value.lstrip("#").upper()
    return default


def _add_disclaimer_page(presentation: Any, style: dict[str, Any]) -> None:
    """Append a quality-disclaimer slide as the final page."""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt

    blank_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string(style["background"])

    # Title
    title_box = _add_textbox(
        slide, "关于本报告",
        x=style["margin_left"], y=1.5, w=12.2, h=0.7,
        font_size=style["title_size"], bold=True,
        color=style["primary"], font_name=style["font_name"],
    )
    title_box.name = "text:disclaimer:title"
    # Body
    body_lines = [
        "本 PPT 由 MeowClaw PPT Smith 自动生成。",
        "因源文档质量差异，可能导致内容不完整、数据提取不足或表述偏差。",
        "建议在正式使用前对关键数据和结论进行人工复核。",
    ]
    for i, line in enumerate(body_lines):
        body_box = _add_textbox(
            slide, line,
            x=style["margin_left"], y=2.35 + i * 0.35, w=11.5, h=0.35,
            font_size=style["body_size"], color=style["text_secondary"],
            font_name=style["font_name"],
        )
        body_box.name = f"text:disclaimer:body:{i + 1:02d}"
    # Feedback channels
    feedback_box = _add_textbox(
        slide, "📮 反馈与建议",
        x=style["margin_left"], y=3.8, w=11.5, h=0.4,
        font_size=style["body_size"] + 1, bold=True,
        color=style["primary"], font_name=style["font_name"],
    )
    feedback_box.name = "text:disclaimer:feedback"
    channels = [
        "GitHub Issues：github.com/MeowClawLab（技术问题 / 功能请求）",
        "公众号「夜猫子弦月」留言（内容建议 / 商务合作）",
    ]
    for i, ch in enumerate(channels):
        channel_box = _add_textbox(
            slide, ch,
            x=style["margin_left"], y=4.35 + i * 0.35, w=11.5, h=0.35,
            font_size=style["table_size"], color=style["text_secondary"],
            font_name=style["font_name"],
        )
        channel_box.name = f"text:disclaimer:channel:{i + 1:02d}"
    # Footer line
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    divider = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(style["margin_left"]), Inches(5.3), Inches(4.0), Inches(0.02),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor.from_string(style["border"])
    divider.line.fill.background()
    footer_box = _add_textbox(
        slide, "MeowClaw Lab · 用 AI 做好每一页",
        x=style["margin_left"], y=5.5, w=6.0, h=0.3,
        font_size=style["footnote_size"], color=style["text_secondary"],
        font_name=style["font_name"],
    )
    footer_box.name = "footer:disclaimer:brand"


def _resolve_style(contract: dict[str, Any]) -> dict[str, Any]:
    colors = contract.get("colors") if isinstance(contract.get("colors"), dict) else {}
    typography = contract.get("typography") if isinstance(contract.get("typography"), dict) else {}
    grid = contract.get("grid") if isinstance(contract.get("grid"), dict) else {}
    title_sizes = typography.get("title_sizes_pt") if isinstance(typography.get("title_sizes_pt"), dict) else {}
    body_sizes = typography.get("body_sizes_pt") if isinstance(typography.get("body_sizes_pt"), dict) else {}
    footer_tokens = contract.get("footer_tokens") if isinstance(contract.get("footer_tokens"), dict) else {}
    fonts = typography.get("font_primary") if isinstance(typography.get("font_primary"), list) else []
    data_series = colors.get("data_series")
    ds_colors: list[str] = []
    if isinstance(data_series, list):
        ds_colors = [c.lstrip("#") for c in data_series if isinstance(c, str)]
    if not ds_colors:
        ds_colors = ["#1F77B4", "#FF7F0E", "#2CA02C", "#D62728", "#9467BD", "#8C564B", "#E377C2", "#7F7F7F"]
    return {
        "background": _hex(colors.get("background"), "FFFFFF"), "primary": _hex(colors.get("primary"), "2D3340"),
        "surface": _hex(colors.get("surface_1"), "F4F1EA"), "surface_2": _hex(colors.get("surface_2", colors.get("surface_1")), "EBEDF0"),
        "border": _hex(colors.get("border"), "E4E0D7"),
        "text_primary": _hex(colors.get("text_primary"), "2D3340"), "text_secondary": _hex(colors.get("text_secondary"), "6E6A60"),
        "font_name": next((font for font in fonts if isinstance(font, str) and font not in {"sans-serif", "serif", "monospace"}), "Aptos"),
        "title_size": float(title_sizes.get("slide", 22)), "body_size": float(body_sizes.get("normal", 13)),
        "table_size": float(body_sizes.get("small", 10)), "footnote_size": float(body_sizes.get("footnote", 8.5)),
        "margin_left": float(grid.get("margin_left_in", 0.55)),
        "margin_top": float(grid.get("margin_top_in", 0.35)),
        "margin_bottom": float(grid.get("margin_bottom_in", 0.35)),
        "footer_enabled": bool(footer_tokens.get("enabled", False)),
        "footer_height": float(footer_tokens.get("height_in", 0.24)),
        "card_gap": _spacing_value(contract, "card_gap", 0.16),
        "data_series": ds_colors,
    }


def _spacing_value(contract: dict[str, Any], rule: str, default: float) -> float:
    spacing = contract.get("spacing") if isinstance(contract.get("spacing"), dict) else {}
    rules = spacing.get("rules") if isinstance(spacing.get("rules"), dict) else {}
    scale = spacing.get("scale") if isinstance(spacing.get("scale"), dict) else {}
    token = rules.get(rule)
    value = scale.get(token) if isinstance(token, str) else None
    return float(value) if isinstance(value, (int, float)) else default


def _unsupported_style_warnings(contract: dict[str, Any]) -> list[str]:
    spacing = contract.get("spacing") if isinstance(contract.get("spacing"), dict) else {}
    rules = spacing.get("rules") if isinstance(spacing.get("rules"), dict) else {}
    card_gap_token = rules.get("card_gap")
    consumed_scale = {card_gap_token: True} if isinstance(card_gap_token, str) else {}
    consumed_paths: dict[str, Any] = {
        "colors": {
            "background": True,
            "primary": True,
            "surface_1": True,
            "border": True,
            "text_primary": True,
            "text_secondary": True,
        },
        "typography": {
            "font_primary": True,
            "title_sizes_pt": {"slide": True},
            "body_sizes_pt": {"normal": True, "small": True, "footnote": True},
        },
        "grid": {"margin_left_in": True, "margin_top_in": True, "margin_bottom_in": True},
        "footer_tokens": True,
        "spacing": {
            "scale": consumed_scale,
            "rules": {"card_gap": True},
        },
    }
    warning_paths: list[str] = []
    _collect_unconsumed_style_paths(contract, consumed_paths, "", warning_paths)
    return [f"Unsupported style field ignored by PythonPptxAdapter: {path}" for path in sorted(set(warning_paths))]


def _collect_unconsumed_style_paths(value: Any, consumed_paths: Any, prefix: str, warnings: list[str]) -> None:
    if consumed_paths is True:
        return
    if isinstance(value, dict):
        if not value and prefix:
            warnings.append(prefix)
            return
        support = consumed_paths if isinstance(consumed_paths, dict) else {}
        for key, child in value.items():
            path = f"{prefix}.{key}" if prefix else key
            _collect_unconsumed_style_paths(child, support.get(key), path, warnings)
        return
    if prefix:
        warnings.append(prefix)
