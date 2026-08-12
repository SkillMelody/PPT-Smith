"""Render Plan v4 -> .pptx via python-pptx. Builders execute, never design:
every position, size, color and font arrives resolved in the plan.

Unsupported plan features fail closed with a clear error instead of being
approximated (v3 connector contract carried forward)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

_SHAPES = {
    "rect": MSO_SHAPE.RECTANGLE,
    "round_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "pill": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "diamond": MSO_SHAPE.DIAMOND,
    "chevron": MSO_SHAPE.CHEVRON,
}
_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
_VALIGN = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


class BuildError(RuntimeError):
    pass


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.lstrip("#"))


def _fill_paragraphs(text_frame, paragraphs: list[dict]) -> None:
    text_frame.word_wrap = True
    for idx, para in enumerate(paragraphs):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.alignment = _ALIGN.get(para.get("align", "left"), PP_ALIGN.LEFT)
        if para.get("line_height_pct"):
            p.line_spacing = para["line_height_pct"] / 100
        if para.get("space_after_cpt"):
            p.space_after = Pt(para["space_after_cpt"] / 100)
        for run_spec in para["runs"]:
            run = p.add_run()
            run.text = run_spec["text"]  # bullet glyphs arrive baked into the text
            font = run.font
            font.name = run_spec["font"]
            font.size = Pt(run_spec["size_cpt"] / 100)
            font.color.rgb = _rgb(run_spec["color"])
            font.bold = run_spec.get("weight", 400) >= 600
            font.italic = bool(run_spec.get("italic"))
            # Ensure CJK glyphs use the same declared face.
            rPr = run._r.get_or_add_rPr()
            ea = rPr.find(qn("a:ea"))
            if ea is None:
                ea = rPr.makeelement(qn("a:ea"), {})
                rPr.append(ea)
            ea.set("typeface", run_spec["font"])


def _apply_frame(shape, frame: dict) -> None:
    shape.left, shape.top = Emu(frame["x"]), Emu(frame["y"])
    shape.width, shape.height = Emu(frame["w"]), Emu(frame["h"])


def _add_textbox(slide, element: dict) -> object:
    box = slide.shapes.add_textbox(Emu(element["frame"]["x"]), Emu(element["frame"]["y"]),
                                   Emu(element["frame"]["w"]), Emu(element["frame"]["h"]))
    tf = box.text_frame
    tf.vertical_anchor = _VALIGN.get(element.get("valign", "top"), MSO_ANCHOR.TOP)
    _fill_paragraphs(tf, element.get("paragraphs", []))
    return box


def _add_shape(slide, element: dict) -> object:
    kind = _SHAPES.get(element.get("shape", "rect"))
    if kind is None:
        raise BuildError(f"unsupported shape: {element.get('shape')}")
    frame = element["frame"]
    shape = slide.shapes.add_shape(kind, Emu(frame["x"]), Emu(frame["y"]),
                                   Emu(frame["w"]), Emu(frame["h"]))
    if element.get("shape") in ("round_rect", "pill") and element.get("corner_radius_emu"):
        try:
            shape.adjustments[0] = min(
                element["corner_radius_emu"] / max(min(frame["w"], frame["h"]), 1), 0.5)
        except (IndexError, ValueError):
            pass
    fill = element.get("fill")
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill["color"])
    else:
        shape.fill.background()
    stroke = element.get("stroke")
    if stroke:
        shape.line.color.rgb = _rgb(stroke["color"])
        shape.line.width = Emu(stroke["width_emu"])
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    if element.get("paragraphs"):
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Emu(152400)
        tf.margin_top = tf.margin_bottom = Emu(101600)
        _fill_paragraphs(tf, element["paragraphs"])
    return shape


def _add_table(slide, element: dict) -> object:
    frame = element["frame"]
    rows, cols = len(element["rows"]), len(element["col_widths_emu"])
    graphic = slide.shapes.add_table(rows, cols, Emu(frame["x"]), Emu(frame["y"]),
                                     Emu(frame["w"]), Emu(frame["h"]))
    table = graphic.table
    for c_idx, width in enumerate(element["col_widths_emu"]):
        table.columns[c_idx].width = Emu(width)
    for r_idx, row_spec in enumerate(element["rows"]):
        if row_spec.get("height_emu"):
            table.rows[r_idx].height = Emu(row_spec["height_emu"])
        for c_idx in range(cols):
            cell = table.cell(r_idx, c_idx)
            specs = row_spec["cells"]
            spec = specs[c_idx] if c_idx < len(specs) else {"paragraphs": []}
            if spec.get("fill"):
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(spec["fill"]["color"])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = cell.margin_bottom = Emu(45720)
            if spec.get("paragraphs"):
                _fill_paragraphs(cell.text_frame, spec["paragraphs"])
    return graphic


def _add_connector(slide, element: dict, registry: dict) -> object:
    begin = registry.get(element["from_element"])
    end = registry.get(element["to_element"])
    if begin is None or end is None:
        raise BuildError(f"connector endpoints missing: {element['element_id']}")
    if element.get("route", "straight") != "straight":
        raise BuildError("only straight connectors are implemented in the v4 MVP builder")
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0, 0, 0, 0)
    connector.begin_connect(begin, 3)
    connector.end_connect(end, 1)
    stroke = element["stroke"]
    connector.line.color.rgb = _rgb(stroke["color"])
    connector.line.width = Emu(stroke["width_emu"])
    if element.get("arrowhead", "none") in ("end", "both"):
        line_el = connector.line._get_or_add_ln()
        for tag, wanted in (("a:tailEnd", True),
                            ("a:headEnd", element["arrowhead"] == "both")):
            if wanted:
                arrow = line_el.find(qn(tag))
                if arrow is None:
                    arrow = line_el.makeelement(qn(tag), {})
                    line_el.append(arrow)
                arrow.set("type", "arrow")
    label = element.get("label")
    if label:
        _add_textbox(slide, {"frame": label["frame"], "paragraphs": label["paragraphs"],
                             "valign": "middle"})
    return connector


def build_pptx(plan: dict, output_path: str | Path) -> Path:
    prs = Presentation()
    prs.slide_width = Emu(plan["canvas"]["width_emu"])
    prs.slide_height = Emu(plan["canvas"]["height_emu"])
    blank = prs.slide_layouts[6]

    for slide_spec in plan["slides"]:
        slide = prs.slides.add_slide(blank)
        background = slide_spec.get("background")
        if background:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = _rgb(background["color"])
        registry: dict[str, object] = {}
        connectors = []
        for element in slide_spec["elements"]:
            etype = element["type"]
            if etype == "textbox":
                registry[element["element_id"]] = _add_textbox(slide, element)
            elif etype == "shape":
                registry[element["element_id"]] = _add_shape(slide, element)
            elif etype == "table":
                registry[element["element_id"]] = _add_table(slide, element)
            elif etype == "connector":
                connectors.append(element)  # after endpoints exist
            elif etype == "image":
                path = Path(element["path"])
                if not path.is_file():
                    raise BuildError(f"image asset missing: {path}")
                pic = slide.shapes.add_picture(str(path), Emu(element["frame"]["x"]),
                                               Emu(element["frame"]["y"]),
                                               Emu(element["frame"]["w"]),
                                               Emu(element["frame"]["h"]))
                registry[element["element_id"]] = pic
            else:
                raise BuildError(f"unsupported element type in MVP builder: {etype}")
        for element in connectors:
            _add_connector(slide, element, registry)

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return output
