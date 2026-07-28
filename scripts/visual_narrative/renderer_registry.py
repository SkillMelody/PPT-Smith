from __future__ import annotations

import json
import importlib
import sys
import tempfile
from pathlib import Path
from typing import Any


_RENDERER_BINDINGS = {
    ("python_pptx", "comparison_card"): {
        "handler": "builders.python_pptx_adapter._add_card",
        "schema_versions": {"1.0"},
    },
    ("python_pptx", "metric_card"): {
        "handler": "builders.python_pptx_adapter._add_card",
        "schema_versions": {"1.0"},
    },
    ("python_pptx", "summary_action_card"): {
        "handler": "builders.python_pptx_adapter._add_card",
        "schema_versions": {"1.0"},
    },
    ("python_pptx", "bar_chart"): {
        "handler": "builders.python_pptx_adapter._add_chart",
        "schema_versions": {"1.0"},
    },
    ("python_pptx", "process_flow"): {
        "handler": "builders.python_pptx_adapter._add_process",
        "schema_versions": {"1.0"},
    },
    ("python_pptx", "native_table"): {
        "handler": "builders.python_pptx_adapter._add_table",
        "schema_versions": {"1.0"},
    },
    ("python_pptx", "heat_matrix"): {
        "handler": "scripts.visual_narrative.renderers.heat_matrix.render",
        "schema_versions": {"1.0"},
    },
    ("python_pptx", "layered_architecture"): {
        "handler": "scripts.visual_narrative.renderers.layered_architecture.render",
        "schema_versions": {"1.0"},
    },
    ("python_pptx", "drill_down_stair"): {
        "handler": "scripts.visual_narrative.renderers.drill_down_stair.render",
        "schema_versions": {"1.0"},
    },
    ("python_pptx", "phase_roadmap"): {
        "handler": "scripts.visual_narrative.renderers.phase_roadmap.render",
        "schema_versions": {"1.0"},
    },
    **{
        ("pptxgenjs", component_type): {
            "handler": "builders.pptxgenjs_adapter.PptxGenJsAdapter.build",
            "schema_versions": {"1.0"},
        }
        for component_type in (
            "comparison_card",
            "metric_card",
            "summary_action_card",
            "bar_chart",
            "process_flow",
            "native_table",
        )
    },
}
_SUPPORT_LEVELS = {"visual_only": 0, "partial": 1, "full": 2}


def smoke_test_renderer(builder: str, component_type: str, schema_version: str) -> bool:
    binding = _RENDERER_BINDINGS.get((builder, component_type))
    if not binding or schema_version not in binding["schema_versions"]:
        return False

    skill_root = Path(__file__).resolve().parents[2]
    inserted_path = str(skill_root) not in sys.path
    if inserted_path:
        sys.path.insert(0, str(skill_root))
    try:
        if builder == "pptxgenjs":
            return _smoke_test_pptxgenjs(component_type)

        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        if component_type in {
            "comparison_card",
            "metric_card",
            "summary_action_card",
            "bar_chart",
            "process_flow",
            "native_table",
        }:
            return _smoke_test_python_renderer(slide, component_type)

        module_name, function_name = binding["handler"].rsplit(".", 1)
        renderer = getattr(importlib.import_module(module_name), function_name)
        result = renderer(
            slide,
            _smoke_object(component_type),
            {},
            {},
        )
        return (
            result.get("actual_route") == "native_diagram"
            and bool(result.get("object_names"))
            and len(slide.shapes) >= _minimum_smoke_shape_count(component_type)
        )
    except Exception:
        return False
    finally:
        if inserted_path:
            sys.path.remove(str(skill_root))


def _smoke_test_python_renderer(slide: Any, component_type: str) -> bool:
    from builders.python_pptx_adapter import (
        _add_card,
        _add_chart,
        _add_process,
        _add_table,
        _resolve_style,
    )

    style = _resolve_style({})
    obj = _smoke_object(component_type)
    if component_type == "native_table":
        _add_table(slide, obj, x=0.5, y=0.5, w=2.0, h=1.0, style=style)
        return len(slide.shapes) == 1 and slide.shapes[0].has_table
    if component_type in {"comparison_card", "metric_card", "summary_action_card"}:
        _add_card(slide, obj, x=0.5, y=0.5, w=2.0, h=1.0, style=style)
        return len(slide.shapes) == 1 and bool(slide.shapes[0].has_text_frame)
    if component_type == "bar_chart":
        _add_chart(slide, obj, x=0.5, y=0.5, w=5.8, h=2.5, style=style)
        return len(slide.shapes) == 1 and bool(slide.shapes[0].has_chart)
    _add_process(slide, obj, x=0.5, y=0.5, w=5.8, h=1.0, style=style)
    return len(slide.shapes) >= 3


def _smoke_test_pptxgenjs(component_type: str) -> bool:
    from builders.base import BuildPlan
    from builders.pptxgenjs_adapter import PptxGenJsAdapter

    route = {
        "bar_chart": "native_chart",
        "process_flow": "native_diagram",
        "native_table": "native_table",
    }.get(component_type, "native_ppt")
    plan = BuildPlan(
        builder="pptxgenjs",
        slides=[
            {
                "id": "S01",
                "title": "Renderer smoke",
                "objects": [
                    {
                        **_smoke_object(component_type),
                        "editability": "native_required",
                        "delivery_plan": {"selected_route": route},
                    }
                ],
            }
        ],
    )
    with tempfile.TemporaryDirectory(prefix="pptsmith-renderer-smoke-") as temp_dir:
        result = PptxGenJsAdapter(timeout_seconds=15).build(plan, Path(temp_dir))
        deck = Path(result.pptx) if result.pptx else None
        return (
            result.status == "created"
            and deck is not None
            and deck.is_file()
            and deck.stat().st_size > 0
            and len(result.object_results) == 1
            and not result.fallbacks
        )


def _smoke_object(component_type: str) -> dict[str, Any]:
    content_by_component: dict[str, dict[str, Any]] = {
        "comparison_card": {"title": "Before", "claim": "Editable comparison"},
        "metric_card": {"title": "Coverage", "value": "100%"},
        "summary_action_card": {"title": "Next", "claim": "Ship verified output"},
        "bar_chart": {
            "categories": ["A", "B"],
            "series": [{"name": "Score", "values": [2, 5]}],
        },
        "process_flow": {
            "nodes": [{"label": "Plan"}, {"label": "Build"}, {"label": "Verify"}],
        },
        "native_table": {"rows": [["smoke"]]},
        "heat_matrix": {
            "rows": ["A", "B"],
            "columns": ["X", "Y"],
            "values": [[1, 2], [3, 4]],
        },
        "layered_architecture": {
            "boundary_label": "Platform",
            "layers": [
                {"label": "Experience", "nodes": ["Portal"]},
                {"label": "Data", "nodes": ["Lake"]},
            ],
        },
        "drill_down_stair": {
            "levels": [
                {"label": "Group"},
                {"label": "Record"},
            ],
        },
        "phase_roadmap": {
            "phases": [
                {"label": "Phase 1"},
                {"label": "Phase 2"},
            ],
            "milestones": [{"label": "Gate", "phase": 0}],
        },
    }
    return {
        "id": f"smoke-{component_type}",
        "component_type": component_type,
        "content": content_by_component[component_type],
    }


def _minimum_smoke_shape_count(component_type: str) -> int:
    return {
        "heat_matrix": 10,
        "layered_architecture": 4,
        "drill_down_stair": 3,
        "phase_roadmap": 4,
    }[component_type]


def load_renderer_capabilities(path: Path) -> dict[str, Any]:
    document = json.loads(path.read_text(encoding="utf-8"))
    if document.get("schema_version") != "1.0":
        raise ValueError("RENDERER_CAPABILITY_SCHEMA_UNSUPPORTED")
    return document


def effective_component_support(
    registry: dict[str, Any],
    renderer_capabilities: dict[str, Any],
    builder: str,
) -> dict[str, str]:
    handlers = (renderer_capabilities.get("renderers") or {}).get(builder, {})
    registry_schema_version = registry.get("schema_version")
    support: dict[str, str] = {}
    for component in registry.get("components", []):
        component_type = component["component_type"]
        declared = ((component.get("builder_support") or {}).get(builder) or {}).get("level", "unknown")
        handler = handlers.get(component_type)
        binding = _RENDERER_BINDINGS.get((builder, component_type))
        capability_level = handler.get("level") if isinstance(handler, dict) else None
        schema_versions = handler.get("schema_versions") if isinstance(handler, dict) else None
        if (
            not binding
            or not isinstance(handler, dict)
            or handler.get("handler") != binding["handler"]
            or declared not in _SUPPORT_LEVELS
            or capability_level not in _SUPPORT_LEVELS
            or _SUPPORT_LEVELS[capability_level] < _SUPPORT_LEVELS[declared]
            or not isinstance(schema_versions, list)
            or registry_schema_version not in schema_versions
            or handler.get("smoke_test") != "passed"
            or not smoke_test_renderer(builder, component_type, registry_schema_version)
        ):
            support[component_type] = "unsupported"
            continue
        support[component_type] = declared
    return support
