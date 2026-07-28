from __future__ import annotations

from dataclasses import dataclass, field
from typing import Any


@dataclass
class RenderResult:
    actual_route: str
    object_names: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        return {
            "actual_route": self.actual_route,
            "object_names": list(self.object_names),
            "warnings": list(self.warnings),
        }


def resolve_style(style: dict[str, Any]) -> dict[str, str]:
    semantic = style.get("color_tokens") if isinstance(style.get("color_tokens"), dict) else {}
    colors = style.get("colors") if isinstance(style.get("colors"), dict) else {}
    return {
        "ink_navy": _hex(semantic.get("ink_navy") or colors.get("primary"), "12233F"),
        "signal_blue": _hex(semantic.get("signal_blue") or colors.get("accent"), "2563EB"),
        "insight_blue": _hex(semantic.get("insight_blue"), "4C82F7"),
        "mist_blue": _hex(semantic.get("mist_blue") or colors.get("surface_2"), "EAF2FF"),
        "paper_white": _hex(semantic.get("paper_white") or colors.get("background"), "F8FAFC"),
        "slate_text": _hex(semantic.get("slate_text") or colors.get("text_secondary"), "475569"),
        "hairline_grey": _hex(semantic.get("hairline_grey") or colors.get("border"), "CBD5E1"),
        "warning": _hex(
            (style.get("semantic_colors") or {}).get("warning") if isinstance(style.get("semantic_colors"), dict) else colors.get("warning"),
            "D97706",
        ),
        "success": _hex(
            (style.get("semantic_colors") or {}).get("success") if isinstance(style.get("semantic_colors"), dict) else colors.get("positive"),
            "15856C",
        ),
    }


def content_box(intent: dict[str, Any]) -> tuple[float, float, float, float]:
    render_box = intent.get("render_box") if isinstance(intent.get("render_box"), dict) else {}
    return (
        _number(render_box.get("x"), 0.8),
        _number(render_box.get("y"), 1.75),
        _number(render_box.get("w"), 11.7),
        _number(render_box.get("h"), 4.9),
    )


def add_box(
    slide: Any,
    *,
    name: str,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str | None,
    line: str,
    text: str = "",
    text_color: str = "12233F",
    font_size: float = 11,
    bold: bool = False,
    radius: bool = True,
    line_width: float = 1.0,
) -> Any:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches, Pt

    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if radius
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.name = name
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.color.rgb = RGBColor.from_string(line)
    shape.line.width = Pt(line_width)
    if text:
        set_shape_text(
            shape,
            text,
            color=text_color,
            font_size=font_size,
            bold=bold,
        )
    return shape


def add_circle(
    slide: Any,
    *,
    name: str,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str,
    line: str,
    text: str,
    text_color: str,
    font_size: float = 9,
) -> Any:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.color.rgb = RGBColor.from_string(line)
    shape.line.width = Pt(1.0)
    set_shape_text(shape, text, color=text_color, font_size=font_size, bold=True)
    return shape


def add_text(
    slide: Any,
    *,
    name: str,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    color: str,
    font_size: float = 10,
    bold: bool = False,
) -> Any:
    from pptx.util import Inches

    shape = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.name = name
    set_shape_text(shape, text, color=color, font_size=font_size, bold=bold)
    return shape


def set_shape_text(
    shape: Any,
    text: str,
    *,
    color: str,
    font_size: float,
    bold: bool = False,
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Pt

    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Pt(4)
    frame.margin_right = Pt(4)
    frame.margin_top = Pt(2)
    frame.margin_bottom = Pt(2)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def add_bound_connector(
    slide: Any,
    start: Any,
    end: Any,
    *,
    name: str,
    color: str,
    vertical: bool = False,
    width_pt: float = 1.6,
) -> Any:
    from lxml import etree
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn
    from pptx.util import Pt

    if vertical:
        start_x = start.left + start.width // 2
        start_y = start.top + start.height
        end_x = end.left + end.width // 2
        end_y = end.top
        start_site, end_site = 2, 0
    else:
        start_x = start.left + start.width
        start_y = start.top + start.height // 2
        end_x = end.left
        end_y = end.top + end.height // 2
        start_site, end_site = 3, 1
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        start_x,
        start_y,
        end_x,
        end_y,
    )
    connector.name = name
    connector.line.color.rgb = RGBColor.from_string(color)
    connector.line.width = Pt(width_pt)
    connector.begin_connect(start, start_site)
    connector.end_connect(end, end_site)
    tail_end = etree.Element(qn("a:tailEnd"))
    tail_end.set("type", "triangle")
    connector.line._get_or_add_ln().append(tail_end)
    return connector


def blend(low: str, high: str, ratio: float) -> str:
    bounded = max(0.0, min(1.0, ratio))
    low_rgb = tuple(int(low[index : index + 2], 16) for index in (0, 2, 4))
    high_rgb = tuple(int(high[index : index + 2], 16) for index in (0, 2, 4))
    return "".join(
        f"{round(left + (right - left) * bounded):02X}"
        for left, right in zip(low_rgb, high_rgb)
    )


def _hex(value: Any, default: str) -> str:
    if isinstance(value, str):
        normalized = value.lstrip("#").upper()
        if len(normalized) == 6 and all(char in "0123456789ABCDEF" for char in normalized):
            return normalized
    return default


def _number(value: Any, default: float) -> float:
    return float(value) if isinstance(value, (int, float)) else default
